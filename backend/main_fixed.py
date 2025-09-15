from fastapi import FastAPI, Response, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any
import html
import json
import uuid
import os
import base64
import subprocess
import tempfile
import logging
import traceback
from datetime import datetime
from pathlib import Path
import requests
import google.generativeai as genai

# Document processing imports
import PyPDF2
import openpyxl
from pptx import Presentation

# Import diagrams for Azure architecture generation
from diagrams import Diagram, Cluster, Edge
from diagrams.azure.compute import VM, AKS, AppServices, FunctionApps, ContainerInstances, ServiceFabricClusters, BatchAccounts
from diagrams.azure.network import VirtualNetworks, ApplicationGateway, LoadBalancers, Firewall, ExpressrouteCircuits, VirtualNetworkGateways
from diagrams.azure.storage import StorageAccounts, BlobStorage, DataLakeStorage
from diagrams.azure.database import SQLDatabases, CosmosDb, DatabaseForMysqlServers, DatabaseForPostgresqlServers
from diagrams.azure.security import KeyVaults, SecurityCenter, Sentinel
from diagrams.azure.identity import ActiveDirectory
from diagrams.azure.analytics import SynapseAnalytics, DataFactories, Databricks, StreamAnalyticsJobs, EventHubs
from diagrams.azure.integration import LogicApps, ServiceBus, EventGridTopics, APIManagement
from diagrams.azure.devops import Devops, Pipelines
from diagrams.azure.general import Subscriptions, Resourcegroups
from diagrams.azure.web import AppServices as WebApps

app = FastAPI(
    title="Azure Landing Zone Agent",
    description="Professional Azure Landing Zone Architecture Generator",
    version="1.0.0"
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # open for dev, restrict later
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure Google Gemini API
GEMINI_API_KEY = "AIzaSyCuYYvGh5wjwNniv9ZQ1QC-5pxwdj5lCWQ"
genai.configure(api_key=GEMINI_API_KEY)

# Initialize Gemini model
try:
    gemini_model = genai.GenerativeModel('gemini-1.5-pro')
    logger.info("Google Gemini API configured successfully")
except Exception as e:
    logger.error(f"Failed to configure Gemini API: {e}")
    gemini_model = None


class CustomerInputs(BaseModel):
    # Business Requirements
    business_objective: Optional[str] = Field(None, description="Primary business objective")
    regulatory: Optional[str] = Field(None, description="Regulatory requirements")
    industry: Optional[str] = Field(None, description="Industry vertical")
    
    # Organization Structure
    org_structure: Optional[str] = Field(None, description="Organization structure")
    governance: Optional[str] = Field(None, description="Governance model")
    
    # Identity & Access Management
    identity: Optional[str] = Field(None, description="Identity management approach")
    
    # Networking & Connectivity
    connectivity: Optional[str] = Field(None, description="Connectivity requirements")
    network_model: Optional[str] = Field(None, description="Network topology model")
    ip_strategy: Optional[str] = Field(None, description="IP address strategy")
    
    # Security
    security_zone: Optional[str] = Field(None, description="Security zone requirements")
    security_posture: Optional[str] = Field(None, description="Security posture")
    key_vault: Optional[str] = Field(None, description="Key management approach")
    threat_protection: Optional[str] = Field(None, description="Threat protection strategy")
    
    # Legacy fields for backward compatibility
    workload: Optional[str] = Field(None, description="Primary workload type")
    architecture_style: Optional[str] = Field(None, description="Architecture style")
    scalability: Optional[str] = Field(None, description="Scalability requirements")
    
    # Operations
    ops_model: Optional[str] = Field(None, description="Operations model")
    monitoring: Optional[str] = Field(None, description="Monitoring strategy")
    backup: Optional[str] = Field(None, description="Backup and recovery strategy")
    
    # Infrastructure
    topology_pattern: Optional[str] = Field(None, description="Topology pattern")
    
    # Migration & Cost
    migration_scope: Optional[str] = Field(None, description="Migration scope")
    cost_priority: Optional[str] = Field(None, description="Cost optimization priority")
    iac: Optional[str] = Field(None, description="Infrastructure as Code preference")
    
    # Specific Azure Service Selections
    # Compute Services
    compute_services: Optional[List[str]] = Field(default_factory=list, description="Selected compute services")
    
    # Networking Services
    network_services: Optional[List[str]] = Field(default_factory=list, description="Selected networking services")
    
    # Storage Services
    storage_services: Optional[List[str]] = Field(default_factory=list, description="Selected storage services")
    
    # Database Services
    database_services: Optional[List[str]] = Field(default_factory=list, description="Selected database services")
    
    # Security Services
    security_services: Optional[List[str]] = Field(default_factory=list, description="Selected security services")
    
    # Monitoring Services
    monitoring_services: Optional[List[str]] = Field(default_factory=list, description="Selected monitoring services")
    
    # AI/ML Services
    ai_services: Optional[List[str]] = Field(default_factory=list, description="Selected AI/ML services")
    
    # Analytics Services
    analytics_services: Optional[List[str]] = Field(default_factory=list, description="Selected analytics services")
    
    # Integration Services
    integration_services: Optional[List[str]] = Field(default_factory=list, description="Selected integration services")
    
    # DevOps Services
    devops_services: Optional[List[str]] = Field(default_factory=list, description="Selected DevOps services")
    
    # Backup Services
    backup_services: Optional[List[str]] = Field(default_factory=list, description="Selected backup services")
    
    # Enhanced Input Fields for AI Integration
    free_text_input: Optional[str] = Field(None, description="Free-form text input for additional requirements and context")
    url_input: Optional[str] = Field(None, description="URL for web content analysis")
    uploaded_files_info: Optional[List[Dict[str, Any]]] = Field(default_factory=list, description="Information about uploaded files")


# Azure Architecture Templates and Patterns
AZURE_TEMPLATES = {
    "enterprise": {
        "name": "Enterprise Scale Landing Zone",
        "management_groups": ["Root", "Platform", "Landing Zones", "Sandbox", "Decommissioned"],
        "subscriptions": ["Connectivity", "Identity", "Management", "Production", "Development"],
        "core_services": ["Azure AD", "Azure Policy", "Azure Monitor", "Key Vault", "Security Center"]
    },
    "small_medium": {
        "name": "Small-Medium Enterprise Landing Zone", 
        "management_groups": ["Root", "Platform", "Workloads"],
        "subscriptions": ["Platform", "Production", "Development"],
        "core_services": ["Azure AD", "Azure Monitor", "Key Vault"]
    },
    "startup": {
        "name": "Startup Landing Zone",
        "management_groups": ["Root", "Workloads"],
        "subscriptions": ["Production", "Development"],
        "core_services": ["Azure AD", "Azure Monitor"]
    }
}

AZURE_SERVICES_MAPPING = {
    # Compute Services
    "virtual_machines": {"name": "Azure Virtual Machines", "icon": "🖥️", "drawio_shape": "virtual_machine", "diagram_class": VM, "category": "compute"},
    "aks": {"name": "Azure Kubernetes Service", "icon": "☸️", "drawio_shape": "kubernetes_service", "diagram_class": AKS, "category": "compute"},
    "app_services": {"name": "Azure App Services", "icon": "🌐", "drawio_shape": "app_services", "diagram_class": AppServices, "category": "compute"},
    "web_apps": {"name": "Azure Web Apps", "icon": "🌐", "drawio_shape": "app_services", "diagram_class": WebApps, "category": "compute"},
    "functions": {"name": "Azure Functions", "icon": "⚡", "drawio_shape": "function_app", "diagram_class": FunctionApps, "category": "compute"},
    "container_instances": {"name": "Container Instances", "icon": "📦", "drawio_shape": "container_instances", "diagram_class": ContainerInstances, "category": "compute"},
    "service_fabric": {"name": "Service Fabric", "icon": "🏗️", "drawio_shape": "service_fabric", "diagram_class": ServiceFabricClusters, "category": "compute"},
    "batch": {"name": "Azure Batch", "icon": "⚙️", "drawio_shape": "batch_accounts", "diagram_class": BatchAccounts, "category": "compute"},
    
    # Networking Services
    "virtual_network": {"name": "Virtual Network", "icon": "🌐", "drawio_shape": "virtual_network", "diagram_class": VirtualNetworks, "category": "network"},
    "vpn_gateway": {"name": "VPN Gateway", "icon": "🔒", "drawio_shape": "vpn_gateway", "diagram_class": VirtualNetworkGateways, "category": "network"},
    "expressroute": {"name": "ExpressRoute", "icon": "⚡", "drawio_shape": "expressroute_circuits", "diagram_class": ExpressrouteCircuits, "category": "network"},
    "load_balancer": {"name": "Load Balancer", "icon": "⚖️", "drawio_shape": "load_balancer", "diagram_class": LoadBalancers, "category": "network"},
    "application_gateway": {"name": "Application Gateway", "icon": "🚪", "drawio_shape": "application_gateway", "diagram_class": ApplicationGateway, "category": "network"},
    "firewall": {"name": "Azure Firewall", "icon": "🛡️", "drawio_shape": "firewall", "diagram_class": Firewall, "category": "network"},
    "waf": {"name": "Web Application Firewall", "icon": "🛡️", "drawio_shape": "application_gateway", "diagram_class": ApplicationGateway, "category": "network"},
    "cdn": {"name": "Content Delivery Network", "icon": "🌍", "drawio_shape": "cdn_profiles", "diagram_class": None, "category": "network"},
    "traffic_manager": {"name": "Traffic Manager", "icon": "🚦", "drawio_shape": "traffic_manager_profiles", "diagram_class": None, "category": "network"},
    "virtual_wan": {"name": "Virtual WAN", "icon": "🌐", "drawio_shape": "virtual_wan", "diagram_class": VirtualNetworks, "category": "network"},
    
    # Storage Services
    "storage_accounts": {"name": "Storage Accounts", "icon": "💾", "drawio_shape": "storage_accounts", "diagram_class": StorageAccounts, "category": "storage"},
    "blob_storage": {"name": "Blob Storage", "icon": "📄", "drawio_shape": "blob_storage", "diagram_class": BlobStorage, "category": "storage"},
    "file_storage": {"name": "Azure Files", "icon": "📁", "drawio_shape": "files", "diagram_class": StorageAccounts, "category": "storage"},
    "disk_storage": {"name": "Managed Disks", "icon": "💿", "drawio_shape": "managed_disks", "diagram_class": StorageAccounts, "category": "storage"},
    "data_lake": {"name": "Data Lake Storage", "icon": "🏞️", "drawio_shape": "data_lake_storage", "diagram_class": DataLakeStorage, "category": "storage"},
    
    # Database Services
    "sql_database": {"name": "Azure SQL Database", "icon": "🗄️", "drawio_shape": "sql_database", "diagram_class": SQLDatabases, "category": "database"},
    "sql_managed_instance": {"name": "SQL Managed Instance", "icon": "🗄️", "drawio_shape": "sql_managed_instance", "diagram_class": SQLDatabases, "category": "database"},
    "cosmos_db": {"name": "Cosmos DB", "icon": "🌍", "drawio_shape": "cosmos_db", "diagram_class": CosmosDb, "category": "database"},
    "mysql": {"name": "Azure Database for MySQL", "icon": "🐬", "drawio_shape": "database_for_mysql_servers", "diagram_class": DatabaseForMysqlServers, "category": "database"},
    "postgresql": {"name": "Azure Database for PostgreSQL", "icon": "🐘", "drawio_shape": "database_for_postgresql_servers", "diagram_class": DatabaseForPostgresqlServers, "category": "database"},
    "mariadb": {"name": "Azure Database for MariaDB", "icon": "🗄️", "drawio_shape": "database_for_mariadb_servers", "diagram_class": DatabaseForMysqlServers, "category": "database"},
    "redis_cache": {"name": "Azure Cache for Redis", "icon": "⚡", "drawio_shape": "cache_redis", "diagram_class": None, "category": "database"},
    
    # Security Services
    "key_vault": {"name": "Azure Key Vault", "icon": "🔐", "drawio_shape": "key_vault", "diagram_class": KeyVaults, "category": "security"},
    "active_directory": {"name": "Azure Active Directory", "icon": "👤", "drawio_shape": "azure_active_directory", "diagram_class": ActiveDirectory, "category": "security"},
    "security_center": {"name": "Azure Security Center", "icon": "🛡️", "drawio_shape": "security_center", "diagram_class": SecurityCenter, "category": "security"},
    "sentinel": {"name": "Azure Sentinel", "icon": "👁️", "drawio_shape": "sentinel", "diagram_class": Sentinel, "category": "security"},
    "defender": {"name": "Microsoft Defender", "icon": "🛡️", "drawio_shape": "defender_easm", "diagram_class": SecurityCenter, "category": "security"},
    "information_protection": {"name": "Azure Information Protection", "icon": "🔒", "drawio_shape": "information_protection", "diagram_class": None, "category": "security"},
    
    # Monitoring & Management
    "monitor": {"name": "Azure Monitor", "icon": "📊", "drawio_shape": "monitor", "diagram_class": None, "category": "monitoring"},
    "log_analytics": {"name": "Log Analytics", "icon": "📋", "drawio_shape": "log_analytics_workspaces", "diagram_class": None, "category": "monitoring"},
    "application_insights": {"name": "Application Insights", "icon": "📈", "drawio_shape": "application_insights", "diagram_class": None, "category": "monitoring"},
    "service_health": {"name": "Service Health", "icon": "❤️", "drawio_shape": "service_health", "diagram_class": None, "category": "monitoring"},
    "advisor": {"name": "Azure Advisor", "icon": "💡", "drawio_shape": "advisor", "diagram_class": None, "category": "monitoring"},
    
    # AI/ML Services  
    "cognitive_services": {"name": "Cognitive Services", "icon": "🧠", "drawio_shape": "cognitive_services", "diagram_class": None, "category": "ai"},
    "machine_learning": {"name": "Azure Machine Learning", "icon": "🤖", "drawio_shape": "machine_learning", "diagram_class": None, "category": "ai"},
    "bot_service": {"name": "Bot Service", "icon": "🤖", "drawio_shape": "bot_services", "diagram_class": None, "category": "ai"},
    "form_recognizer": {"name": "Form Recognizer", "icon": "📄", "drawio_shape": "form_recognizer", "diagram_class": None, "category": "ai"},
    
    # Data & Analytics
    "synapse": {"name": "Azure Synapse Analytics", "icon": "📊", "drawio_shape": "synapse_analytics", "diagram_class": SynapseAnalytics, "category": "analytics"},
    "data_factory": {"name": "Azure Data Factory", "icon": "🏭", "drawio_shape": "data_factory", "diagram_class": DataFactories, "category": "analytics"},
    "databricks": {"name": "Azure Databricks", "icon": "📊", "drawio_shape": "databricks", "diagram_class": Databricks, "category": "analytics"},
    "stream_analytics": {"name": "Stream Analytics", "icon": "🌊", "drawio_shape": "stream_analytics", "diagram_class": StreamAnalyticsJobs, "category": "analytics"},
    "power_bi": {"name": "Power BI", "icon": "📊", "drawio_shape": "power_bi", "diagram_class": None, "category": "analytics"},
    
    # Integration Services
    "logic_apps": {"name": "Logic Apps", "icon": "🔗", "drawio_shape": "logic_apps", "diagram_class": LogicApps, "category": "integration"},
    "service_bus": {"name": "Service Bus", "icon": "🚌", "drawio_shape": "service_bus", "diagram_class": ServiceBus, "category": "integration"},
    "event_grid": {"name": "Event Grid", "icon": "⚡", "drawio_shape": "event_grid_topics", "diagram_class": EventGridTopics, "category": "integration"},
    "event_hubs": {"name": "Event Hubs", "icon": "📡", "drawio_shape": "event_hubs", "diagram_class": EventHubs, "category": "integration"},
    "api_management": {"name": "API Management", "icon": "🔌", "drawio_shape": "api_management", "diagram_class": APIManagement, "category": "integration"},
    
    # DevOps & Management
    "devops": {"name": "Azure DevOps", "icon": "⚙️", "drawio_shape": "devops", "diagram_class": Devops, "category": "devops"},
    "automation": {"name": "Azure Automation", "icon": "🤖", "drawio_shape": "automation_accounts", "diagram_class": None, "category": "devops"},
    "policy": {"name": "Azure Policy", "icon": "📋", "drawio_shape": "policy", "diagram_class": None, "category": "governance"},
    "blueprints": {"name": "Azure Blueprints", "icon": "📐", "drawio_shape": "blueprints", "diagram_class": None, "category": "governance"},
    "resource_manager": {"name": "Azure Resource Manager", "icon": "🏗️", "drawio_shape": "resource_groups", "diagram_class": Resourcegroups, "category": "governance"},
    
    # Backup & Recovery
    "backup": {"name": "Azure Backup", "icon": "💾", "drawio_shape": "backup", "diagram_class": None, "category": "backup"},
    "site_recovery": {"name": "Azure Site Recovery", "icon": "🔄", "drawio_shape": "site_recovery", "diagram_class": None, "category": "backup"},
}

def get_safe_output_directory() -> str:
    """Get a safe directory for output files with fallback options"""
    directories_to_try = [
        "/tmp",
        tempfile.gettempdir(),
        os.path.expanduser("~/tmp"),
        "./tmp"
    ]
    
    for directory in directories_to_try:
        try:
            # Create directory if it doesn't exist
            Path(directory).mkdir(parents=True, exist_ok=True)
            
            # Test if we can write to it
            test_file = os.path.join(directory, f"test_write_{uuid.uuid4().hex[:8]}.tmp")
            with open(test_file, 'w') as f:
                f.write("test")
            os.remove(test_file)
            
            logger.info(f"Using output directory: {directory}")
            return directory
            
        except Exception as e:
            logger.warning(f"Cannot use directory {directory}: {e}")
            continue
    
    raise Exception("No writable output directory found. Tried: " + ", ".join(directories_to_try))

def cleanup_old_files(directory: str, max_age_hours: int = 24):
    """Clean up old generated files to prevent disk space issues"""
    try:
        current_time = datetime.now().timestamp()
        max_age_seconds = max_age_hours * 3600
        
        for filename in os.listdir(directory):
            if filename.startswith("azure_landing_zone_") and filename.endswith(".png"):
                filepath = os.path.join(directory, filename)
                try:
                    file_age = current_time - os.path.getmtime(filepath)
                    if file_age > max_age_seconds:
                        os.remove(filepath)
                        logger.info(f"Cleaned up old file: {filename}")
                except Exception as e:
                    logger.warning(f"Failed to clean up file {filename}: {e}")
                    
    except Exception as e:
        logger.warning(f"Failed to perform cleanup in {directory}: {e}")

# Google Gemini AI Integration Functions
def analyze_url_content(url: str) -> str:
    """Fetch and analyze URL content using Gemini AI"""
    try:
        if not gemini_model:
            return "Gemini AI not available for URL analysis"
            
        # Fetch URL content
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        content = response.text[:10000]  # Limit content size
        
        prompt = f"""
        Analyze the following web content for Azure architecture planning:
        
        URL: {url}
        Content: {content}
        
        Please provide insights for:
        1. Relevant Azure services mentioned or implied
        2. Architecture patterns or requirements
        3. Security and compliance considerations
        4. Scalability and performance requirements
        5. Integration requirements
        
        Format your response as a structured analysis.
        """
        
        result = gemini_model.generate_content(prompt)
        return result.text
        
    except Exception as e:
        logger.error(f"Error analyzing URL {url}: {e}")
        return f"Error analyzing URL: {str(e)}"

def process_uploaded_document(file_content: bytes, filename: str, file_type: str) -> str:
    """Process uploaded document using Gemini AI"""
    try:
        if not gemini_model:
            return "Gemini AI not available for document analysis"
            
        text_content = ""
        
        # Extract text based on file type
        if file_type.lower() == 'pdf':
            text_content = extract_pdf_text(file_content)
        elif file_type.lower() in ['xlsx', 'xls']:
            text_content = extract_excel_text(file_content)
        elif file_type.lower() in ['pptx', 'ppt']:
            text_content = extract_pptx_text(file_content)
        else:
            return f"Unsupported file type: {file_type}"
        
        if not text_content.strip():
            return "No readable text found in the document"
            
        prompt = f"""
        Analyze the following document content for Azure Landing Zone architecture planning:
        
        Document: {filename}
        Type: {file_type}
        Content: {text_content[:8000]}  # Limit content size
        
        Please provide insights for:
        1. Current architecture mentioned in the document
        2. Business requirements and objectives
        3. Compliance and regulatory requirements
        4. Security requirements
        5. Recommended Azure services and patterns
        6. Migration considerations
        7. Governance and operational requirements
        
        Format your response as a structured analysis for enterprise architecture planning.
        """
        
        result = gemini_model.generate_content(prompt)
        return result.text
        
    except Exception as e:
        logger.error(f"Error processing document {filename}: {e}")
        return f"Error processing document: {str(e)}"

def extract_pdf_text(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        import io
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""

def extract_excel_text(file_content: bytes) -> str:
    """Extract text from Excel file"""
    try:
        import io
        excel_file = io.BytesIO(file_content)
        workbook = openpyxl.load_workbook(excel_file)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting Excel text: {e}")
        return ""

def extract_pptx_text(file_content: bytes) -> str:
    """Extract text from PowerPoint file"""
    try:
        import io
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        text = ""
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PowerPoint text: {e}")
        return ""

def generate_ai_enhanced_recommendations(inputs: CustomerInputs, url_analysis: str = "", doc_analysis: str = "") -> str:
    """Generate AI-enhanced architecture recommendations using Gemini"""
    try:
        if not gemini_model:
            return "Standard recommendations (AI enhancement not available)"
            
        # Build context from inputs
        context = f"""
        Business Objective: {inputs.business_objective or 'Not specified'}
        Industry: {inputs.industry or 'General'}
        Organization Structure: {inputs.org_structure or 'Not specified'}
        Regulatory Requirements: {inputs.regulatory or 'Standard'}
        Security Requirements: {inputs.security_posture or 'Standard'}
        Scalability Requirements: {inputs.scalability or 'Standard'}
        Free-text Input: {inputs.free_text_input or 'None provided'}
        
        Selected Services:
        - Compute: {', '.join(inputs.compute_services or [])}
        - Network: {', '.join(inputs.network_services or [])}
        - Storage: {', '.join(inputs.storage_services or [])}
        - Database: {', '.join(inputs.database_services or [])}
        - Security: {', '.join(inputs.security_services or [])}
        """
        
        if url_analysis:
            context += f"\n\nURL Analysis Results:\n{url_analysis}"
            
        if doc_analysis:
            context += f"\n\nDocument Analysis Results:\n{doc_analysis}"
        
        prompt = f"""
        Based on the following customer requirements and analysis, provide comprehensive Azure Landing Zone architecture recommendations:
        
        {context}
        
        Please provide:
        1. Recommended Azure Landing Zone template (Enterprise Scale, Small Scale, etc.)
        2. Detailed architecture recommendations with specific Azure services
        3. Security and compliance strategy
        4. Network topology and connectivity recommendations
        5. Identity and access management strategy
        6. Operations and monitoring approach
        7. Cost optimization strategies
        8. Migration roadmap and phases
        9. Governance and policy recommendations
        10. Risk assessment and mitigation strategies
        
        Format your response as a comprehensive enterprise architecture document.
        """
        
        result = gemini_model.generate_content(prompt)
        return result.text
        
    except Exception as e:
        logger.error(f"Error generating AI recommendations: {e}")
        return f"Error generating AI recommendations: {str(e)}"

def analyze_free_text_requirements(free_text: str) -> dict:
    """Analyze free text input to extract comprehensive service requirements using AI"""
    try:
        if not gemini_model or not free_text:
            return {"services": [], "reasoning": "No analysis available"}
            
        prompt = """
        You are an ENTERPRISE AZURE SOLUTIONS ARCHITECT with deep expertise in designing production-ready, scalable, and secure Azure Landing Zone architectures. You have extensive experience with Azure Well-Architected Framework, Azure CAF (Cloud Adoption Framework), and enterprise-grade architectural patterns.

        User Requirement: "{}"

        ENTERPRISE ANALYSIS INSTRUCTIONS:
        
        1. COMPREHENSIVE BUSINESS ANALYSIS:
           - Analyze the business requirement to understand the complete technical and operational context
           - Consider enterprise-grade scalability, security, compliance, and governance needs
           - Think about future growth, disaster recovery, and business continuity requirements
           - Identify the appropriate Azure architectural pattern (e.g., N-tier, microservices, event-driven, etc.)

        2. INTELLIGENT SERVICE SELECTION:
           - Select ONLY services that are truly needed for the specific requirement
           - DO NOT include common services like virtual_machines, virtual_network, storage_accounts, or key_vault by default
           - DO NOT include random or default services unless they are justified by the requirement
           - For minimal requirements, provide minimal service sets - avoid over-engineering
           - For each service selected, provide clear architectural reasoning
           - Consider service dependencies and integration patterns
           - Ensure services work together as a cohesive architectural solution

        3. ENTERPRISE ARCHITECTURE PATTERNS:
           - Apply proven enterprise patterns: Hub-and-Spoke networking, Microservices, Event-driven architecture
           - Consider cross-cutting concerns: Security, Monitoring, Logging, Backup, Disaster Recovery
           - Design for High Availability (99.9%+), Scalability, and Performance
           - Include appropriate governance and compliance services when needed

        4. CONNECTIVITY AND INTEGRATION:
           - Design proper service-to-service communication patterns
           - Consider network security groups, private endpoints, and secure communication
           - Plan for API management, service mesh, and integration patterns
           - Design proper data flow and processing pipelines

        5. SECURITY BY DESIGN:
           - Implement Zero Trust security model where appropriate
           - Consider identity and access management, data encryption, network security
           - Include security monitoring and threat detection capabilities
           - Plan for compliance requirements (GDPR, HIPAA, SOC 2, etc.)

        AVAILABLE AZURE SERVICES BY CATEGORY:
        
        Compute: virtual_machines, aks, app_services, functions, container_instances, service_fabric, batch
        Network: virtual_network, vpn_gateway, expressroute, load_balancer, application_gateway, firewall, waf, cdn, traffic_manager, virtual_wan
        Storage: storage_accounts, blob_storage, data_lake_storage, file_storage, disk_storage
        Database: sql_database, cosmos_db, mysql, postgresql, redis, synapse_dedicated_pools
        Security: key_vault, security_center, sentinel, azure_ad_b2c, azure_firewall, ddos_protection, active_directory
        Monitoring: azure_monitor, log_analytics, application_insights, azure_advisor
        Analytics: synapse_analytics, data_factory, databricks, stream_analytics, event_hubs, power_bi
        Integration: logic_apps, service_bus, event_grid, api_management
        DevOps: devops, pipelines, container_registry, azure_artifacts
        Backup: backup_vault, site_recovery, azure_backup

        RESPONSE FORMAT - Provide a DETAILED JSON response with:

        1. "services" - Array of Azure service keys that form a complete, enterprise-ready solution
        2. "reasoning" - Comprehensive architectural explanation (minimum 200 words) covering:
           - Why each service was selected and how it fits the architectural pattern
           - How services integrate and communicate with each other
           - Security, scalability, and operational considerations
           - Alternatives considered and why they were not chosen
           
        3. "architecture_pattern" - Detailed description of the overall architecture pattern and design principles
        4. "connectivity_requirements" - Specific connectivity patterns, network topology, and communication flows
        5. "security_considerations" - Comprehensive security architecture including identity, network, data, and application security
        6. "scalability_design" - How the architecture scales to meet demand
        7. "operational_excellence" - Monitoring, logging, automation, and maintenance considerations
        8. "cost_optimization" - Cost-effective design choices and optimization strategies
        9. "needs_confirmation" - Set to true if clarification is needed from the user
        10. "suggested_additions" - Additional services or patterns to consider for enhancement

        EXAMPLE APPROACH - DO NOT COPY THIS LIST, analyze the actual requirement:
        For requirement "I need a simple web application with database":
        {{
          "services": ["app_services", "sql_database"],
          "reasoning": "Selected App Services for hosting the web application as it provides managed hosting with automatic scaling capabilities. SQL Database chosen for relational data storage with built-in high availability and backup features. These two services form a complete solution for the stated simple web application requirement without unnecessary complexity.",
          "architecture_pattern": "Simple two-tier web application architecture with presentation layer and data layer",
          "connectivity_requirements": "App Service connects to SQL Database via connection string with SSL encryption",
          "security_considerations": "SQL Database configured with firewall rules, App Service uses managed identity for database access",
          "scalability_design": "App Service auto-scaling based on CPU/memory, SQL Database can scale compute and storage independently",
          "operational_excellence": "Built-in monitoring for both services, automated backups for SQL Database",
          "cost_optimization": "Use appropriate service tiers based on expected load, SQL Database DTU-based pricing for predictable workloads",
          "needs_confirmation": false,
          "suggested_additions": ["Azure Monitor for enhanced monitoring", "Application Gateway for advanced routing if needed"]
        }}

        CRITICAL REQUIREMENTS:
        - ABSOLUTELY DO NOT include default or commonly used services unless explicitly required
        - DO NOT include virtual_machines, virtual_network, storage_accounts, or key_vault unless the requirement specifically needs them
        - ONLY suggest services that are specifically needed for the stated requirement
        - If the requirement is minimal, suggest minimal services - do not over-engineer
        - Provide detailed architectural reasoning for every service selection
        - Ensure all services work together as an integrated solution
        - Consider enterprise-grade non-functional requirements only when relevant to the requirement
        - Return ONLY valid JSON format with no additional text

        """.format(free_text)
        
        result = gemini_model.generate_content(prompt)
        response_text = result.text.strip()
        
        # Try to extract JSON from the response
        import json
        import re
        
        # Look for JSON content
        json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
        if json_match:
            json_str = json_match.group()
            try:
                analysis = json.loads(json_str)
                # Ensure backward compatibility by adding missing fields
                if "architecture_pattern" not in analysis:
                    analysis["architecture_pattern"] = "Custom architecture"
                if "connectivity_requirements" not in analysis:
                    analysis["connectivity_requirements"] = "Standard Azure networking"
                if "security_considerations" not in analysis:
                    analysis["security_considerations"] = "Standard security practices"
                if "scalability_design" not in analysis:
                    analysis["scalability_design"] = "Standard scalability patterns"
                if "operational_excellence" not in analysis:
                    analysis["operational_excellence"] = "Standard monitoring and operations"
                if "cost_optimization" not in analysis:
                    analysis["cost_optimization"] = "Standard cost optimization"
                if "needs_confirmation" not in analysis:
                    analysis["needs_confirmation"] = False
                if "suggested_additions" not in analysis:
                    analysis["suggested_additions"] = []
                return analysis
            except json.JSONDecodeError:
                pass
        
        # Fallback parsing if JSON parsing fails
        logger.warning(f"Failed to parse AI response as JSON: {response_text}")
        return {
            "services": [], 
            "reasoning": "Could not parse AI analysis",
            "architecture_pattern": "Unknown",
            "connectivity_requirements": "Not specified",
            "security_considerations": "Not specified",
            "scalability_design": "Not specified",
            "operational_excellence": "Not specified",
            "cost_optimization": "Not specified",
            "needs_confirmation": True,
            "suggested_additions": []
        }
        
    except Exception as e:
        logger.error(f"Error analyzing free text requirements: {e}")
        return {
            "services": [], 
            "reasoning": f"Analysis error: {str(e)}",
            "architecture_pattern": "Error",
            "connectivity_requirements": "Error",
            "security_considerations": "Error",
            "scalability_design": "Error",
            "operational_excellence": "Error",
            "cost_optimization": "Error",
            "needs_confirmation": True,
            "suggested_additions": []
        }

def validate_customer_inputs(inputs: CustomerInputs) -> None:
    """Validate customer inputs to prevent potential errors"""
    # Check for extremely long strings that might cause issues
    string_fields = [
        inputs.business_objective, inputs.regulatory, inputs.industry,
        inputs.org_structure, inputs.governance, inputs.identity,
        inputs.connectivity, inputs.network_model, inputs.ip_strategy,
        inputs.security_zone, inputs.security_posture, inputs.key_vault,
        inputs.threat_protection, inputs.workload, inputs.architecture_style,
        inputs.scalability, inputs.ops_model, inputs.monitoring, inputs.backup,
        inputs.topology_pattern, inputs.migration_scope, inputs.cost_priority, inputs.iac,
        inputs.url_input
    ]
    
    for field in string_fields:
        if field and len(field) > 1000:  # Reasonable limit for most fields
            raise ValueError(f"Input field too long: {len(field)} characters (max 1000)")
    
    # Special validation for free-text input (allowing more characters)
    if inputs.free_text_input and len(inputs.free_text_input) > 10000:
        raise ValueError(f"Free text input too long: {len(inputs.free_text_input)} characters (max 10000)")
    
    # Check service lists for reasonable sizes
    service_lists = [
        inputs.compute_services, inputs.network_services, inputs.storage_services,
        inputs.database_services, inputs.security_services, inputs.monitoring_services,
        inputs.ai_services, inputs.analytics_services, inputs.integration_services,
        inputs.devops_services, inputs.backup_services
    ]
    
    for service_list in service_lists:
        if service_list and len(service_list) > 50:  # Reasonable limit
            raise ValueError(f"Too many services selected: {len(service_list)} (max 50)")
    
    # Validate URL format if provided
    if inputs.url_input:
        if not inputs.url_input.startswith(('http://', 'https://')):
            raise ValueError("URL must start with http:// or https://")
    
    # Validate uploaded files info
    if inputs.uploaded_files_info:
        if len(inputs.uploaded_files_info) > 10:  # Reasonable limit
            raise ValueError(f"Too many uploaded files: {len(inputs.uploaded_files_info)} (max 10)")

def generate_azure_architecture_diagram(inputs: CustomerInputs, output_dir: str = None, format: str = "png") -> str:
    """Generate Azure architecture diagram using the Python Diagrams library with proper Azure icons"""
    
    logger.info("Starting Azure architecture diagram generation")
    
    try:
        # Validate inputs first
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Get safe output directory
        if output_dir is None:
            output_dir = get_safe_output_directory()
        
        # Clean up old files to prevent disk space issues
        cleanup_old_files(output_dir)
        
        # Verify Graphviz availability before proceeding
        try:
            result = subprocess.run(['dot', '-V'], capture_output=True, text=True, timeout=10)
            if result.returncode != 0:
                raise Exception(f"Graphviz 'dot' command failed with return code {result.returncode}. stderr: {result.stderr}")
            logger.info(f"Graphviz version: {result.stderr.strip()}")
        except subprocess.TimeoutExpired:
            raise Exception("Graphviz 'dot' command timed out. Graphviz may be unresponsive.")
        except FileNotFoundError:
            raise Exception("Graphviz is not installed or not accessible. Please install Graphviz: sudo apt-get install -y graphviz graphviz-dev")
        except subprocess.SubprocessError as e:
            raise Exception(f"Graphviz check failed: {str(e)}. Please install Graphviz: sudo apt-get install -y graphviz graphviz-dev")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = str(uuid.uuid4())[:8]  # Use first 8 chars of UUID for uniqueness
        filename = f"azure_landing_zone_{timestamp}_{unique_id}"
        filepath = os.path.join(output_dir, filename)
        
        logger.info(f"Generating diagram with filename: {filename}")
        
        # Verify output directory is writable
        if not os.access(output_dir, os.W_OK):
            raise Exception(f"Output directory {output_dir} is not writable")
        
        # Determine organization template
        template = generate_architecture_template(inputs)
        org_name = inputs.org_structure or "Enterprise"
        
        logger.info(f"Using template: {template['template']['name']}")
        
        try:
            # Set the format based on the requested output format
            output_format = "svg" if format.lower() == "svg" else "png"
            
            with Diagram(
                f"Azure Landing Zone - {template['template']['name']}", 
                filename=filepath, 
                show=False, 
                direction="TB",
                outformat=output_format,
                graph_attr={
                    "fontsize": "16",
                    "fontname": "Arial",
                    "rankdir": "TB",
                    "nodesep": "1.0",
                    "ranksep": "1.5",
                    "bgcolor": "#ffffff",
                    "margin": "0.5"
                },
                node_attr={
                    "fontsize": "12",
                    "fontname": "Arial"
                },
                edge_attr={
                    "fontsize": "10",
                    "fontname": "Arial"
                }
            ):
                
                logger.info("Creating diagram structure...")
                
                # Only create management groups if governance services are explicitly selected
                created_resources = []
                
                # Identity and Security Services - ONLY if explicitly selected or AI recommended
                identity_resources = []
                if (inputs.security_services and ("active_directory" in inputs.security_services or "key_vault" in inputs.security_services)) or \
                   (template.get("ai_services") and any(service in ["active_directory", "key_vault"] for service in template["ai_services"])):
                    
                    with Cluster("Identity & Security", graph_attr={"bgcolor": "#e8f4f8", "style": "rounded"}):
                        if (inputs.security_services and "active_directory" in inputs.security_services) or \
                           (template.get("ai_services") and "active_directory" in template["ai_services"]):
                            aad = ActiveDirectory("Azure Active Directory")
                            identity_resources.append(aad)
                            created_resources.append(aad)
                            
                        if (inputs.security_services and "key_vault" in inputs.security_services) or \
                           (template.get("ai_services") and "key_vault" in template["ai_services"]):
                            key_vault = KeyVaults("Key Vault")
                            identity_resources.append(key_vault)
                            created_resources.append(key_vault)
                            
                        if inputs.security_services and "security_center" in inputs.security_services:
                            sec_center = SecurityCenter("Security Center")
                            identity_resources.append(sec_center)
                            created_resources.append(sec_center)
                            
                        if inputs.security_services and "sentinel" in inputs.security_services:
                            sentinel = Sentinel("Sentinel")
                            identity_resources.append(sentinel)
                            created_resources.append(sentinel)
                
                # Azure Subscription - ONLY if we actually have services to put in it
                subscription = None
                total_services = sum([
                    len(inputs.compute_services or []),
                    len(inputs.network_services or []),
                    len(inputs.storage_services or []),
                    len(inputs.database_services or []),
                    len(inputs.security_services or []),
                    len(inputs.monitoring_services or []),
                    len(inputs.ai_services or []),
                    len(inputs.analytics_services or []),
                    len(inputs.integration_services or []),
                    len(inputs.devops_services or []),
                    len(inputs.backup_services or []),
                    len(template.get("ai_services", []))
                ])
                
                if total_services > 0:
                    with Cluster("Azure Subscription", graph_attr={"bgcolor": "#f0f8ff", "style": "rounded"}):
                        subscription = Subscriptions("Azure Subscription")
                        created_resources.append(subscription)
                
                # Create network services only if explicitly selected
                network_resources_by_type = {}
                if inputs.network_services:
                    with Cluster("Network Architecture", graph_attr={"bgcolor": "#f0fff0", "style": "rounded"}):
                        network_resources = []
                        for service in inputs.network_services:
                            if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                service_name = AZURE_SERVICES_MAPPING[service]["name"]
                                network_resource = diagram_class(service_name)
                                network_resources.append(network_resource)
                                created_resources.append(network_resource)
                                
                                # Track network resources by type
                                if service not in network_resources_by_type:
                                    network_resources_by_type[service] = []
                                network_resources_by_type[service].append(network_resource)
                        
                        # Intelligent network connections based on Azure patterns
                        _add_intelligent_network_connections(network_resources, network_resources_by_type)
                
                # Create other service clusters and get resource tracking
                all_resources, resources_by_type = _add_selected_service_clusters(inputs, created_resources)
                
                # Merge network resources into the overall tracking
                resources_by_type.update(network_resources_by_type)
                
                # Add intelligent connections between all services
                _add_intelligent_service_connections(inputs, resources_by_type, template)
                
                # If no services are selected at all, create a placeholder message
                if len(all_resources) == 0:
                    from diagrams.generic.blank import Blank
                    placeholder = Blank("No Services Selected\n\nPlease specify Azure services\nto generate architecture")
                    all_resources.append(placeholder)
                
                logger.info("Diagram structure and intelligent connections created successfully")
        
        except Exception as e:
            logger.error(f"Error during diagram creation: {str(e)}")
            logger.error(traceback.format_exc())
            raise Exception(f"Error generating Azure architecture diagram: {str(e)}")
        
        # Return the file path of the generated diagram
        if format.lower() == "svg":
            # Check for SVG file generated directly by diagrams library
            svg_path = f"{filepath}.svg"
            if os.path.exists(svg_path):
                file_size = os.path.getsize(svg_path)
                logger.info(f"SVG diagram generated successfully: {svg_path} (size: {file_size} bytes)")
                return svg_path
            else:
                # Fallback: try to generate SVG using dot command from gv file
                dot_path = f"{filepath}.gv" 
                if os.path.exists(dot_path):
                    try:
                        # Convert dot file to SVG
                        result = subprocess.run(['dot', '-Tsvg', dot_path, '-o', svg_path], 
                                              capture_output=True, text=True, timeout=30)
                        if result.returncode != 0:
                            raise Exception(f"SVG generation failed: {result.stderr}")
                        
                        if os.path.exists(svg_path):
                            file_size = os.path.getsize(svg_path)
                            logger.info(f"SVG diagram generated successfully via dot: {svg_path} (size: {file_size} bytes)")
                            return svg_path
                        else:
                            raise Exception(f"SVG generation failed - file not found: {svg_path}")
                    except subprocess.TimeoutExpired:
                        raise Exception("SVG generation timed out")
                    except Exception as e:
                        raise Exception(f"Failed to generate SVG: {str(e)}")
                else:
                    raise Exception(f"Neither SVG nor dot file found: {svg_path}, {dot_path}")
        else:
            # Default PNG generation
            png_path = f"{filepath}.png"
            if os.path.exists(png_path):
                file_size = os.path.getsize(png_path)
                logger.info(f"Diagram generated successfully: {png_path} (size: {file_size} bytes)")
                return png_path
            else:
                raise Exception(f"Diagram generation failed - PNG file not found: {png_path}")
            
    except Exception as e:
        logger.error(f"Failed to generate Azure architecture diagram: {str(e)}")
        logger.error(traceback.format_exc())
        raise

def generate_simple_svg_diagram(inputs: CustomerInputs) -> str:
    """Generate a minimal SVG diagram based only on selected services"""
    
    template = generate_architecture_template(inputs)
    template_name = template['template']['name']
    
    # Create a simple SVG representation
    svg_width = 800
    svg_height = 600
    
    svg_content = f'''<svg width="{svg_width}" height="{svg_height}" xmlns="http://www.w3.org/2000/svg">
    <defs>
        <style>
            .title {{ font-family: Arial, sans-serif; font-size: 18px; font-weight: bold; fill: #0078d4; }}
            .group-title {{ font-family: Arial, sans-serif; font-size: 14px; font-weight: bold; fill: #323130; }}
            .service {{ font-family: Arial, sans-serif; font-size: 12px; fill: #605e5c; }}
            .subscription {{ fill: #e1f5fe; stroke: #0078d4; stroke-width: 2; }}
            .resource-group {{ fill: #f3e5f5; stroke: #6b69d6; stroke-width: 2; }}
            .service-box {{ fill: #fff3e0; stroke: #d83b01; stroke-width: 1; cursor: pointer; }}
            .service-box:hover {{ fill: #ffebdd; }}
            .network-box {{ fill: #e8f5e8; stroke: #107c10; stroke-width: 1; cursor: pointer; }}
            .network-box:hover {{ fill: #f3fdf3; }}
        </style>
    </defs>
    
    <!-- Background -->
    <rect width="100%" height="100%" fill="#f8f9fa"/>
    
    <!-- Title -->
    <text x="400" y="30" class="title" text-anchor="middle">Azure Architecture - {template_name}</text>
    
    <!-- Azure Subscription -->
    <rect x="100" y="80" width="600" height="450" class="subscription" rx="5"/>
    <text x="110" y="100" class="group-title">Azure Subscription</text>
    
    <!-- Resource Group -->
    <rect x="120" y="120" width="560" height="380" class="resource-group" rx="5"/>
    <text x="130" y="140" class="group-title">Resource Group</text>'''
    
    # Add selected services
    y_offset = 170
    x_start = 150
    services_added = []
    
    if inputs.compute_services:
        svg_content += f'''
    <!-- Compute Services -->
    <rect x="{x_start}" y="{y_offset}" width="200" height="100" class="service-box" rx="5"/>
    <text x="{x_start + 10}" y="{y_offset + 20}" class="group-title">Compute</text>'''
        
        for i, service in enumerate(inputs.compute_services[:2]):  # Max 2 services
            service_info = AZURE_SERVICES_MAPPING.get(service, {"name": service, "icon": "🖥️"})
            svg_content += f'''
    <rect x="{x_start + 10}" y="{y_offset + 30 + (i * 30)}" width="180" height="25" class="service-box" rx="3"/>
    <text x="{x_start + 20}" y="{y_offset + 47 + (i * 30)}" class="service" font-size="11">{service_info["icon"]} {service_info["name"][:15]}</text>'''
        services_added.append("compute")
    
    if inputs.network_services:
        x_pos = x_start + 220 if "compute" in services_added else x_start
        svg_content += f'''
    <!-- Network Services -->
    <rect x="{x_pos}" y="{y_offset}" width="200" height="100" class="network-box" rx="5"/>
    <text x="{x_pos + 10}" y="{y_offset + 20}" class="group-title">Network</text>'''
        
        for i, service in enumerate(inputs.network_services[:2]):  # Max 2 services
            service_info = AZURE_SERVICES_MAPPING.get(service, {"name": service, "icon": "🌐"})
            svg_content += f'''
    <rect x="{x_pos + 10}" y="{y_offset + 30 + (i * 30)}" width="180" height="25" class="network-box" rx="3"/>
    <text x="{x_pos + 20}" y="{y_offset + 47 + (i * 30)}" class="service" font-size="11">{service_info["icon"]} {service_info["name"][:15]}</text>'''
        services_added.append("network")
    
    # Add other service types if needed
    if inputs.storage_services:
        x_pos = x_start + (220 * len(services_added))
        if x_pos < 500:  # Only if there's space
            svg_content += f'''
    <!-- Storage Services -->
    <rect x="{x_pos}" y="{y_offset}" width="200" height="100" class="service-box" rx="5"/>
    <text x="{x_pos + 10}" y="{y_offset + 20}" class="group-title">Storage</text>'''
            
            for i, service in enumerate(inputs.storage_services[:2]):
                service_info = AZURE_SERVICES_MAPPING.get(service, {"name": service, "icon": "💾"})
                svg_content += f'''
    <rect x="{x_pos + 10}" y="{y_offset + 30 + (i * 30)}" width="180" height="25" class="service-box" rx="3"/>
    <text x="{x_pos + 20}" y="{y_offset + 47 + (i * 30)}" class="service" font-size="11">{service_info["icon"]} {service_info["name"][:15]}</text>'''
    
    # If no services selected, show message
    if not any([inputs.compute_services, inputs.network_services, inputs.storage_services,
                inputs.database_services, inputs.security_services, inputs.monitoring_services]):
        svg_content += f'''
    <!-- No Services Message -->
    <rect x="{x_start}" y="{y_offset}" width="400" height="80" fill="#f8f8f8" stroke="#ddd" stroke-width="1" rx="5"/>
    <text x="{x_start + 200}" y="{y_offset + 35}" class="group-title" text-anchor="middle">No Services Selected</text>
    <text x="{x_start + 200}" y="{y_offset + 55}" class="service" text-anchor="middle">Please select Azure services to generate architecture</text>'''
    
    svg_content += '''
    <!-- Footer -->
    <text x="400" y="570" class="service" text-anchor="middle" fill="#8a8886">Generated by Azure Landing Zone Agent - Minimal Mode</text>
</svg>'''
    
    return svg_content

def _add_service_clusters(inputs: CustomerInputs, prod_vnet, workloads_mg):
    """Helper method to add service clusters to avoid code duplication"""
    try:
        # Compute and Application Services
        if inputs.compute_services or inputs.workload:
            with Cluster("Compute & Applications", graph_attr={"bgcolor": "#fff8dc", "style": "rounded"}):
                compute_services = []
                
                # Add selected compute services
                if inputs.compute_services:
                    for service in inputs.compute_services:
                        if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                            diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                            service_name = AZURE_SERVICES_MAPPING[service]["name"]
                            compute_services.append(diagram_class(service_name))
                
                # Fallback to workload if no specific compute services
                if not compute_services and inputs.workload:
                    if inputs.workload in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[inputs.workload]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[inputs.workload]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[inputs.workload]["name"]
                        compute_services.append(diagram_class(service_name))
                
                # Default to App Services if nothing specified
                if not compute_services:
                    compute_services.append(AppServices("Azure App Services"))
                
                # Connect compute services to production VNet
                for cs in compute_services:
                    prod_vnet >> cs
                    workloads_mg >> cs
        
        # Storage Services
        if inputs.storage_services:
            with Cluster("Storage & Data", graph_attr={"bgcolor": "#f5f5dc", "style": "rounded"}):
                storage_services = []
                for service in inputs.storage_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        storage_services.append(diagram_class(service_name))
                
                if not storage_services:
                    storage_services.append(StorageAccounts("Storage Accounts"))
                
                # Connect storage to production VNet and workloads
                for ss in storage_services:
                    prod_vnet >> ss
                    workloads_mg >> ss
        
        # Database Services
        if inputs.database_services:
            with Cluster("Databases", graph_attr={"bgcolor": "#e6f3ff", "style": "rounded"}):
                database_services = []
                for service in inputs.database_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        database_services.append(diagram_class(service_name))
                
                # Connect databases to production VNet and workloads
                for ds in database_services:
                    prod_vnet >> ds
                    workloads_mg >> ds
        
        # Analytics Services
        if inputs.analytics_services:
            with Cluster("Analytics & AI", graph_attr={"bgcolor": "#f0e6ff", "style": "rounded"}):
                analytics_services = []
                for service in inputs.analytics_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        analytics_services.append(diagram_class(service_name))
                
                # Connect analytics to production VNet and workloads
                for as_service in analytics_services:
                    prod_vnet >> as_service
                    workloads_mg >> as_service
        
        # Integration Services
        if inputs.integration_services:
            with Cluster("Integration", graph_attr={"bgcolor": "#fff0e6", "style": "rounded"}):
                integration_services = []
                for service in inputs.integration_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        integration_services.append(diagram_class(service_name))
                
                # Connect integration services to production VNet and workloads
                for is_service in integration_services:
                    prod_vnet >> is_service
                    workloads_mg >> is_service
        
        # DevOps Services  
        if inputs.devops_services:
            with Cluster("DevOps & Automation", graph_attr={"bgcolor": "#f5f5f5", "style": "rounded"}):
                devops_services = []
                for service in inputs.devops_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        devops_services.append(diagram_class(service_name))
                
                # Connect DevOps services to management
                for ds in devops_services:
                    workloads_mg >> ds
                    
    except Exception as e:
        logger.warning(f"Error adding service clusters: {str(e)}")
        # Don't fail the entire diagram generation for service cluster issues


def _add_intelligent_network_connections(network_resources: list, network_resources_by_type: dict):
    """Add intelligent connections between network services based on Azure patterns"""
    try:
        # Hub-spoke pattern: VNet as hub, connect others to it
        vnets = network_resources_by_type.get("virtual_network", [])
        if len(vnets) > 1:
            hub_vnet = vnets[0]  # First VNet as hub
            for spoke_vnet in vnets[1:]:
                hub_vnet >> Edge(style="dashed", color="blue", label="VNet Peering") >> spoke_vnet
        
        # Application Gateway -> Load Balancer pattern
        app_gateways = network_resources_by_type.get("application_gateway", [])
        load_balancers = network_resources_by_type.get("load_balancer", [])
        for ag in app_gateways:
            for lb in load_balancers:
                ag >> Edge(color="green", label="HTTP/HTTPS") >> lb
        
        # Firewall -> Application Gateway pattern (secure traffic flow)
        firewalls = network_resources_by_type.get("firewall", [])
        for fw in firewalls:
            for ag in app_gateways:
                fw >> Edge(color="red", label="Filtered Traffic") >> ag
        
        # VPN/ExpressRoute -> VNet pattern
        vpn_gateways = network_resources_by_type.get("vpn_gateway", [])
        expressroutes = network_resources_by_type.get("expressroute", [])
        
        for vpn in vpn_gateways:
            for vnet in vnets:
                vpn >> Edge(color="orange", label="Site-to-Site") >> vnet
        
        for er in expressroutes:
            for vnet in vnets:
                er >> Edge(color="purple", label="Private Connection") >> vnet
                
        logger.debug("Intelligent network connections added")
        
    except Exception as e:
        logger.warning(f"Error adding intelligent network connections: {e}")


def _add_intelligent_service_connections(inputs: CustomerInputs, all_resources: dict, template: dict):
    """Add intelligent connections between services based on Azure architecture patterns"""
    try:
        logger.info("Adding intelligent service connections based on architecture patterns")
        
        # Get AI analysis for connectivity guidance
        connectivity_guidance = ""
        if hasattr(template, 'get') and template.get("connectivity_requirements"):
            connectivity_guidance = template["connectivity_requirements"]
        
        # Define enterprise-grade connection patterns based on Azure architectural best practices
        connection_patterns = {
            # Web tier patterns - Front-end to application connectivity
            "web_tier": {
                "sources": ["application_gateway", "load_balancer", "traffic_manager", "cdn"],
                "targets": ["app_services", "virtual_machines", "aks", "functions"]
            },
            
            # Data flow patterns - Application to data tier
            "data_flow": {
                "sources": ["app_services", "virtual_machines", "aks", "functions"],
                "targets": ["sql_database", "cosmos_db", "mysql", "postgresql", "storage_accounts", "blob_storage", "redis"]
            },
            
            # Network infrastructure patterns
            "network_flow": {
                "sources": ["vpn_gateway", "expressroute", "virtual_wan"],
                "targets": ["virtual_network", "firewall", "application_gateway"]
            },
            
            # Security and secrets management patterns
            "security_flow": {
                "sources": ["key_vault", "active_directory"],
                "targets": ["app_services", "virtual_machines", "aks", "functions", "sql_database", "cosmos_db"]
            },
            
            # Monitoring and observability patterns
            "monitoring_flow": {
                "sources": ["azure_monitor", "log_analytics", "application_insights"],
                "targets": ["app_services", "virtual_machines", "aks", "sql_database", "cosmos_db", "functions", "storage_accounts"]
            },
            
            # API and integration patterns
            "integration_flow": {
                "sources": ["api_management"],
                "targets": ["app_services", "functions", "logic_apps", "service_bus", "event_grid", "aks"]
            },
            
            # DevOps and deployment patterns
            "devops_flow": {
                "sources": ["devops", "pipelines", "container_registry", "azure_artifacts"],
                "targets": ["aks", "app_services", "virtual_machines", "functions"]
            },
            
            # Backup and disaster recovery patterns
            "backup_flow": {
                "sources": ["backup_vault", "site_recovery"],
                "targets": ["virtual_machines", "sql_database", "storage_accounts", "cosmos_db"]
            },
            
            # Caching patterns for performance optimization
            "cache_flow": {
                "sources": ["redis", "cdn"],
                "targets": ["app_services", "aks", "functions"]
            },
            
            # Event-driven architecture patterns
            "event_flow": {
                "sources": ["event_grid", "service_bus", "event_hubs"],
                "targets": ["functions", "logic_apps", "app_services", "aks"]
            },
            
            # Analytics and data processing patterns
            "analytics_flow": {
                "sources": ["data_factory", "databricks", "synapse_analytics", "stream_analytics"],
                "targets": ["data_lake_storage", "storage_accounts", "sql_database", "cosmos_db", "event_hubs"]
            },
            
            # Content delivery patterns
            "content_delivery": {
                "sources": ["cdn", "storage_accounts"],
                "targets": ["application_gateway", "app_services"]
            }
        }
        
        # Apply connection patterns with intelligent relationship detection
        for pattern_name, pattern in connection_patterns.items():
            for source_type in pattern["sources"]:
                for target_type in pattern["targets"]:
                    # Find matching resources
                    source_resources = all_resources.get(source_type, [])
                    target_resources = all_resources.get(target_type, [])
                    
                    # Create connections between matching resources
                    for source in source_resources:
                        for target in target_resources:
                            if source != target:  # Don't connect to self
                                try:
                                    # Create connection with appropriate style
                                    connection_style = _get_connection_style(pattern_name)
                                    
                                    # Determine if this should be a bidirectional connection
                                    is_bidirectional = _is_bidirectional_connection(source_type, target_type, pattern_name)
                                    
                                    if is_bidirectional:
                                        # Create bidirectional connection
                                        connection_style["dir"] = "both"
                                        connection_style["arrowtail"] = connection_style["arrowhead"]
                                        source >> Edge(**connection_style) >> target
                                    else:
                                        # Create unidirectional connection
                                        connection_style["dir"] = "forward"
                                        source >> Edge(**connection_style) >> target
                                    
                                    logger.debug(f"Connected {source_type} -> {target_type} ({pattern_name}) - {'Bidirectional' if is_bidirectional else 'Unidirectional'}")
                                except Exception as conn_error:
                                    logger.warning(f"Failed to connect {source_type} -> {target_type}: {conn_error}")
        
        # Add hub-spoke network connections if virtual networks exist
        vnets = all_resources.get("virtual_network", [])
        if len(vnets) > 1:
            # Connect first VNet (hub) to all others (spokes)
            hub = vnets[0]
            for spoke in vnets[1:]:
                try:
                    hub >> Edge(style="dashed", color="blue", label="Peering") >> spoke
                    logger.debug("Added hub-spoke VNet peering connection")
                except Exception as e:
                    logger.warning(f"Failed to add VNet peering: {e}")
                    
        logger.info("Intelligent service connections added successfully")
        
    except Exception as e:
        logger.warning(f"Error adding intelligent service connections: {str(e)}")
        # Don't fail diagram generation if connections fail


def _get_connection_style(pattern_name: str) -> dict:
    """Get enterprise-grade connection style based on architectural pattern type
    
    Implements professional architectural diagram conventions:
    - Solid lines: Direct communication/data flow
    - Dashed lines: Indirect/asynchronous communication
    - Dotted lines: Monitoring/logging/non-critical flows
    - Bold lines: Primary/critical data paths
    - Arrows: Directional flow indication
    - Colors: Semantic meaning for different flow types
    """
    styles = {
        # Web tier - Solid green arrows for HTTP/HTTPS traffic flow
        "web_tier": {
            "color": "#2E8B57",  # Sea Green
            "style": "solid",
            "arrowhead": "normal",
            "arrowsize": "1.0",
            "penwidth": "2.0",
            "label": "HTTP/HTTPS",
            "fontcolor": "#2E8B57",
            "fontsize": "10"
        },
        
        # Data flow - Bold blue lines for critical data pathways
        "data_flow": {
            "color": "#1E90FF",  # Dodger Blue
            "style": "solid",
            "arrowhead": "normal",
            "arrowsize": "1.2",
            "penwidth": "3.0",
            "label": "Data Flow",
            "fontcolor": "#1E90FF",
            "fontsize": "10"
        },
        
        # Network flow - Bold orange lines for network connectivity
        "network_flow": {
            "color": "#FF8C00",  # Dark Orange
            "style": "solid",
            "arrowhead": "diamond",
            "arrowsize": "1.0",
            "penwidth": "2.5",
            "label": "Network",
            "fontcolor": "#FF8C00",
            "fontsize": "10"
        },
        
        # Security flow - Dashed red lines for security/secrets management
        "security_flow": {
            "color": "#DC143C",  # Crimson
            "style": "dashed",
            "arrowhead": "vee",
            "arrowsize": "0.8",
            "penwidth": "2.0",
            "label": "Secure Access",
            "fontcolor": "#DC143C",
            "fontsize": "9"
        },
        
        # Monitoring flow - Dotted purple lines for observability
        "monitoring_flow": {
            "color": "#9932CC",  # Dark Orchid
            "style": "dotted",
            "arrowhead": "open",
            "arrowsize": "0.7",
            "penwidth": "1.5",
            "label": "Telemetry",
            "fontcolor": "#9932CC",
            "fontsize": "9"
        },
        
        # Integration flow - Solid teal lines for API communication
        "integration_flow": {
            "color": "#008B8B",  # Dark Cyan
            "style": "solid",
            "arrowhead": "normal",
            "arrowsize": "1.0",
            "penwidth": "2.0",
            "label": "API Calls",
            "fontcolor": "#008B8B",
            "fontsize": "10"
        },
        
        # DevOps flow - Dashed gray lines for deployment pipelines
        "devops_flow": {
            "color": "#696969",  # Dim Gray
            "style": "dashed",
            "arrowhead": "box",
            "arrowsize": "0.9",
            "penwidth": "1.8",
            "label": "CI/CD",
            "fontcolor": "#696969",
            "fontsize": "9"
        },
        
        # Backup flow - Dotted brown lines for backup operations
        "backup_flow": {
            "color": "#8B4513",  # Saddle Brown
            "style": "dotted",
            "arrowhead": "tee",
            "arrowsize": "0.8",
            "penwidth": "1.5",
            "label": "Backup",
            "fontcolor": "#8B4513",
            "fontsize": "9"
        },
        
        # Cache flow - Dashed light blue for caching patterns
        "cache_flow": {
            "color": "#87CEEB",  # Sky Blue
            "style": "dashed",
            "arrowhead": "curve",
            "arrowsize": "0.8",
            "penwidth": "1.5",
            "label": "Cache",
            "fontcolor": "#4682B4",
            "fontsize": "9"
        },
        
        # Event flow - Wavy lines for event-driven communication
        "event_flow": {
            "color": "#FFD700",  # Gold
            "style": "solid",
            "arrowhead": "dot",
            "arrowsize": "1.0",
            "penwidth": "2.0",
            "label": "Events",
            "fontcolor": "#DAA520",
            "fontsize": "10"
        }
    }
    
    # Default style for unknown patterns
    default_style = {
        "color": "#2F4F4F",  # Dark Slate Gray
        "style": "solid",
        "arrowhead": "normal",
        "arrowsize": "0.8",
        "penwidth": "1.0",
        "fontcolor": "#2F4F4F",
        "fontsize": "9"
    }
    
    return styles.get(pattern_name, default_style)


def _is_bidirectional_connection(source_type: str, target_type: str, pattern_name: str) -> bool:
    """Determine if a connection should be bidirectional based on architectural patterns"""
    
    # Bidirectional connection patterns based on Azure architectural best practices
    bidirectional_patterns = {
        # Network infrastructure typically has bidirectional communication
        "network_flow": [
            ("vpn_gateway", "virtual_network"),
            ("expressroute", "virtual_network"),
            ("virtual_network", "firewall"),
            ("virtual_wan", "virtual_network")
        ],
        
        # Cache flows are typically bidirectional (read/write)
        "cache_flow": [
            ("redis", "app_services"),
            ("redis", "aks"),
            ("redis", "functions")
        ],
        
        # Event flows can be bidirectional for event sourcing patterns
        "event_flow": [
            ("service_bus", "functions"),
            ("service_bus", "logic_apps"),
            ("event_grid", "functions"),
            ("event_hubs", "functions")
        ],
        
        # API integration often involves bidirectional communication
        "integration_flow": [
            ("api_management", "app_services"),
            ("api_management", "functions"),
            ("service_bus", "aks")
        ],
        
        # Analytics patterns can involve bidirectional data flow
        "analytics_flow": [
            ("databricks", "data_lake_storage"),
            ("synapse_analytics", "storage_accounts"),
            ("stream_analytics", "event_hubs")
        ]
    }
    
    # Check if this specific connection should be bidirectional
    pattern_connections = bidirectional_patterns.get(pattern_name, [])
    
    for source, target in pattern_connections:
        if (source_type == source and target_type == target) or \
           (source_type == target and target_type == source):
            return True
    
    # Special cases for specific service combinations regardless of pattern
    bidirectional_service_pairs = [
        ("virtual_network", "firewall"),
        ("application_gateway", "aks"),
        ("load_balancer", "virtual_machines"),
        ("redis", "app_services"),
        ("service_bus", "functions"),
        ("event_grid", "logic_apps")
    ]
    
    for source, target in bidirectional_service_pairs:
        if (source_type == source and target_type == target) or \
           (source_type == target and target_type == source):
            return True
    
    return False


def _add_selected_service_clusters(inputs: CustomerInputs, existing_resources: list):
    """Add only explicitly selected services without defaults or hardcoded assumptions"""
    try:
        all_created_resources = existing_resources.copy()
        
        # Track resources by type for intelligent connections
        resources_by_type = {}
        
        # Compute Services - only if explicitly selected
        if inputs.compute_services:
            with Cluster("Compute Services", graph_attr={"bgcolor": "#fff8dc", "style": "rounded"}):
                for service in inputs.compute_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Storage Services - only if explicitly selected  
        if inputs.storage_services:
            with Cluster("Storage Services", graph_attr={"bgcolor": "#f5f5dc", "style": "rounded"}):
                for service in inputs.storage_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Database Services - only if explicitly selected
        if inputs.database_services:
            with Cluster("Database Services", graph_attr={"bgcolor": "#e6f3ff", "style": "rounded"}):
                for service in inputs.database_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Security Services - only if explicitly selected
        if inputs.security_services:
            with Cluster("Security Services", graph_attr={"bgcolor": "#ffe6e6", "style": "rounded"}):
                for service in inputs.security_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Monitoring Services - only if explicitly selected
        if inputs.monitoring_services:
            with Cluster("Monitoring Services", graph_attr={"bgcolor": "#f0f8ff", "style": "rounded"}):
                for service in inputs.monitoring_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # AI/ML Services - only if explicitly selected
        if inputs.ai_services:
            with Cluster("AI & Machine Learning", graph_attr={"bgcolor": "#f0fff0", "style": "rounded"}):
                for service in inputs.ai_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Analytics Services - only if explicitly selected
        if inputs.analytics_services:
            with Cluster("Analytics Services", graph_attr={"bgcolor": "#fff8f0", "style": "rounded"}):
                for service in inputs.analytics_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Integration Services - only if explicitly selected
        if inputs.integration_services:
            with Cluster("Integration Services", graph_attr={"bgcolor": "#f8f8ff", "style": "rounded"}):
                for service in inputs.integration_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # DevOps Services - only if explicitly selected
        if inputs.devops_services:
            with Cluster("DevOps Services", graph_attr={"bgcolor": "#f5f5f5", "style": "rounded"}):
                for service in inputs.devops_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Backup Services - only if explicitly selected
        if inputs.backup_services:
            with Cluster("Backup & Recovery", graph_attr={"bgcolor": "#fff0f5", "style": "rounded"}):
                for service in inputs.backup_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        resource = diagram_class(service_name)
                        all_created_resources.append(resource)
                        
                        # Track by service type
                        if service not in resources_by_type:
                            resources_by_type[service] = []
                        resources_by_type[service].append(resource)
        
        # Return both the list and the type mapping for intelligent connections
        return all_created_resources, resources_by_type
                        
    except Exception as e:
        logger.warning(f"Error adding selected service clusters: {str(e)}")
        # Don't fail the entire diagram generation for service cluster issues
        return existing_resources, {}


def generate_architecture_template(inputs: CustomerInputs) -> Dict[str, Any]:
    """Generate architecture template based on inputs with AI enhancement for free text"""
    
    # First, check if we should use AI analysis for service selection
    ai_services = []
    ai_reasoning = ""
    needs_confirmation = False
    suggested_additions = []
    
    # Count explicitly selected services
    total_selected_services = sum([
        len(inputs.compute_services or []),
        len(inputs.network_services or []),
        len(inputs.storage_services or []),
        len(inputs.database_services or []),
        len(inputs.security_services or []),
        len(inputs.monitoring_services or []),
        len(inputs.ai_services or []),
        len(inputs.analytics_services or []),
        len(inputs.integration_services or []),
        len(inputs.devops_services or []),
        len(inputs.backup_services or [])
    ])
    
    # Enhanced AI analysis fields
    architecture_pattern = ""
    connectivity_requirements = ""
    security_considerations = ""
    
    # Only use AI if user provided free text
    if inputs.free_text_input:
        logger.info("Using comprehensive AI analysis for architecture design based on free text input")
        ai_analysis = analyze_free_text_requirements(inputs.free_text_input)
        ai_services = ai_analysis.get("services", [])
        ai_reasoning = ai_analysis.get("reasoning", "")
        architecture_pattern = ai_analysis.get("architecture_pattern", "")
        connectivity_requirements = ai_analysis.get("connectivity_requirements", "")
        security_considerations = ai_analysis.get("security_considerations", "")
        needs_confirmation = ai_analysis.get("needs_confirmation", False)
        suggested_additions = ai_analysis.get("suggested_additions", [])
        
        logger.info(f"AI analysis completed: {len(ai_services)} services identified for {architecture_pattern} pattern")
    
    # Use minimal template by default - no complex management groups unless specifically requested
    minimal_template = {
        "name": "Minimal Azure Architecture",
        "management_groups": ["Root"],  # Only root, no complex hierarchy
        "subscriptions": ["Production"],  # Only one subscription
        "core_services": []  # No automatic core services
    }
    
    # Build architecture components with enhanced AI insights
    components = {
        "template": minimal_template,
        "identity": inputs.identity or None,
        "network_model": inputs.network_model or "simple",
        "workload": inputs.workload or None,
        "security": inputs.security_posture or None,
        "monitoring": inputs.monitoring or None,
        "governance": inputs.governance or None,
        "ai_services": ai_services,
        "ai_reasoning": ai_reasoning,
        "architecture_pattern": architecture_pattern,
        "connectivity_requirements": connectivity_requirements,
        "security_considerations": security_considerations,
        "needs_confirmation": needs_confirmation,
        "suggested_additions": suggested_additions
    }
    
    return components

def generate_professional_mermaid(inputs: CustomerInputs) -> str:
    """Generate minimal Mermaid diagram for Azure Landing Zone based only on selected services"""
    
    template = generate_architecture_template(inputs)
    
    lines = [
        "graph TB",
        "    subgraph \"Azure Architecture\"",
    ]
    
    # Check if we have any services selected (including AI suggested ones)
    ai_services = template.get("ai_services", [])
    total_selected_services = sum([
        len(inputs.compute_services or []),
        len(inputs.network_services or []),
        len(inputs.storage_services or []),
        len(inputs.database_services or []),
        len(inputs.security_services or []),
        len(inputs.monitoring_services or []),
        len(inputs.ai_services or []),
        len(inputs.analytics_services or []),
        len(inputs.integration_services or []),
        len(inputs.devops_services or []),
        len(inputs.backup_services or []),
        len(ai_services)
    ])
    
    # Only add structure if we have services
    if total_selected_services > 0:
        lines.extend([
            "        SUBSCRIPTION[\"📦 Azure Subscription\"]"
        ])
        
        # Add Resource Group only if we have resources
        lines.extend([
            "        RESOURCEGROUP[\"📁 Resource Group\"]",
            "        SUBSCRIPTION --> RESOURCEGROUP"
        ])
        
        service_count = 0
        
        # Add compute services if selected
        if inputs.compute_services or (ai_services and any(service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "compute" for service in ai_services)):
            lines.extend([
                "        subgraph \"Compute Services\""
            ])
            
            service_connections = []
            # Add explicitly selected compute services
            for service in (inputs.compute_services or []):
                if service in AZURE_SERVICES_MAPPING:
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"COMPUTE{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            # Add AI-suggested compute services
            for service in ai_services:
                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "compute":
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"COMPUTE{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            lines.extend(service_connections)
            lines.append("        end")
        
        # Add network services if selected
        if inputs.network_services or (ai_services and any(service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "network" for service in ai_services)):
            lines.extend([
                "        subgraph \"Network Services\""
            ])
            
            service_connections = []
            # Add explicitly selected network services
            for service in (inputs.network_services or []):
                if service in AZURE_SERVICES_MAPPING:
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"NETWORK{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            # Add AI-suggested network services
            for service in ai_services:
                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "network":
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"NETWORK{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            lines.extend(service_connections)
            lines.append("        end")
        
        # Add storage services if selected
        if inputs.storage_services or (ai_services and any(service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "storage" for service in ai_services)):
            lines.extend([
                "        subgraph \"Storage Services\""
            ])
            
            service_connections = []
            # Add explicitly selected storage services
            for service in (inputs.storage_services or []):
                if service in AZURE_SERVICES_MAPPING:
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"STORAGE{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            # Add AI-suggested storage services
            for service in ai_services:
                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "storage":
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"STORAGE{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            lines.extend(service_connections)
            lines.append("        end")
        
        # Add database services if selected
        if inputs.database_services or (ai_services and any(service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "database" for service in ai_services)):
            lines.extend([
                "        subgraph \"Database Services\""
            ])
            
            service_connections = []
            # Add explicitly selected database services
            for service in (inputs.database_services or []):
                if service in AZURE_SERVICES_MAPPING:
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"DATABASE{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            # Add AI-suggested database services
            for service in ai_services:
                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "database":
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"DATABASE{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            lines.extend(service_connections)
            lines.append("        end")
        
        # Add security services if selected
        if inputs.security_services or (ai_services and any(service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "security" for service in ai_services)):
            lines.extend([
                "        subgraph \"Security Services\""
            ])
            
            service_connections = []
            # Add explicitly selected security services
            for service in (inputs.security_services or []):
                if service in AZURE_SERVICES_MAPPING:
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"SECURITY{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            # Add AI-suggested security services
            for service in ai_services:
                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "security":
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"SECURITY{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            lines.extend(service_connections)
            lines.append("        end")
        
        # Add monitoring services if selected
        if inputs.monitoring_services or (ai_services and any(service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "monitoring" for service in ai_services)):
            lines.extend([
                "        subgraph \"Monitoring Services\""
            ])
            
            service_connections = []
            # Add explicitly selected monitoring services
            for service in (inputs.monitoring_services or []):
                if service in AZURE_SERVICES_MAPPING:
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"MONITORING{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            # Add AI-suggested monitoring services
            for service in ai_services:
                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["category"] == "monitoring":
                    service_info = AZURE_SERVICES_MAPPING[service]
                    service_id = f"MONITORING{service_count}"
                    lines.append(f"            {service_id}[\"{service_info['icon']} {service_info['name']}\"]")
                    service_connections.append(f"            RESOURCEGROUP --> {service_id}")
                    service_count += 1
            
            lines.extend(service_connections)
            lines.append("        end")
    
    else:
        # If no services selected, show minimal placeholder
        lines.extend([
            "        PLACEHOLDER[\"📝 No Services Selected\"]",
            "        PLACEHOLDER2[\"ℹ️ Please select Azure services to generate architecture\"]"
        ])
    
    lines.append("    end")
    
    # Add styling
    lines.extend([
        "",
        "    classDef subscription fill:#e1f5fe,stroke:#01579b,stroke-width:2px;",
        "    classDef resourceGroup fill:#f3e5f5,stroke:#4a148c,stroke-width:2px;",
        "    classDef compute fill:#fff3e0,stroke:#e65100,stroke-width:2px;",
        "    classDef network fill:#e8f5e8,stroke:#1b5e20,stroke-width:2px;",
        "    classDef storage fill:#fce4ec,stroke:#880e4f,stroke-width:2px;",
        "    classDef database fill:#e3f2fd,stroke:#0d47a1,stroke-width:2px;",
        "    classDef security fill:#ffebee,stroke:#b71c1c,stroke-width:2px;",
        "    classDef monitoring fill:#f1f8e9,stroke:#33691e,stroke-width:2px;",
        "",
        "    class SUBSCRIPTION subscription;",
        "    class RESOURCEGROUP resourceGroup;"
    ])
    
    return "\n".join(lines)

def generate_enhanced_drawio_xml(inputs: CustomerInputs) -> str:
    """Generate enhanced Draw.io XML with comprehensive Azure stencils based on user selections"""
    
    def esc(s): 
        return html.escape(s) if s else ""
    
    template = generate_architecture_template(inputs)
    diagram_id = str(uuid.uuid4())
    
    # Base layout coordinates
    y_start = 100
    current_y = y_start
    section_height = 350
    service_width = 100
    service_height = 80
    
    # Build dynamic XML content
    xml_parts = [
        f"""<mxfile host="app.diagrams.net" modified="2024-01-01T00:00:00.000Z" agent="Azure Landing Zone Agent" version="1.0.0">
  <diagram name="Azure Landing Zone Architecture" id="azure-lz-{diagram_id}">
    <mxGraphModel dx="1422" dy="794" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="2400" pageHeight="1600" math="0" shadow="0">
      <root>
        <mxCell id="0" />
        <mxCell id="1" parent="0" />
        
        <!-- Azure Tenant Container -->
        <mxCell id="tenant" value="Azure Tenant - {esc(inputs.org_structure or 'Enterprise')}" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#dae8fc;strokeColor=#6c8ebf;fontSize=16;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="50" y="50" width="2300" height="1500" as="geometry" />
        </mxCell>
        
        <!-- Management Groups -->
        <mxCell id="mgmt-groups" value="Management Groups" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e1d5e7;strokeColor=#9673a6;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="500" height="250" as="geometry" />
        </mxCell>"""
    ]
    
    # Management Group structure based on template
    mg_x = 150
    mg_y = current_y + 50
    xml_parts.append(f"""
        <mxCell id="root-mg" value="Root MG" style="shape=mxgraph.azure.management;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{mg_x}" y="{mg_y}" width="80" height="60" as="geometry" />
        </mxCell>""")
    
    mg_x += 120
    for i, mg in enumerate(template['template']['management_groups'][1:3]):  # Platform and Workloads
        mg_id = mg.lower().replace(' ', '-')
        xml_parts.append(f"""
        <mxCell id="{mg_id}-mg" value="{mg}" style="shape=mxgraph.azure.management;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{mg_x}" y="{mg_y}" width="80" height="60" as="geometry" />
        </mxCell>""")
        mg_x += 120
    
    # Subscriptions
    current_y += 300
    xml_parts.append(f"""
        <!-- Subscriptions -->
        <mxCell id="subscriptions" value="Subscriptions" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#d5e8d4;strokeColor=#82b366;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="700" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
    
    sub_x = 750
    sub_y = current_y + 50
    for sub in template['template']['subscriptions'][:4]:  # First 4 subscriptions
        xml_parts.append(f"""
        <mxCell id="{sub.lower().replace(' ', '-')}-sub" value="{sub}" style="shape=mxgraph.azure.subscription;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sub_x}" y="{sub_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
        sub_x += 130
        if sub_x > 1200:  # Wrap to next row
            sub_x = 750
            sub_y += 100
    
    # Network Architecture Section - Only if network services are explicitly selected
    if inputs.network_services:
        current_y += 300
        xml_parts.append(f"""
        <!-- Network Architecture -->
        <mxCell id="network" value="Network Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff2cc;strokeColor=#d6b656;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="800" height="{section_height}" as="geometry" />
        </mxCell>""")
        
        # Add only explicitly selected network services
        net_x = 200
        net_y = current_y + 80
        for i, service in enumerate(inputs.network_services):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="net-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{net_x}" y="{net_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                net_x += 120
                if net_x > 600:  # Wrap to next row
                    net_x = 200
                    net_y += 100
    
    # Compute Services Section
    if inputs.compute_services or inputs.workload:
        current_y += section_height + 50
        xml_parts.append(f"""
        <!-- Compute Services -->
        <mxCell id="compute" value="Compute Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff0e6;strokeColor=#d79b00;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1000" y="{current_y}" width="600" height="{section_height}" as="geometry" />
        </mxCell>""")
        
        comp_x = 1050
        comp_y = current_y + 50
        
        # Add selected compute services
        services_to_add = inputs.compute_services or []
        if inputs.workload and inputs.workload not in services_to_add:
            services_to_add.append(inputs.workload)
            
        for i, service in enumerate(services_to_add[:6]):  # Max 6 compute services
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="compute-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{comp_x}" y="{comp_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                comp_x += 120
                if comp_x > 1450:  # Wrap to next row
                    comp_x = 1050
                    comp_y += 100
    
    # Storage Services Section
    if inputs.storage_services:
        current_y += section_height + 50
        xml_parts.append(f"""
        <!-- Storage Services -->
        <mxCell id="storage" value="Storage Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f0f0f0;strokeColor=#666666;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        
        stor_x = 150
        stor_y = current_y + 50
        for i, service in enumerate(inputs.storage_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'storage_accounts')
                xml_parts.append(f"""
        <mxCell id="storage-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{stor_x}" y="{stor_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                stor_x += 120
                if stor_x > 550:
                    stor_x = 150
                    stor_y += 100
    
    # Database Services Section
    if inputs.database_services:
        db_y = current_y if not inputs.storage_services else current_y
        if inputs.storage_services:
            xml_parts.append(f"""
        <!-- Database Services -->
        <mxCell id="database" value="Database Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e6f3ff;strokeColor=#0066cc;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="800" y="{db_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        else:
            current_y += section_height + 50
            xml_parts.append(f"""
        <!-- Database Services -->
        <mxCell id="database" value="Database Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e6f3ff;strokeColor=#0066cc;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
            db_y = current_y
        
        db_x = 850 if inputs.storage_services else 150
        db_y += 50
        for i, service in enumerate(inputs.database_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'sql_database')
                xml_parts.append(f"""
        <mxCell id="database-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{db_x}" y="{db_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                db_x += 120
                if db_x > (1250 if inputs.storage_services else 550):
                    db_x = 850 if inputs.storage_services else 150
                    db_y += 100
    
    # Security Services Section (always present)
    current_y += 300
    xml_parts.append(f"""
        <!-- Security & Identity Services -->
        <mxCell id="security" value="Security &amp; Identity Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#ffebee;strokeColor=#c62828;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1700" y="{y_start}" width="600" height="600" as="geometry" />
        </mxCell>""")
    
    # Core security services (always present)
    sec_x = 1750
    sec_y = y_start + 50
    core_security = [
        ('azure-ad', 'Azure AD', 'azure_active_directory'),
        ('key-vault', 'Key Vault', 'key_vault'),
        ('security-center', 'Security Center', 'security_center')
    ]
    
    for sec_id, sec_name, sec_shape in core_security:
        xml_parts.append(f"""
        <mxCell id="{sec_id}" value="{sec_name}" style="shape=mxgraph.azure.{sec_shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sec_x}" y="{sec_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
        sec_x += 120
        if sec_x > 2100:
            sec_x = 1750
            sec_y += 100
    
    # Add additional selected security services
    if inputs.security_services:
        for i, service in enumerate(inputs.security_services):
            if service in AZURE_SERVICES_MAPPING and service not in ['active_directory', 'key_vault', 'security_center']:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="security-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sec_x}" y="{sec_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                sec_x += 120
                if sec_x > 2100:
                    sec_x = 1750
                    sec_y += 100
    
    # Analytics Services Section
    if inputs.analytics_services:
        analytics_y = y_start + 650
        xml_parts.append(f"""
        <!-- Analytics & AI Services -->
        <mxCell id="analytics" value="Analytics &amp; AI Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f3e5f5;strokeColor=#9c27b0;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1700" y="{analytics_y}" width="600" height="300" as="geometry" />
        </mxCell>""")
        
        ana_x = 1750
        ana_y = analytics_y + 50
        for i, service in enumerate(inputs.analytics_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="analytics-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{ana_x}" y="{ana_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                ana_x += 120
                if ana_x > 2100:
                    ana_x = 1750
                    ana_y += 100
    
    # Integration Services Section
    if inputs.integration_services:
        int_y = current_y
        xml_parts.append(f"""
        <!-- Integration Services -->
        <mxCell id="integration" value="Integration Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff8e1;strokeColor=#ff8f00;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{int_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        
        int_x = 150
        int_y += 50
        for i, service in enumerate(inputs.integration_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="integration-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{int_x}" y="{int_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                int_x += 120
                if int_x > 550:
                    int_x = 150
                    int_y += 100
    
    # DevOps Services Section
    if inputs.devops_services:
        devops_y = current_y if not inputs.integration_services else current_y
        devops_x_offset = 800 if inputs.integration_services else 100
        xml_parts.append(f"""
        <!-- DevOps Services -->
        <mxCell id="devops" value="DevOps Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f5f5f5;strokeColor=#666666;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="{devops_x_offset}" y="{devops_y}" width="400" height="250" as="geometry" />
        </mxCell>""")
        
        dev_x = devops_x_offset + 50
        dev_y = devops_y + 50
        for i, service in enumerate(inputs.devops_services[:3]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="devops-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{dev_x}" y="{dev_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                dev_x += 120
                if dev_x > (devops_x_offset + 250):
                    dev_x = devops_x_offset + 50
                    dev_y += 100
    
    # Add basic connections
    xml_parts.append("""
        <!-- Key Connections -->
        <mxCell id="conn1" edge="1" source="root-mg" target="platform-mg" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn2" edge="1" source="root-mg" target="landing-zones-mg" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn3" edge="1" source="hub-vnet" target="spoke1-vnet" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn4" edge="1" source="hub-vnet" target="spoke2-vnet" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>""")
    
    return "".join(xml_parts)


def generate_professional_documentation(inputs: CustomerInputs) -> Dict[str, str]:
    """Generate comprehensive enterprise-grade TSD, HLD, and LLD documentation (10+ pages)"""
    
    template = generate_architecture_template(inputs)
    timestamp = datetime.now().strftime("%Y-%m-%d")
    
    # Generate AI insights if additional inputs are provided
    url_analysis = ""
    doc_analysis = ""
    ai_recommendations = ""
    
    # Get comprehensive AI analysis results
    architecture_pattern = template.get("architecture_pattern", "Custom Architecture")
    connectivity_requirements = template.get("connectivity_requirements", "Standard Azure networking")
    security_considerations = template.get("security_considerations", "Standard security practices")
    ai_reasoning = template.get("ai_reasoning", "Standard architectural decisions applied")
    scalability_design = template.get("scalability_design", "Standard scalability patterns")
    operational_excellence = template.get("operational_excellence", "Standard monitoring and operations")
    cost_optimization = template.get("cost_optimization", "Standard cost optimization")
    ai_services = template.get("ai_services", [])
    
    try:
        if inputs.url_input:
            url_analysis = analyze_url_content(inputs.url_input)
            
        if inputs.uploaded_files_info:
            doc_analysis = "Document analysis results incorporated from uploaded files."
            
        # Generate comprehensive AI-enhanced recommendations
        ai_recommendations = generate_ai_enhanced_recommendations(inputs, url_analysis, doc_analysis)
    except Exception as e:
        logger.warning(f"AI enhancement failed: {e}")
        ai_recommendations = "AI enhancement not available - using standard recommendations."
    
    # Count total services
    all_selected_services = []
    all_selected_services.extend(inputs.compute_services or [])
    all_selected_services.extend(inputs.network_services or [])
    all_selected_services.extend(inputs.storage_services or [])
    all_selected_services.extend(inputs.database_services or [])
    all_selected_services.extend(inputs.security_services or [])
    all_selected_services.extend(inputs.monitoring_services or [])
    all_selected_services.extend(inputs.ai_services or [])
    all_selected_services.extend(inputs.analytics_services or [])
    all_selected_services.extend(inputs.integration_services or [])
    all_selected_services.extend(inputs.devops_services or [])
    all_selected_services.extend(inputs.backup_services or [])
    all_selected_services.extend(ai_services)
    
    # Remove duplicates and get service details
    unique_services = list(set(all_selected_services))
    
    # Technical Specification Document (TSD) - Enterprise Edition (10+ Pages)
    tsd = f"""# Technical Specification Document (TSD)
## Azure Landing Zone Architecture - Enterprise Edition

**Document Version:** 4.0 (AI-Enhanced Enterprise Architecture)
**Date:** {timestamp}
**Classification:** Internal Use Only
**Document Owner:** Azure Solutions Architecture Team
**Review Cycle:** Quarterly

---

## Table of Contents

1. [Executive Summary](#executive-summary)
2. [Business Requirements Analysis](#business-requirements-analysis)
3. [AI-Powered Architecture Analysis](#ai-powered-architecture-analysis)
4. [Architecture Framework & Patterns](#architecture-framework--patterns)
5. [Service Architecture & Connectivity](#service-architecture--connectivity)
6. [Security Architecture](#security-architecture)
7. [Network Architecture & Connectivity](#network-architecture--connectivity)
8. [Data Architecture & Management](#data-architecture--management)
9. [Scalability & Performance Design](#scalability--performance-design)
10. [Operational Excellence Framework](#operational-excellence-framework)
11. [Cost Optimization Strategy](#cost-optimization-strategy)
12. [Compliance & Governance](#compliance--governance)
13. [Risk Assessment & Mitigation](#risk-assessment--mitigation)
14. [Implementation Roadmap](#implementation-roadmap)
15. [Appendices](#appendices)

---

## 1. Executive Summary

### 1.1 Document Purpose
This Technical Specification Document (TSD) provides comprehensive technical specifications for implementing an enterprise-ready Azure Landing Zone architecture. The design is based on advanced AI-powered analysis of business requirements, Azure Well-Architected Framework principles, and industry best practices.

### 1.2 Business Context
**Primary Business Objective:** {inputs.business_objective or 'Enterprise digital transformation with cloud-native architecture'}

The proposed architecture addresses critical business needs including:
- Scalable and resilient cloud infrastructure
- Enterprise-grade security and compliance
- Operational efficiency and automation
- Cost optimization and resource management
- Future-proof technology stack

### 1.3 Architecture Overview
**Selected Architecture Pattern:** {architecture_pattern}

The solution encompasses {len(unique_services)} carefully selected Azure services that form an integrated, enterprise-ready platform designed for:
- High availability (99.9% uptime SLA)
- Horizontal and vertical scalability
- Security-first design principles
- Operational excellence and monitoring
- Cost-effective resource utilization

### 1.4 Key Success Metrics
- **Availability:** 99.9% uptime target
- **Performance:** Sub-second response times for critical operations
- **Security:** Zero successful security breaches
- **Cost:** 15-20% reduction in operational costs compared to on-premises
- **Compliance:** 100% adherence to regulatory requirements

---

## 2. Business Requirements Analysis

### 2.1 Functional Requirements
**Primary Business Objective:** {inputs.business_objective or 'Cost optimization and operational efficiency with enterprise-grade reliability'}

**Industry Context:** {inputs.industry or 'Multi-industry applicable with focus on scalability and security'}

**Regulatory Framework:** {inputs.regulatory or 'Standard compliance including data protection, privacy regulations, and industry-specific requirements'}

### 2.2 Non-Functional Requirements

#### 2.2.1 Performance Requirements
- Response time: < 1 second for 95% of user requests
- Throughput: Support for 10,000+ concurrent users
- Data processing: Real-time analytics with < 100ms latency

#### 2.2.2 Availability Requirements
- System uptime: 99.9% (8.77 hours downtime per year)
- Recovery time objective (RTO): < 4 hours
- Recovery point objective (RPO): < 1 hour

#### 2.2.3 Scalability Requirements
- Horizontal scaling: Support for 10-times traffic spikes
- Geographic distribution: Multi-region deployment capability
- Resource elasticity: Auto-scaling based on demand

#### 2.2.4 Security Requirements
- Data encryption: At rest and in transit
- Identity management: Role-based access control (RBAC)
- Network security: Zero-trust architecture
- Compliance: GDPR, SOC 2, ISO 27001 readiness

### 2.3 Organizational Context
**Organization Structure:** {inputs.org_structure or 'Enterprise with distributed teams and centralized governance'}

**Governance Model:** {inputs.governance or 'Centralized governance with delegated operational permissions'}

**Operational Model:** {inputs.ops_model or 'DevOps with automated CI/CD pipelines and infrastructure as code'}

---

## 3. AI-Powered Architecture Analysis

### 3.1 Intelligent Architecture Design
The architecture design leverages advanced AI analysis to ensure optimal service selection and configuration based on specific business requirements.

**AI Architecture Reasoning:**
{ai_reasoning}

### 3.2 Service Selection Rationale
The AI analysis identified {len(unique_services)} critical Azure services based on:

1. **Functional Requirements Mapping:** Direct correlation between business needs and technical capabilities
2. **Integration Patterns:** Optimal service combinations for seamless connectivity
3. **Performance Optimization:** Services selected for maximum efficiency and scalability
4. **Cost-Benefit Analysis:** Balanced approach between functionality and cost
5. **Future-Proofing:** Architecture designed for evolution and expansion

### 3.3 Architecture Pattern Analysis
**Identified Pattern:** {architecture_pattern}

This pattern was selected because:
- Aligns with business scalability requirements
- Provides optimal separation of concerns
- Enables independent service scaling
- Supports microservices architecture principles
- Facilitates DevOps and continuous deployment

### 3.4 AI-Enhanced Recommendations
{ai_recommendations if ai_recommendations else "Standard enterprise recommendations applied with focus on security, scalability, and operational excellence."}

---

## 4. Architecture Framework & Patterns

### 4.1 Azure Well-Architected Framework Alignment

#### 4.1.1 Reliability
- Multi-region deployment strategy
- Automated failover mechanisms
- Health monitoring and alerting
- Disaster recovery planning

#### 4.1.2 Security
- Zero-trust network architecture
- Identity and access management
- Data encryption and key management
- Security monitoring and threat detection

#### 4.1.3 Cost Optimization
- Resource right-sizing
- Reserved capacity planning
- Automated cost monitoring
- Continuous optimization

#### 4.1.4 Operational Excellence
- Infrastructure as Code (IaC)
- Automated deployment pipelines
- Comprehensive monitoring and logging
- Performance optimization

#### 4.1.5 Performance Efficiency
- Auto-scaling configurations
- Caching strategies
- Content delivery optimization
- Database performance tuning

### 4.2 Connectivity Architecture
**Connectivity Requirements:** {connectivity_requirements}

The connectivity architecture implements:
- Hub-and-spoke network topology
- Private endpoint connectivity
- API gateway patterns
- Service mesh for microservices
- Load balancing and traffic management

---

## 5. Service Architecture & Connectivity

### 5.1 Selected Azure Services Portfolio

#### 5.1.1 Compute Services
**Selected:** {', '.join(inputs.compute_services) if inputs.compute_services else 'None explicitly selected (AI may recommend based on requirements)'}

**AI-Recommended Compute:** {', '.join([s for s in ai_services if s in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[s]['category'] == 'compute']) if ai_services else 'None recommended'}

#### 5.1.2 Network Services
**Selected:** {', '.join(inputs.network_services) if inputs.network_services else 'None explicitly selected (AI may recommend based on requirements)'}

**AI-Recommended Network:** {', '.join([s for s in ai_services if s in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[s]['category'] == 'network']) if ai_services else 'None recommended'}

#### 5.1.3 Storage Services
**Selected:** {', '.join(inputs.storage_services) if inputs.storage_services else 'None explicitly selected (AI may recommend based on requirements)'}

**AI-Recommended Storage:** {', '.join([s for s in ai_services if s in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[s]['category'] == 'storage']) if ai_services else 'None recommended'}

#### 5.1.4 Database Services
**Selected:** {', '.join(inputs.database_services) if inputs.database_services else 'None explicitly selected (AI may recommend based on requirements)'}

**AI-Recommended Database:** {', '.join([s for s in ai_services if s in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[s]['category'] == 'database']) if ai_services else 'None recommended'}

#### 5.1.5 Security Services
**Selected:** {', '.join(inputs.security_services) if inputs.security_services else 'None explicitly selected (AI may recommend based on requirements)'}

**AI-Recommended Security:** {', '.join([s for s in ai_services if s in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[s]['category'] == 'security']) if ai_services else 'None recommended'}

### 5.2 Service Integration Patterns

#### 5.2.1 API Management Strategy
- Centralized API gateway for external access
- Rate limiting and throttling policies
- API versioning and lifecycle management
- Security token validation and transformation

#### 5.2.2 Data Flow Architecture
- Event-driven messaging patterns
- Asynchronous processing pipelines
- Real-time streaming analytics
- Batch processing workflows

#### 5.2.3 Microservices Communication
- Service-to-service authentication
- Circuit breaker patterns
- Retry and timeout policies
- Distributed tracing and monitoring

---

## 6. Security Architecture

### 6.1 Security Framework
**Security Considerations:** {security_considerations}

### 6.2 Zero Trust Architecture
- Identity-based access controls
- Least privilege principle
- Continuous verification
- Network micro-segmentation

### 6.3 Data Protection Strategy
- Encryption at rest using customer-managed keys
- Encryption in transit with TLS 1.3
- Data classification and labeling
- Data loss prevention (DLP) policies

### 6.4 Identity and Access Management
- Azure Active Directory integration
- Multi-factor authentication (MFA)
- Privileged identity management (PIM)
- Conditional access policies

### 6.5 Network Security
- Network security groups (NSGs)
- Azure Firewall with threat intelligence
- Web application firewall (WAF)
- DDoS protection standard

---

## 7. Network Architecture & Connectivity

### 7.1 Network Topology
- Hub-and-spoke architecture
- Virtual network peering
- ExpressRoute connectivity (if applicable)
- VPN gateway for hybrid scenarios

### 7.2 Connectivity Matrix
```
Service Tier 1 (Web) -> Service Tier 2 (Application) -> Service Tier 3 (Data)
     ↓                          ↓                           ↓
Load Balancer           ->  App Services            ->   Database Services
Application Gateway     ->  Container Services      ->   Storage Services
CDN                     ->  Function Apps           ->   Cache Services
```

### 7.3 Private Connectivity
- Private endpoints for PaaS services
- Service endpoints where applicable
- Private DNS zones for name resolution
- Network policies for pod-to-pod communication

---

## 8. Data Architecture & Management

### 8.1 Data Strategy
- Data lifecycle management
- Backup and archival policies
- Data governance framework
- Analytics and reporting strategy

### 8.2 Database Design
- Horizontal and vertical partitioning
- Read replicas for performance
- Cross-region replication
- Automated backup strategies

---

## 9. Scalability & Performance Design

### 9.1 Scalability Strategy
**Scalability Design:** {scalability_design}

### 9.2 Performance Optimization
- Caching strategies (Redis, CDN)
- Database performance tuning
- Application performance monitoring
- Resource optimization

### 9.3 Auto-scaling Configuration
- Horizontal pod autoscaling
- Virtual machine scale sets
- Application Gateway autoscaling
- Database auto-scaling policies

---

## 10. Operational Excellence Framework

### 10.1 Operations Strategy
**Operational Excellence:** {operational_excellence}

### 10.2 Monitoring and Observability
- Comprehensive logging strategy
- Application performance monitoring
- Infrastructure monitoring
- Security monitoring and alerting

### 10.3 DevOps Integration
- Infrastructure as Code (IaC)
- CI/CD pipeline automation
- Automated testing strategies
- Configuration management

---

## 11. Cost Optimization Strategy

### 11.1 Cost Management
**Cost Optimization:** {cost_optimization}

### 11.2 Resource Optimization
- Right-sizing recommendations
- Reserved capacity planning
- Spot instance utilization
- Resource scheduling automation

### 11.3 Financial Operations (FinOps)
- Cost allocation and chargeback
- Budget monitoring and alerting
- Cost optimization automation
- Regular cost review processes

---

## 12. Compliance & Governance

### 12.1 Governance Framework
- Azure Policy implementation
- Role-based access control (RBAC)
- Resource tagging strategy
- Compliance monitoring

### 12.2 Regulatory Compliance
- GDPR compliance measures
- Industry-specific regulations
- Data residency requirements
- Audit and reporting capabilities

---

## 13. Risk Assessment & Mitigation

### 13.1 Risk Analysis
- Technical risks and mitigation strategies
- Security risks and controls
- Operational risks and procedures
- Business continuity planning

### 13.2 Disaster Recovery
- Recovery time objectives (RTO)
- Recovery point objectives (RPO)
- Backup and restore procedures
- Business continuity testing

---

## 14. Implementation Roadmap

### 14.1 Phase 1: Foundation (Weeks 1-4)
- Core infrastructure deployment
- Network and security setup
- Identity and access management
- Basic monitoring implementation

### 14.2 Phase 2: Application Services (Weeks 5-8)
- Application tier deployment
- Database setup and configuration
- API gateway implementation
- Performance optimization

### 14.3 Phase 3: Advanced Features (Weeks 9-12)
- Advanced monitoring and alerting
- Disaster recovery testing
- Performance tuning
- Security hardening

### 14.4 Phase 4: Production Readiness (Weeks 13-16)
- Production deployment
- Load testing and optimization
- Documentation completion
- Knowledge transfer

---

## 15. Appendices

### Appendix A: Service Configuration Details
[Detailed configuration parameters for each Azure service]

### Appendix B: Network Diagrams
[Comprehensive network topology diagrams]

### Appendix C: Security Policies
[Detailed security policy configurations]

### Appendix D: Monitoring Configurations
[Monitoring and alerting configuration details]

### Appendix E: Cost Estimates
[Detailed cost breakdown and projections]

---

**Document Control:**
"""
