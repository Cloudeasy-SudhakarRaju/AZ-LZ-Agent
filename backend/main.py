from fastapi import FastAPI, Response, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any
import html
import json
import uuid
import os
import base64
import subprocess
import tempfile
import logging
import traceback
from datetime import datetime
from pathlib import Path
import requests
import google.generativeai as genai

# Document processing imports
import PyPDF2
import openpyxl
from pptx import Presentation

# Import diagrams for Azure architecture generation
from diagrams import Diagram, Cluster, Edge
from diagrams.azure.compute import VM, AKS, AppServices, FunctionApps, ContainerInstances, ServiceFabricClusters, BatchAccounts
from diagrams.azure.network import VirtualNetworks, ApplicationGateway, LoadBalancers, Firewall, ExpressrouteCircuits, VirtualNetworkGateways, CDNProfiles, TrafficManagerProfiles
from diagrams.azure.storage import StorageAccounts, BlobStorage, DataLakeStorage
from diagrams.azure.database import SQLDatabases, CosmosDb, DatabaseForMysqlServers, DatabaseForPostgresqlServers, CacheForRedis
from diagrams.azure.security import KeyVaults, SecurityCenter, Sentinel
from diagrams.azure.identity import ActiveDirectory
from diagrams.azure.analytics import SynapseAnalytics, DataFactories, Databricks, StreamAnalyticsJobs, EventHubs, LogAnalyticsWorkspaces
from diagrams.azure.integration import LogicApps, ServiceBus, EventGridTopics, APIManagement
from diagrams.azure.devops import Devops, Pipelines, ApplicationInsights
from diagrams.azure.general import Subscriptions, Resourcegroups, Servicehealth, Templates, Support
from diagrams.azure.web import AppServices as WebApps
# Additional imports for missing services
from diagrams.azure.ml import CognitiveServices, MachineLearningServiceWorkspaces, BotServices
from diagrams.azure.compute import ACR

# Import intelligent diagram generator
from intelligent_diagram_generator import IntelligentArchitectureDiagramGenerator, DiagramGenerationResult

# Import LangGraph orchestrator
from langgraph_workflow import AzureLandingZoneOrchestrator, categorize_services_by_hub_spoke, create_orchestrator

# Import architecture validation
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from validate_architecture import (
    validate_architecture, validate_resource, generate_diagram_structure,
    ValidationResult, DiagramStructure, AZ_LZ_RULES
)

app = FastAPI(
    title="Azure Landing Zone Agent",
    description="Professional Azure Landing Zone Architecture Generator",
    version="1.0.0"
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # open for dev, restrict later
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure Google Gemini API
GEMINI_API_KEY = "AIzaSyCuYYvGh5wjwNniv9ZQ1QC-5pxwdj5lCWQ"
genai.configure(api_key=GEMINI_API_KEY)

# Initialize Gemini model
try:
    gemini_model = genai.GenerativeModel('gemini-1.5-pro')
    logger.info("Google Gemini API configured successfully")
except Exception as e:
    logger.error(f"Failed to configure Gemini API: {e}")
    gemini_model = None

# Initialize intelligent diagram generator
try:
    # Try to get OpenAI API key from environment
    openai_api_key = os.getenv('OPENAI_API_KEY')
    if not openai_api_key:
        logger.warning("OpenAI API key not found. Intelligent diagram generator will not be available.")
        intelligent_generator = None
    else:
        intelligent_generator = IntelligentArchitectureDiagramGenerator(openai_api_key)
        logger.info("Intelligent Architecture Diagram Generator initialized successfully")
except Exception as e:
    logger.error(f"Failed to initialize intelligent generator: {e}")
    intelligent_generator = None

# Initialize LangGraph orchestrator
try:
    langgraph_orchestrator = create_orchestrator()
    logger.info("LangGraph Hub and Spoke Orchestrator initialized")
except Exception as e:
    logger.error(f"Failed to initialize LangGraph orchestrator: {e}")
    langgraph_orchestrator = None


class CustomerInputs(BaseModel):
    # Business Requirements
    business_objective: Optional[str] = Field(None, description="Primary business objective")
    regulatory: Optional[str] = Field(None, description="Regulatory requirements")
    industry: Optional[str] = Field(None, description="Industry vertical")
    
    # Organization Structure
    org_structure: Optional[str] = Field(None, description="Organization structure")
    governance: Optional[str] = Field(None, description="Governance model")
    
    # Identity & Access Management
    identity: Optional[str] = Field(None, description="Identity management approach")
    
    # Networking & Connectivity
    connectivity: Optional[str] = Field(None, description="Connectivity requirements")
    network_model: Optional[str] = Field(None, description="Network topology model")
    ip_strategy: Optional[str] = Field(None, description="IP address strategy")
    
    # Security
    security_zone: Optional[str] = Field(None, description="Security zone requirements")
    security_posture: Optional[str] = Field(None, description="Security posture")
    key_vault: Optional[str] = Field(None, description="Key management approach")
    threat_protection: Optional[str] = Field(None, description="Threat protection strategy")
    
    # Legacy fields for backward compatibility
    workload: Optional[str] = Field(None, description="Primary workload type")
    architecture_style: Optional[str] = Field(None, description="Architecture style")
    scalability: Optional[str] = Field(None, description="Scalability requirements")
    
    # Operations
    ops_model: Optional[str] = Field(None, description="Operations model")
    monitoring: Optional[str] = Field(None, description="Monitoring strategy")
    backup: Optional[str] = Field(None, description="Backup and recovery strategy")
    
    # Infrastructure
    topology_pattern: Optional[str] = Field(None, description="Topology pattern")
    
    # Migration & Cost
    migration_scope: Optional[str] = Field(None, description="Migration scope")
    cost_priority: Optional[str] = Field(None, description="Cost optimization priority")
    iac: Optional[str] = Field(None, description="Infrastructure as Code preference")
    
    # Specific Azure Service Selections
    # Compute Services
    compute_services: Optional[List[str]] = Field(default_factory=list, description="Selected compute services")
    
    # Networking Services
    network_services: Optional[List[str]] = Field(default_factory=list, description="Selected networking services")
    
    # Storage Services
    storage_services: Optional[List[str]] = Field(default_factory=list, description="Selected storage services")
    
    # Database Services
    database_services: Optional[List[str]] = Field(default_factory=list, description="Selected database services")
    
    # Security Services
    security_services: Optional[List[str]] = Field(default_factory=list, description="Selected security services")
    
    # Monitoring Services
    monitoring_services: Optional[List[str]] = Field(default_factory=list, description="Selected monitoring services")
    
    # AI/ML Services
    ai_services: Optional[List[str]] = Field(default_factory=list, description="Selected AI/ML services")
    
    # Analytics Services
    analytics_services: Optional[List[str]] = Field(default_factory=list, description="Selected analytics services")
    
    # Integration Services
    integration_services: Optional[List[str]] = Field(default_factory=list, description="Selected integration services")
    
    # DevOps Services
    devops_services: Optional[List[str]] = Field(default_factory=list, description="Selected DevOps services")
    
    # Backup Services
    backup_services: Optional[List[str]] = Field(default_factory=list, description="Selected backup services")
    
    # Enhanced Input Fields for AI Integration
    free_text_input: Optional[str] = Field(None, description="Free-form text input for additional requirements and context")
    url_input: Optional[str] = Field(None, description="URL for web content analysis")
    uploaded_files_info: Optional[List[Dict[str, Any]]] = Field(default_factory=list, description="Information about uploaded files")
    
    # Enterprise Resource Auto-Inclusion Preferences
    enterprise_resources_mode: Optional[str] = Field(
        "auto_when_missing", 
        description="Enterprise resources inclusion mode: 'always_include', 'auto_when_missing', or 'never_auto_include'"
    )
    show_enterprise_connections: Optional[bool] = Field(
        True, 
        description="Whether to explicitly show connections between enterprise resources and VM/VNet"
    )


# Azure Architecture Templates and Patterns
AZURE_TEMPLATES = {
    "enterprise": {
        "name": "Enterprise Scale Landing Zone",
        "management_groups": ["Root", "Platform", "Landing Zones", "Sandbox", "Decommissioned"],
        "subscriptions": ["Connectivity", "Identity", "Management", "Production", "Development"],
        "core_services": ["Azure AD", "Azure Policy", "Azure Monitor", "Key Vault", "Security Center"]
    },
    "small_medium": {
        "name": "Small-Medium Enterprise Landing Zone", 
        "management_groups": ["Root", "Platform", "Workloads"],
        "subscriptions": ["Platform", "Production", "Development"],
        "core_services": ["Azure AD", "Azure Monitor", "Key Vault"]
    },
    "startup": {
        "name": "Startup Landing Zone",
        "management_groups": ["Root", "Workloads"],
        "subscriptions": ["Production", "Development"],
        "core_services": ["Azure AD", "Azure Monitor"]
    }
}

AZURE_SERVICES_MAPPING = {
    # Compute Services
    "virtual_machines": {"name": "Azure Virtual Machines", "icon": "🖥️", "drawio_shape": "virtual_machine", "diagram_class": VM, "category": "compute"},
    "aks": {"name": "Azure Kubernetes Service", "icon": "☸️", "drawio_shape": "kubernetes_service", "diagram_class": AKS, "category": "compute"},
    "app_services": {"name": "Azure App Services", "icon": "🌐", "drawio_shape": "app_services", "diagram_class": AppServices, "category": "compute"},
    "web_apps": {"name": "Azure Web Apps", "icon": "🌐", "drawio_shape": "app_services", "diagram_class": WebApps, "category": "compute"},
    "functions": {"name": "Azure Functions", "icon": "⚡", "drawio_shape": "function_app", "diagram_class": FunctionApps, "category": "compute"},
    "container_instances": {"name": "Container Instances", "icon": "📦", "drawio_shape": "container_instances", "diagram_class": ContainerInstances, "category": "compute"},
    "service_fabric": {"name": "Service Fabric", "icon": "🏗️", "drawio_shape": "service_fabric", "diagram_class": ServiceFabricClusters, "category": "compute"},
    "batch": {"name": "Azure Batch", "icon": "⚙️", "drawio_shape": "batch_accounts", "diagram_class": BatchAccounts, "category": "compute"},
    
    # Networking Services
    "virtual_network": {"name": "Virtual Network", "icon": "🌐", "drawio_shape": "virtual_network", "diagram_class": VirtualNetworks, "category": "network"},
    "vpn_gateway": {"name": "VPN Gateway", "icon": "🔒", "drawio_shape": "vpn_gateway", "diagram_class": VirtualNetworkGateways, "category": "network"},
    "expressroute": {"name": "ExpressRoute", "icon": "⚡", "drawio_shape": "expressroute_circuits", "diagram_class": ExpressrouteCircuits, "category": "network"},
    "load_balancer": {"name": "Load Balancer", "icon": "⚖️", "drawio_shape": "load_balancer", "diagram_class": LoadBalancers, "category": "network"},
    "application_gateway": {"name": "Application Gateway", "icon": "🚪", "drawio_shape": "application_gateway", "diagram_class": ApplicationGateway, "category": "network"},
    "firewall": {"name": "Azure Firewall", "icon": "🛡️", "drawio_shape": "firewall", "diagram_class": Firewall, "category": "network"},
    "waf": {"name": "Web Application Firewall", "icon": "🛡️", "drawio_shape": "application_gateway", "diagram_class": ApplicationGateway, "category": "network"},
    "cdn": {"name": "Content Delivery Network", "icon": "🌍", "drawio_shape": "cdn_profiles", "diagram_class": CDNProfiles, "category": "network"},
    "traffic_manager": {"name": "Traffic Manager", "icon": "🚦", "drawio_shape": "traffic_manager_profiles", "diagram_class": TrafficManagerProfiles, "category": "network"},
    "virtual_wan": {"name": "Virtual WAN", "icon": "🌐", "drawio_shape": "virtual_wan", "diagram_class": VirtualNetworks, "category": "network"},
    
    # Storage Services
    "storage_accounts": {"name": "Storage Accounts", "icon": "💾", "drawio_shape": "storage_accounts", "diagram_class": StorageAccounts, "category": "storage"},
    "blob_storage": {"name": "Blob Storage", "icon": "📄", "drawio_shape": "blob_storage", "diagram_class": BlobStorage, "category": "storage"},
    "file_storage": {"name": "Azure Files", "icon": "📁", "drawio_shape": "files", "diagram_class": StorageAccounts, "category": "storage"},
    "disk_storage": {"name": "Managed Disks", "icon": "💿", "drawio_shape": "managed_disks", "diagram_class": StorageAccounts, "category": "storage"},
    "data_lake": {"name": "Data Lake Storage", "icon": "🏞️", "drawio_shape": "data_lake_storage", "diagram_class": DataLakeStorage, "category": "storage"},
    
    # Database Services
    "sql_database": {"name": "Azure SQL Database", "icon": "🗄️", "drawio_shape": "sql_database", "diagram_class": SQLDatabases, "category": "database"},
    "sql_managed_instance": {"name": "SQL Managed Instance", "icon": "🗄️", "drawio_shape": "sql_managed_instance", "diagram_class": SQLDatabases, "category": "database"},
    "cosmos_db": {"name": "Cosmos DB", "icon": "🌍", "drawio_shape": "cosmos_db", "diagram_class": CosmosDb, "category": "database"},
    "mysql": {"name": "Azure Database for MySQL", "icon": "🐬", "drawio_shape": "database_for_mysql_servers", "diagram_class": DatabaseForMysqlServers, "category": "database"},
    "postgresql": {"name": "Azure Database for PostgreSQL", "icon": "🐘", "drawio_shape": "database_for_postgresql_servers", "diagram_class": DatabaseForPostgresqlServers, "category": "database"},
    "mariadb": {"name": "Azure Database for MariaDB", "icon": "🗄️", "drawio_shape": "database_for_mariadb_servers", "diagram_class": DatabaseForMysqlServers, "category": "database"},
    "redis_cache": {"name": "Azure Cache for Redis", "icon": "⚡", "drawio_shape": "cache_redis", "diagram_class": CacheForRedis, "category": "database"},
    
    # Security Services
    "key_vault": {"name": "Azure Key Vault", "icon": "🔐", "drawio_shape": "key_vault", "diagram_class": KeyVaults, "category": "security"},
    "active_directory": {"name": "Azure Active Directory", "icon": "👤", "drawio_shape": "azure_active_directory", "diagram_class": ActiveDirectory, "category": "security"},
    "security_center": {"name": "Azure Security Center", "icon": "🛡️", "drawio_shape": "security_center", "diagram_class": SecurityCenter, "category": "security"},
    "sentinel": {"name": "Azure Sentinel", "icon": "👁️", "drawio_shape": "sentinel", "diagram_class": Sentinel, "category": "security"},
    "defender": {"name": "Microsoft Defender", "icon": "🛡️", "drawio_shape": "defender_easm", "diagram_class": SecurityCenter, "category": "security"},
    "information_protection": {"name": "Azure Information Protection", "icon": "🔒", "drawio_shape": "information_protection", "diagram_class": SecurityCenter, "category": "security"},
    
    # Monitoring & Management
    "monitor": {"name": "Azure Monitor", "icon": "📊", "drawio_shape": "monitor", "diagram_class": Servicehealth, "category": "monitoring"},
    "log_analytics": {"name": "Log Analytics", "icon": "📋", "drawio_shape": "log_analytics_workspaces", "diagram_class": LogAnalyticsWorkspaces, "category": "monitoring"},
    "application_insights": {"name": "Application Insights", "icon": "📈", "drawio_shape": "application_insights", "diagram_class": ApplicationInsights, "category": "monitoring"},
    "service_health": {"name": "Service Health", "icon": "❤️", "drawio_shape": "service_health", "diagram_class": Servicehealth, "category": "monitoring"},
    "advisor": {"name": "Azure Advisor", "icon": "💡", "drawio_shape": "advisor", "diagram_class": Support, "category": "monitoring"},
    
    # AI/ML Services  
    "cognitive_services": {"name": "Cognitive Services", "icon": "🧠", "drawio_shape": "cognitive_services", "diagram_class": CognitiveServices, "category": "ai"},
    "machine_learning": {"name": "Azure Machine Learning", "icon": "🤖", "drawio_shape": "machine_learning", "diagram_class": MachineLearningServiceWorkspaces, "category": "ai"},
    "bot_service": {"name": "Bot Service", "icon": "🤖", "drawio_shape": "bot_services", "diagram_class": BotServices, "category": "ai"},
    "form_recognizer": {"name": "Form Recognizer", "icon": "📄", "drawio_shape": "form_recognizer", "diagram_class": CognitiveServices, "category": "ai"},
    
    # Data & Analytics
    "synapse": {"name": "Azure Synapse Analytics", "icon": "📊", "drawio_shape": "synapse_analytics", "diagram_class": SynapseAnalytics, "category": "analytics"},
    "data_factory": {"name": "Azure Data Factory", "icon": "🏭", "drawio_shape": "data_factory", "diagram_class": DataFactories, "category": "analytics"},
    "databricks": {"name": "Azure Databricks", "icon": "📊", "drawio_shape": "databricks", "diagram_class": Databricks, "category": "analytics"},
    "stream_analytics": {"name": "Stream Analytics", "icon": "🌊", "drawio_shape": "stream_analytics", "diagram_class": StreamAnalyticsJobs, "category": "analytics"},
    "power_bi": {"name": "Power BI", "icon": "📊", "drawio_shape": "power_bi", "diagram_class": SynapseAnalytics, "category": "analytics"},
    
    # Integration Services
    "logic_apps": {"name": "Logic Apps", "icon": "🔗", "drawio_shape": "logic_apps", "diagram_class": LogicApps, "category": "integration"},
    "service_bus": {"name": "Service Bus", "icon": "🚌", "drawio_shape": "service_bus", "diagram_class": ServiceBus, "category": "integration"},
    "event_grid": {"name": "Event Grid", "icon": "⚡", "drawio_shape": "event_grid_topics", "diagram_class": EventGridTopics, "category": "integration"},
    "event_hubs": {"name": "Event Hubs", "icon": "📡", "drawio_shape": "event_hubs", "diagram_class": EventHubs, "category": "integration"},
    "api_management": {"name": "API Management", "icon": "🔌", "drawio_shape": "api_management", "diagram_class": APIManagement, "category": "integration"},
    
    # DevOps & Management
    "devops": {"name": "Azure DevOps", "icon": "⚙️", "drawio_shape": "devops", "diagram_class": Devops, "category": "devops"},
    "automation": {"name": "Azure Automation", "icon": "🤖", "drawio_shape": "automation_accounts", "diagram_class": Devops, "category": "devops"},
    "policy": {"name": "Azure Policy", "icon": "📋", "drawio_shape": "policy", "diagram_class": Templates, "category": "governance"},
    "blueprints": {"name": "Azure Blueprints", "icon": "📐", "drawio_shape": "blueprints", "diagram_class": Templates, "category": "governance"},
    "resource_manager": {"name": "Azure Resource Manager", "icon": "🏗️", "drawio_shape": "resource_groups", "diagram_class": Resourcegroups, "category": "governance"},
    
    # Backup & Recovery
    "backup": {"name": "Azure Backup", "icon": "💾", "drawio_shape": "backup", "diagram_class": StorageAccounts, "category": "backup"},
    "site_recovery": {"name": "Azure Site Recovery", "icon": "🔄", "drawio_shape": "site_recovery", "diagram_class": StorageAccounts, "category": "backup"},
    
    # Additional Platform Services
    "container_registry": {"name": "Azure Container Registry", "icon": "📦", "drawio_shape": "container_registries", "diagram_class": ACR, "category": "compute"},
    "spring_cloud": {"name": "Azure Spring Cloud", "icon": "🍃", "drawio_shape": "spring_cloud", "diagram_class": AppServices, "category": "compute"},
    "notification_hubs": {"name": "Notification Hubs", "icon": "📱", "drawio_shape": "notification_hubs", "diagram_class": EventGridTopics, "category": "integration"},
    "iot_hub": {"name": "IoT Hub", "icon": "🌐", "drawio_shape": "iot_hub", "diagram_class": EventHubs, "category": "integration"},
    "digital_twins": {"name": "Azure Digital Twins", "icon": "👥", "drawio_shape": "digital_twins", "diagram_class": CognitiveServices, "category": "ai"},
    "purview": {"name": "Microsoft Purview", "icon": "🔍", "drawio_shape": "purview", "diagram_class": SecurityCenter, "category": "governance"},
}

def get_safe_output_directory() -> str:
    """Get a safe directory for output files with fallback options"""
    directories_to_try = [
        "/tmp",
        tempfile.gettempdir(),
        os.path.expanduser("~/tmp"),
        "./tmp"
    ]
    
    for directory in directories_to_try:
        try:
            # Create directory if it doesn't exist
            Path(directory).mkdir(parents=True, exist_ok=True)
            
            # Test if we can write to it
            test_file = os.path.join(directory, f"test_write_{uuid.uuid4().hex[:8]}.tmp")
            with open(test_file, 'w') as f:
                f.write("test")
            os.remove(test_file)
            
            logger.info(f"Using output directory: {directory}")
            return directory
            
        except Exception as e:
            logger.warning(f"Cannot use directory {directory}: {e}")
            continue
    
    raise Exception("No writable output directory found. Tried: " + ", ".join(directories_to_try))

def cleanup_old_files(directory: str, max_age_hours: int = 24):
    """Clean up old generated files to prevent disk space issues"""
    try:
        current_time = datetime.now().timestamp()
        max_age_seconds = max_age_hours * 3600
        
        for filename in os.listdir(directory):
            if filename.startswith("azure_landing_zone_") and filename.endswith(".png"):
                filepath = os.path.join(directory, filename)
                try:
                    file_age = current_time - os.path.getmtime(filepath)
                    if file_age > max_age_seconds:
                        os.remove(filepath)
                        logger.info(f"Cleaned up old file: {filename}")
                except Exception as e:
                    logger.warning(f"Failed to clean up file {filename}: {e}")
                    
    except Exception as e:
        logger.warning(f"Failed to perform cleanup in {directory}: {e}")

# Google Gemini AI Integration Functions
def analyze_url_content(url: str) -> str:
    """Fetch and analyze URL content using Gemini AI"""
    try:
        if not gemini_model:
            return "Gemini AI not available for URL analysis"
            
        # Fetch URL content
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        content = response.text[:10000]  # Limit content size
        
        prompt = f"""
        Analyze the following web content for Azure architecture planning:
        
        URL: {url}
        Content: {content}
        
        Please provide insights for:
        1. Relevant Azure services mentioned or implied
        2. Architecture patterns or requirements
        3. Security and compliance considerations
        4. Scalability and performance requirements
        5. Integration requirements
        
        Format your response as a structured analysis.
        """
        
        result = gemini_model.generate_content(prompt)
        return result.text
        
    except Exception as e:
        logger.error(f"Error analyzing URL {url}: {e}")
        return f"Error analyzing URL: {str(e)}"

def process_uploaded_document(file_content: bytes, filename: str, file_type: str) -> str:
    """Process uploaded document using Gemini AI"""
    try:
        if not gemini_model:
            return "Gemini AI not available for document analysis"
            
        text_content = ""
        
        # Extract text based on file type
        if file_type.lower() == 'pdf':
            text_content = extract_pdf_text(file_content)
        elif file_type.lower() in ['xlsx', 'xls']:
            text_content = extract_excel_text(file_content)
        elif file_type.lower() in ['pptx', 'ppt']:
            text_content = extract_pptx_text(file_content)
        else:
            return f"Unsupported file type: {file_type}"
        
        if not text_content.strip():
            return "No readable text found in the document"
            
        prompt = f"""
        Analyze the following document content for Azure Landing Zone architecture planning:
        
        Document: {filename}
        Type: {file_type}
        Content: {text_content[:8000]}  # Limit content size
        
        Please provide insights for:
        1. Current architecture mentioned in the document
        2. Business requirements and objectives
        3. Compliance and regulatory requirements
        4. Security requirements
        5. Recommended Azure services and patterns
        6. Migration considerations
        7. Governance and operational requirements
        
        Format your response as a structured analysis for enterprise architecture planning.
        """
        
        result = gemini_model.generate_content(prompt)
        return result.text
        
    except Exception as e:
        logger.error(f"Error processing document {filename}: {e}")
        return f"Error processing document: {str(e)}"

def extract_pdf_text(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        import io
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""

def extract_excel_text(file_content: bytes) -> str:
    """Extract text from Excel file"""
    try:
        import io
        excel_file = io.BytesIO(file_content)
        workbook = openpyxl.load_workbook(excel_file)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting Excel text: {e}")
        return ""

def extract_pptx_text(file_content: bytes) -> str:
    """Extract text from PowerPoint file"""
    try:
        import io
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        text = ""
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PowerPoint text: {e}")
        return ""

def generate_ai_enhanced_recommendations(inputs: CustomerInputs, url_analysis: str = "", doc_analysis: str = "") -> str:
    """Generate AI-enhanced architecture recommendations using Gemini"""
    try:
        if not gemini_model:
            return "Standard recommendations (AI enhancement not available)"
            
        # Build context from inputs
        context = f"""
        Business Objective: {inputs.business_objective or 'Not specified'}
        Industry: {inputs.industry or 'General'}
        Organization Structure: {inputs.org_structure or 'Not specified'}
        Regulatory Requirements: {inputs.regulatory or 'Standard'}
        Security Requirements: {inputs.security_posture or 'Standard'}
        Scalability Requirements: {inputs.scalability or 'Standard'}
        Free-text Input: {inputs.free_text_input or 'None provided'}
        
        Selected Services:
        - Compute: {', '.join(inputs.compute_services or [])}
        - Network: {', '.join(inputs.network_services or [])}
        - Storage: {', '.join(inputs.storage_services or [])}
        - Database: {', '.join(inputs.database_services or [])}
        - Security: {', '.join(inputs.security_services or [])}
        """
        
        if url_analysis:
            context += f"\n\nURL Analysis Results:\n{url_analysis}"
            
        if doc_analysis:
            context += f"\n\nDocument Analysis Results:\n{doc_analysis}"
        
        prompt = f"""
        Based on the following customer requirements and analysis, provide comprehensive Azure Landing Zone architecture recommendations:
        
        {context}
        
        Please provide:
        1. Recommended Azure Landing Zone template (Enterprise Scale, Small Scale, etc.)
        2. Detailed architecture recommendations with specific Azure services
        3. Security and compliance strategy
        4. Network topology and connectivity recommendations
        5. Identity and access management strategy
        6. Operations and monitoring approach
        7. Cost optimization strategies
        8. Migration roadmap and phases
        9. Governance and policy recommendations
        10. Risk assessment and mitigation strategies
        
        Format your response as a comprehensive enterprise architecture document.
        """
        
        result = gemini_model.generate_content(prompt)
        return result.text
        
    except Exception as e:
        logger.error(f"Error generating AI recommendations: {e}")
        return f"Error generating AI recommendations: {str(e)}"

def convert_customer_inputs_to_architecture(inputs: CustomerInputs) -> Dict[str, Any]:
    """
    Convert CustomerInputs to the architecture format expected by the validation system.
    
    Args:
        inputs: CustomerInputs object from the API
        
    Returns:
        Dictionary representing the architecture with resources for validation
    """
    architecture = {
        "metadata": {
            "name": f"Azure Landing Zone - {inputs.business_objective or 'Enterprise Architecture'}",
            "version": "1.0",
            "compliance_requirements": []
        },
        "resources": []
    }
    
    # Add compliance requirements based on regulatory input
    if inputs.regulatory:
        regulatory_mapping = {
            "pci": "PCI-DSS",
            "hipaa": "HIPAA", 
            "sox": "SOX",
            "gdpr": "GDPR",
            "iso27001": "ISO27001"
        }
        reg_lower = inputs.regulatory.lower()
        for key, value in regulatory_mapping.items():
            if key in reg_lower:
                architecture["metadata"]["compliance_requirements"].append(value)
    
    # Convert service selections to resources
    resource_id_counter = 1
    
    # Process compute services
    if inputs.compute_services:
        for service in inputs.compute_services:
            resource = _create_resource_from_service(service, "compute", resource_id_counter, inputs)
            if resource:
                architecture["resources"].append(resource)
                resource_id_counter += 1
    
    # Process network services
    if inputs.network_services:
        for service in inputs.network_services:
            resource = _create_resource_from_service(service, "network", resource_id_counter, inputs)
            if resource:
                architecture["resources"].append(resource)
                resource_id_counter += 1
    
    # Process storage services
    if inputs.storage_services:
        for service in inputs.storage_services:
            resource = _create_resource_from_service(service, "storage", resource_id_counter, inputs)
            if resource:
                architecture["resources"].append(resource)
                resource_id_counter += 1
    
    # Process database services
    if inputs.database_services:
        for service in inputs.database_services:
            resource = _create_resource_from_service(service, "database", resource_id_counter, inputs)
            if resource:
                architecture["resources"].append(resource)
                resource_id_counter += 1
    
    # Process security services
    if inputs.security_services:
        for service in inputs.security_services:
            resource = _create_resource_from_service(service, "security", resource_id_counter, inputs)
            if resource:
                architecture["resources"].append(resource)
                resource_id_counter += 1
    
    return architecture

def _create_resource_from_service(service: str, category: str, resource_id: int, inputs: CustomerInputs) -> Optional[Dict[str, Any]]:
    """
    Create a resource configuration from a service string and customer inputs.
    """
    service_lower = service.lower().replace("-", "_").replace(" ", "_")
    
    # Map service names to resource types and configurations
    service_mappings = {
        # Compute services
        "virtual_machines": {
            "type": "VM",
            "base_config": {
                "subnet": "private-subnet" if inputs.security_posture == "zero-trust" else "public-subnet",
                "vnet": f"spoke-{category}-vnet",
                "availability_zones": inputs.scalability in ["high", "critical"],
                "network_security_group": True,
                "backup_enabled": inputs.backup in ["comprehensive", "standard"],
                "disk_encryption": True
            }
        },
        "aks": {
            "type": "AKS", 
            "base_config": {
                "private_cluster": inputs.security_posture == "zero-trust",
                "rbac_enabled": True,
                "network_policy": inputs.security_posture == "zero-trust",
                "container_insights": inputs.monitoring in ["azure-monitor", "comprehensive"]
            }
        },
        "app_services": {
            "type": "AppService",
            "base_config": {
                "vnet_integration": inputs.security_posture == "zero-trust",
                "private_endpoint": inputs.security_posture == "zero-trust",
                "https_only": True,
                "managed_identity": True,
                "application_insights": inputs.monitoring in ["azure-monitor", "application-insights"]
            }
        },
        # Database services
        "sql_database": {
            "type": "SQL",
            "base_config": {
                "public_access": inputs.security_posture != "zero-trust",
                "private_endpoint": inputs.security_posture == "zero-trust", 
                "encryption_at_rest": True,
                "encryption_in_transit": True,
                "auditing": inputs.regulatory is not None,
                "backup_retention_days": 30 if inputs.backup in ["comprehensive", "standard"] else 7
            }
        },
        "cosmos_db": {
            "type": "CosmosDB",
            "base_config": {
                "public_network_access": inputs.security_posture != "zero-trust",
                "private_endpoint": inputs.security_posture == "zero-trust",
                "encryption_at_rest": True,
                "firewall_enabled": True
            }
        },
        # Storage services
        "storage_accounts": {
            "type": "Storage", 
            "base_config": {
                "public_blob_access": inputs.security_posture != "zero-trust",
                "private_endpoint": inputs.security_posture == "zero-trust",
                "https_only": True,
                "min_tls_version": "1.2",
                "storage_analytics": inputs.monitoring in ["azure-monitor", "log-analytics"]
            }
        },
        # Network services
        "azure_firewall": {
            "type": "Firewall",
            "base_config": {
                "hub_vnet": True,
                "threat_intelligence": True,
                "diagnostic_logs": inputs.monitoring in ["azure-monitor", "log-analytics"]
            }
        },
        "firewall": {
            "type": "Firewall",
            "base_config": {
                "hub_vnet": True,
                "threat_intelligence": True,
                "diagnostic_logs": inputs.monitoring in ["azure-monitor", "log-analytics"]
            }
        },
        "key_vault": {
            "type": "KeyVault",
            "base_config": {
                "public_network_access": inputs.security_posture != "zero-trust",
                "private_endpoint": inputs.security_posture == "zero-trust"
            }
        }
    }
    
    if service_lower not in service_mappings:
        logger.warning(f"Unknown service type: {service}")
        return None
    
    mapping = service_mappings[service_lower]
    resource_name = f"{service.replace('_', '-')}-{resource_id:02d}"
    
    # Create base resource structure
    resource = {
        "name": resource_name,
        "type": mapping["type"],
        **mapping["base_config"]
    }
    
    # Add common tags based on inputs
    resource["tags"] = {
        "environment": "production",  # Default assumption
        "project": inputs.business_objective or "azure-landing-zone"
    }
    
    # Add owner tag if we can derive it from org structure
    if inputs.org_structure:
        resource["tags"]["owner"] = inputs.org_structure.lower().replace(" ", "")
    
    # Apply security posture configurations
    if inputs.security_posture == "zero-trust":
        _apply_zero_trust_config(resource)
    elif inputs.security_posture == "defense-in-depth":
        _apply_defense_in_depth_config(resource)
    
    return resource

def _apply_zero_trust_config(resource: Dict[str, Any]) -> None:
    """Apply zero trust security configurations to a resource"""
    resource_type = resource.get("type", "")
    
    # Universal zero trust settings
    if "public_access" in resource:
        resource["public_access"] = False
    if "public_blob_access" in resource:
        resource["public_blob_access"] = False
    if "public_network_access" in resource:
        resource["public_network_access"] = False
    
    # Enable private endpoints where applicable
    resource["private_endpoint"] = True
    
    # Type-specific zero trust configurations
    if resource_type == "VM":
        resource["subnet"] = "private-subnet"
    elif resource_type in ["SQL", "CosmosDB", "Storage"]:
        resource["firewall_enabled"] = True

def _apply_defense_in_depth_config(resource: Dict[str, Any]) -> None:
    """Apply defense in depth security configurations to a resource"""
    # More permissive than zero trust but still secure
    resource_type = resource.get("type", "")
    
    if resource_type == "VM":
        resource["network_security_group"] = True
        resource["vulnerability_assessment"] = True
    elif resource_type in ["SQL", "Storage"]:
        resource["firewall_rules"] = "restrictive"

def validate_customer_inputs(inputs: CustomerInputs) -> None:
    """Validate customer inputs to prevent potential errors"""
    # Check for extremely long strings that might cause issues
    string_fields = [
        inputs.business_objective, inputs.regulatory, inputs.industry,
        inputs.org_structure, inputs.governance, inputs.identity,
        inputs.connectivity, inputs.network_model, inputs.ip_strategy,
        inputs.security_zone, inputs.security_posture, inputs.key_vault,
        inputs.threat_protection, inputs.workload, inputs.architecture_style,
        inputs.scalability, inputs.ops_model, inputs.monitoring, inputs.backup,
        inputs.topology_pattern, inputs.migration_scope, inputs.cost_priority, inputs.iac,
        inputs.url_input
    ]
    
    for field in string_fields:
        if field and len(field) > 1000:  # Reasonable limit for most fields
            raise ValueError(f"Input field too long: {len(field)} characters (max 1000)")
    
    # Special validation for free-text input (allowing more characters)
    if inputs.free_text_input and len(inputs.free_text_input) > 10000:
        raise ValueError(f"Free text input too long: {len(inputs.free_text_input)} characters (max 10000)")
    
    # Check service lists for reasonable sizes
    service_lists = [
        inputs.compute_services, inputs.network_services, inputs.storage_services,
        inputs.database_services, inputs.security_services, inputs.monitoring_services,
        inputs.ai_services, inputs.analytics_services, inputs.integration_services,
        inputs.devops_services, inputs.backup_services
    ]
    
    for service_list in service_lists:
        if service_list and len(service_list) > 50:  # Reasonable limit
            raise ValueError(f"Too many services selected: {len(service_list)} (max 50)")
    
    # Validate URL format if provided
    if inputs.url_input:
        if not inputs.url_input.startswith(('http://', 'https://')):
            raise ValueError("URL must start with http:// or https://")
    
    # Validate uploaded files info
    if inputs.uploaded_files_info:
        if len(inputs.uploaded_files_info) > 10:  # Reasonable limit
            raise ValueError(f"Too many uploaded files: {len(inputs.uploaded_files_info)} (max 10)")

def _get_enterprise_resources() -> List[str]:
    """Get the list of enterprise resources that should be auto-included"""
    return ["key_vault", "active_directory", "firewall", "monitor"]

def _should_include_enterprise_resources(inputs: CustomerInputs, enterprise_resources: List[str]) -> bool:
    """Determine if enterprise resources should be included based on user preferences and current input"""
    mode = inputs.enterprise_resources_mode or "auto_when_missing"
    
    if mode == "always_include":
        return True
    elif mode == "never_auto_include":
        return False
    else:  # "auto_when_missing"
        # Check if any enterprise resources are missing from user's explicit selections
        all_selected_services = set()
        all_selected_services.update(inputs.security_services or [])
        all_selected_services.update(inputs.network_services or [])
        all_selected_services.update(inputs.monitoring_services or [])
        
        # Check if any enterprise resources are missing
        missing_resources = [res for res in enterprise_resources if res not in all_selected_services]
        return len(missing_resources) > 0

def _ensure_enterprise_resources_included(inputs: CustomerInputs) -> CustomerInputs:
    """Ensure enterprise resources are included in the input based on user preferences"""
    enterprise_resources = _get_enterprise_resources()
    
    if not _should_include_enterprise_resources(inputs, enterprise_resources):
        return inputs
    
    # Create a copy of inputs to modify
    modified_inputs = inputs.model_copy()
    
    # Ensure lists are initialized
    if not modified_inputs.security_services:
        modified_inputs.security_services = []
    if not modified_inputs.network_services:
        modified_inputs.network_services = []
    if not modified_inputs.monitoring_services:
        modified_inputs.monitoring_services = []
    
    # Add missing enterprise resources
    if "key_vault" not in modified_inputs.security_services:
        modified_inputs.security_services.append("key_vault")
        logger.info("Auto-included Key Vault for enterprise compliance")
    
    if "active_directory" not in modified_inputs.security_services:
        modified_inputs.security_services.append("active_directory")
        logger.info("Auto-included Active Directory for enterprise compliance")
    
    if "firewall" not in modified_inputs.network_services:
        modified_inputs.network_services.append("firewall")
        logger.info("Auto-included Azure Firewall for enterprise compliance")
    
    if "monitor" not in modified_inputs.monitoring_services:
        modified_inputs.monitoring_services.append("monitor")
        logger.info("Auto-included Azure Monitor for enterprise compliance")
    
    return modified_inputs

def _get_user_prompt_for_enterprise_resources() -> str:
    """Generate user prompt about enterprise resource preferences"""
    return ("Would you like to always include these enterprise resources (Key Vault, Active Directory, "
            "Firewall, Azure Monitor) in every diagram, or only when not specified? "
            "Should connections between these resources and VM/VNet be explicitly shown?")

def _add_enterprise_resource_connections(inputs: CustomerInputs, hub_vnet, prod_vnet, dev_vnet, aad, key_vault, network_services):
    """Add comprehensive connections between enterprise resources and infrastructure components"""
    logger.info("Adding enhanced enterprise resource connections")
    
    # Find specific network services
    firewall_service = None
    vpn_gateway = None
    app_gateway = None
    load_balancer = None
    
    for ns in network_services:
        if hasattr(ns, '_name'):
            if 'Firewall' in ns._name:
                firewall_service = ns
            elif 'VPN' in ns._name or 'Gateway' in ns._name:
                vpn_gateway = ns
            elif 'Application' in ns._name or 'AppGateway' in ns._name:
                app_gateway = ns
            elif 'LoadBalancer' in ns._name or 'Balancer' in ns._name:
                load_balancer = ns
    
    # 1. Core Security Connections - Key Vault to all environments
    try:
        key_vault >> Edge(label="Secrets", style="bold", color="orange") >> prod_vnet
        key_vault >> Edge(label="Secrets", style="bold", color="orange") >> dev_vnet
        if hub_vnet:
            key_vault >> Edge(label="Hub Secrets", style="bold", color="orange") >> hub_vnet
        logger.debug("Connected Key Vault to all VNets with labeled connections")
    except Exception as e:
        logger.debug(f"Key Vault connection to VNets: {e}")
    
    # 2. Identity Connections - Azure AD to all environments
    try:
        aad >> Edge(label="Authentication", style="bold", color="blue") >> prod_vnet
        aad >> Edge(label="Authentication", style="bold", color="blue") >> dev_vnet
        if hub_vnet:
            aad >> Edge(label="Hub Identity", style="bold", color="blue") >> hub_vnet
        logger.debug("Connected Active Directory to all VNets with authentication labels")
    except Exception as e:
        logger.debug(f"Active Directory connection to VNets: {e}")
    
    # 3. Enhanced Firewall Connections - Security traffic routing
    if firewall_service:
        try:
            # Firewall as central security gateway
            firewall_service >> Edge(label="Secure Routing", style="bold", color="red") >> prod_vnet
            firewall_service >> Edge(label="Secure Routing", style="bold", color="red") >> dev_vnet
            
            # Bi-directional traffic flow indication
            prod_vnet >> Edge(label="Return Traffic", style="dashed", color="red") >> firewall_service
            dev_vnet >> Edge(label="Return Traffic", style="dashed", color="red") >> firewall_service
            logger.debug("Connected Firewall with bi-directional traffic flow")
        except Exception as e:
            logger.debug(f"Firewall connection to VNets: {e}")
    
    # 4. Security Services Integration - Central security orchestration
    if firewall_service:
        try:
            # Identity-driven security policies
            aad >> Edge(label="Security Policies", style="dotted", color="purple") >> firewall_service
            # Key management for firewall
            key_vault >> Edge(label="Certificates", style="dotted", color="orange") >> firewall_service
            logger.debug("Connected security services with policy and certificate flows")
        except Exception as e:
            logger.debug(f"Security services interconnection: {e}")
    
    # 5. Gateway Connections - Hybrid connectivity
    if vpn_gateway:
        try:
            vpn_gateway >> Edge(label="Hybrid Access", style="bold", color="green") >> hub_vnet
            # Connect gateway to security services
            aad >> Edge(label="VPN Auth", style="dotted", color="blue") >> vpn_gateway
            if firewall_service:
                vpn_gateway >> Edge(label="Gateway Traffic", style="solid", color="green") >> firewall_service
            logger.debug("Connected VPN Gateway with hybrid access patterns")
        except Exception as e:
            logger.debug(f"VPN Gateway connections: {e}")
    
    # 6. Application Gateway Connections - Web application delivery
    if app_gateway:
        try:
            app_gateway >> Edge(label="Web Traffic", style="bold", color="cyan") >> prod_vnet
            key_vault >> Edge(label="SSL Certs", style="dotted", color="orange") >> app_gateway
            if firewall_service:
                app_gateway >> Edge(label="Security Scanning", style="dashed", color="red") >> firewall_service
            logger.debug("Connected Application Gateway with web delivery patterns")
        except Exception as e:
            logger.debug(f"Application Gateway connections: {e}")
    
    # 7. Load Balancer Connections - Traffic distribution
    if load_balancer:
        try:
            load_balancer >> Edge(label="Load Distribution", style="bold", color="navy") >> prod_vnet
            load_balancer >> Edge(label="Health Checks", style="dotted", color="navy") >> dev_vnet
            logger.debug("Connected Load Balancer with traffic distribution patterns")
        except Exception as e:
            logger.debug(f"Load Balancer connections: {e}")
    
    logger.info("Enhanced enterprise resource connections completed")

def _add_service_to_service_connections(inputs: CustomerInputs, service_collections, aad, key_vault):
    """Add intelligent connections between different service tiers with comprehensive patterns"""
    logger.info("Adding service-to-service connections")
    
    # Extract service collections
    compute_services = service_collections.get('compute_services', [])
    storage_services = service_collections.get('storage_services', [])
    database_services = service_collections.get('database_services', [])
    analytics_services = service_collections.get('analytics_services', [])
    monitoring_services = service_collections.get('monitoring_services', [])
    
    try:
        # 1. Compute to Storage Connections - Data persistence patterns
        if compute_services and storage_services:
            for compute in compute_services:
                for storage in storage_services:
                    compute >> Edge(label="Data Storage", style="solid", color="darkgreen") >> storage
                    logger.debug(f"Connected {compute} to {storage} for data persistence")
        
        # 2. Compute to Database Connections - Application data patterns
        if compute_services and database_services:
            for compute in compute_services:
                for database in database_services:
                    compute >> Edge(label="Database Access", style="bold", color="darkblue") >> database
                    logger.debug(f"Connected {compute} to {database} for application data")
        
        # 3. Identity Integration - All services need authentication
        if aad:
            # Connect identity to compute services
            for compute in compute_services:
                aad >> Edge(label="Service Identity", style="dotted", color="blue") >> compute
            
            # Connect identity to data services for access control
            for database in database_services:
                aad >> Edge(label="Data Access Control", style="dotted", color="blue") >> database
            
            for storage in storage_services:
                aad >> Edge(label="Storage Access Control", style="dotted", color="blue") >> storage
        
        # 4. Key Vault Integration - Secrets management for all services
        if key_vault:
            # Connect Key Vault to compute services for application secrets
            for compute in compute_services:
                key_vault >> Edge(label="App Secrets", style="dashed", color="orange") >> compute
            
            # Connect Key Vault to databases for connection strings and keys
            for database in database_services:
                key_vault >> Edge(label="DB Connection Strings", style="dashed", color="orange") >> database
            
            # Connect Key Vault to storage for access keys
            for storage in storage_services:
                key_vault >> Edge(label="Storage Keys", style="dashed", color="orange") >> storage
        
        # 5. Analytics Integration Patterns - Data flows to analytics
        if analytics_services:
            for analytics in analytics_services:
                # Analytics consume data from storage
                for storage in storage_services:
                    storage >> Edge(label="Analytics Data", style="bold", color="purple") >> analytics
                # Analytics consume data from databases
                for database in database_services:
                    database >> Edge(label="Data Pipeline", style="bold", color="purple") >> analytics
                logger.debug(f"Connected data sources to {analytics}")
        
        # 6. AI/ML Service Integration
        if inputs.ai_services:
            for ai_service_key in inputs.ai_services:
                if ai_service_key in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[ai_service_key]["diagram_class"]:
                    logger.debug(f"AI service {ai_service_key} available for data connections")
                    # AI services need data and compute resources
                    # This would be implemented if AI services were in service_collections
        
        # 7. Monitoring has already been handled in the main function
        
        # 8. DevOps Integration Patterns
        if inputs.devops_services:
            logger.debug("DevOps services deploy to compute resources")
            # DevOps connections would be handled at the infrastructure level
        
        # 9. Integration Services Patterns - Event-driven connectivity
        if inputs.integration_services:
            logger.debug("Integration services enable event-driven architecture")
            # These would connect compute services through messaging/events
        
        # 10. Cross-service communication patterns
        logger.debug("Added comprehensive cross-service connectivity patterns")
        
        logger.info("Service-to-service connections completed successfully")
        
    except Exception as e:
        logger.error(f"Error adding service-to-service connections: {e}")
        logger.debug(traceback.format_exc())

def generate_azure_architecture_diagram(inputs: CustomerInputs, output_dir: str = None, format: str = "png") -> str:
    """Generate Azure architecture diagram using the Python Diagrams library with proper Azure icons"""
    
    logger.info("Starting Azure architecture diagram generation")
    
    try:
        # Validate inputs first
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Apply enterprise resource auto-inclusion logic
        inputs = _ensure_enterprise_resources_included(inputs)
        
        # Get safe output directory
        if output_dir is None:
            output_dir = get_safe_output_directory()
        
        # Clean up old files to prevent disk space issues
        cleanup_old_files(output_dir)
        
        # Verify Graphviz availability before proceeding
        try:
            result = subprocess.run(['dot', '-V'], capture_output=True, text=True, timeout=10)
            if result.returncode != 0:
                raise Exception(f"Graphviz 'dot' command failed with return code {result.returncode}. stderr: {result.stderr}")
            logger.info(f"Graphviz version: {result.stderr.strip()}")
        except subprocess.TimeoutExpired:
            raise Exception("Graphviz 'dot' command timed out. Graphviz may be unresponsive.")
        except FileNotFoundError:
            raise Exception("Graphviz is not installed or not accessible. Please install Graphviz: sudo apt-get install -y graphviz graphviz-dev")
        except subprocess.SubprocessError as e:
            raise Exception(f"Graphviz check failed: {str(e)}. Please install Graphviz: sudo apt-get install -y graphviz graphviz-dev")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = str(uuid.uuid4())[:8]  # Use first 8 chars of UUID for uniqueness
        filename = f"azure_landing_zone_{timestamp}_{unique_id}"
        filepath = os.path.join(output_dir, filename)
        
        logger.info(f"Generating diagram with filename: {filename}")
        
        # Verify output directory is writable
        if not os.access(output_dir, os.W_OK):
            raise Exception(f"Output directory {output_dir} is not writable")
        
        # Determine organization template
        template = generate_architecture_template(inputs)
        org_name = inputs.org_structure or "Enterprise"
        
        logger.info(f"Using template: {template['template']['name']}")
        
        try:
            # Set the format based on the requested output format
            output_format = "svg" if format.lower() == "svg" else "png"
            
            with Diagram(
                f"Azure Landing Zone - {template['template']['name']}", 
                filename=filepath, 
                show=False, 
                direction="TB",
                outformat=output_format,
                graph_attr={
                    "fontsize": "16",
                    "fontname": "Arial",
                    "rankdir": "TB",
                    "nodesep": "1.0",
                    "ranksep": "1.5",
                    "bgcolor": "#ffffff",
                    "margin": "0.5"
                },
                node_attr={
                    "fontsize": "12",
                    "fontname": "Arial"
                },
                edge_attr={
                    "fontsize": "10",
                    "fontname": "Arial"
                }
            ):
                
                logger.info("Creating diagram structure...")
                
                # Core Identity and Security Services
                with Cluster("Identity & Security", graph_attr={"bgcolor": "#e8f4f8", "style": "rounded"}):
                    aad = ActiveDirectory("Azure Active Directory")
                    key_vault = KeyVaults("Key Vault")
                    if inputs.security_services and "security_center" in inputs.security_services:
                        sec_center = SecurityCenter("Security Center")
                    if inputs.security_services and "sentinel" in inputs.security_services:
                        sentinel = Sentinel("Sentinel")
                
                # Management Groups and Subscriptions Structure
                with Cluster("Management & Governance", graph_attr={"bgcolor": "#f0f8ff", "style": "rounded"}):
                    root_mg = Subscriptions("Root Management Group")
                    if template['template']['name'] == "Enterprise Scale Landing Zone":
                        platform_mg = Subscriptions("Platform MG")
                        workloads_mg = Subscriptions("Landing Zones MG")
                        root_mg >> [platform_mg, workloads_mg]
                    else:
                        platform_mg = Subscriptions("Platform MG")
                        workloads_mg = Subscriptions("Workloads MG")
                        root_mg >> [platform_mg, workloads_mg]
                
                # Hub-Spoke Network Architecture with Enhanced Visualization
                with Cluster("Network Architecture", graph_attr={"bgcolor": "#f0fff0", "style": "rounded"}):
                    # Hub VNet with dashed border to indicate central hub
                    with Cluster("Hub VNet", graph_attr={
                        "bgcolor": "#e6f7ff", 
                        "style": "dashed", 
                        "color": "#0078d4",
                        "penwidth": "2"
                    }):
                        hub_vnet = VirtualNetworks("Hub VNet\n(Shared Services)")
                        
                        # Network services based on selections - these go in the hub
                        network_services = []
                        if inputs.network_services:
                            for service in inputs.network_services:
                                if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                                    diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                                    service_name = AZURE_SERVICES_MAPPING[service]["name"]
                                    network_services.append(diagram_class(service_name))
                        
                        # Default network services if none specified
                        if not network_services:
                            firewall = Firewall("Azure Firewall")
                            vpn_gw = VirtualNetworkGateways("VPN Gateway")
                            network_services = [firewall, vpn_gw]
                        
                        # Connect network services within hub
                        for ns in network_services:
                            hub_vnet >> Edge(style="solid", color="#0078d4") >> ns
                    
                    # Spoke VNets with dashed borders to indicate spokes
                    with Cluster("Production Spoke", graph_attr={
                        "bgcolor": "#fff7e6", 
                        "style": "dashed", 
                        "color": "#d83b01",
                        "penwidth": "2"
                    }):
                        prod_vnet = VirtualNetworks("Production VNet")
                    
                    with Cluster("Development Spoke", graph_attr={
                        "bgcolor": "#f0f8e6", 
                        "style": "dashed", 
                        "color": "#107c10",
                        "penwidth": "2"
                    }):
                        dev_vnet = VirtualNetworks("Development VNet")
                    
                    # Connect hub to spokes with dashed lines to show hub-spoke topology
                    hub_vnet >> Edge(style="dashed", color="#666666", label="Hub-Spoke\nConnection") >> prod_vnet
                    hub_vnet >> Edge(style="dashed", color="#666666", label="Hub-Spoke\nConnection") >> dev_vnet
                    
                    # Connect platform subscription to hub
                    platform_mg >> hub_vnet
                
                # Add other service clusters based on input and get service collections for connectivity
                service_collections = _add_service_clusters(inputs, prod_vnet, workloads_mg)
                
                # Core security connections
                aad >> key_vault
                platform_mg >> [aad, key_vault]
                
                # Enhanced enterprise resource connections (when enabled)
                if inputs.show_enterprise_connections:
                    _add_enterprise_resource_connections(inputs, hub_vnet, prod_vnet, dev_vnet, aad, key_vault, network_services)
                
                # Add comprehensive service-to-service connections
                _add_service_to_service_connections(
                    inputs, 
                    service_collections,
                    aad, 
                    key_vault
                )
                
                # Add monitoring connections to all services (monitoring everything)
                if service_collections.get('monitoring_services'):
                    for monitor in service_collections['monitoring_services']:
                        # Monitor all compute services
                        for compute in service_collections.get('compute_services', []):
                            monitor >> Edge(label="Monitoring", style="dotted", color="green") >> compute
                        # Monitor all databases  
                        for database in service_collections.get('database_services', []):
                            monitor >> Edge(label="DB Monitoring", style="dotted", color="green") >> database
                        # Monitor all storage
                        for storage in service_collections.get('storage_services', []):
                            monitor >> Edge(label="Storage Monitoring", style="dotted", color="green") >> storage
                
                logger.info("Diagram structure created successfully")
        
        except Exception as e:
            logger.error(f"Error during diagram creation: {str(e)}")
            logger.error(traceback.format_exc())
            raise Exception(f"Error generating Azure architecture diagram: {str(e)}")
        
        # Return the file path of the generated diagram
        if format.lower() == "svg":
            # Check for SVG file generated directly by diagrams library
            svg_path = f"{filepath}.svg"
            if os.path.exists(svg_path):
                file_size = os.path.getsize(svg_path)
                logger.info(f"SVG diagram generated successfully: {svg_path} (size: {file_size} bytes)")
                return svg_path
            else:
                # Fallback: try to generate SVG using dot command from gv file
                dot_path = f"{filepath}.gv" 
                if os.path.exists(dot_path):
                    try:
                        # Convert dot file to SVG
                        result = subprocess.run(['dot', '-Tsvg', dot_path, '-o', svg_path], 
                                              capture_output=True, text=True, timeout=30)
                        if result.returncode != 0:
                            raise Exception(f"SVG generation failed: {result.stderr}")
                        
                        if os.path.exists(svg_path):
                            file_size = os.path.getsize(svg_path)
                            logger.info(f"SVG diagram generated successfully via dot: {svg_path} (size: {file_size} bytes)")
                            return svg_path
                        else:
                            raise Exception(f"SVG generation failed - file not found: {svg_path}")
                    except subprocess.TimeoutExpired:
                        raise Exception("SVG generation timed out")
                    except Exception as e:
                        raise Exception(f"Failed to generate SVG: {str(e)}")
                else:
                    raise Exception(f"Neither SVG nor dot file found: {svg_path}, {dot_path}")
        else:
            # Default PNG generation
            png_path = f"{filepath}.png"
            if os.path.exists(png_path):
                file_size = os.path.getsize(png_path)
                logger.info(f"Diagram generated successfully: {png_path} (size: {file_size} bytes)")
                return png_path
            else:
                raise Exception(f"Diagram generation failed - PNG file not found: {png_path}")
            
    except Exception as e:
        logger.error(f"Failed to generate Azure architecture diagram: {str(e)}")
        logger.error(traceback.format_exc())
        raise

def generate_simple_svg_diagram(inputs: CustomerInputs) -> str:
    """Generate a simple SVG diagram as fallback when Python Diagrams fails"""
    
    template = generate_architecture_template(inputs)
    template_name = template['template']['name']
    
    # Create a simple SVG representation
    svg_width = 800
    svg_height = 600
    
    svg_content = f'''<svg width="{svg_width}" height="{svg_height}" xmlns="http://www.w3.org/2000/svg">
    <defs>
        <style>
            .title {{ font-family: Arial, sans-serif; font-size: 18px; font-weight: bold; fill: #0078d4; }}
            .group-title {{ font-family: Arial, sans-serif; font-size: 14px; font-weight: bold; fill: #323130; }}
            .service {{ font-family: Arial, sans-serif; font-size: 12px; fill: #605e5c; }}
            .mgmt-group {{ fill: #e1f5fe; stroke: #0078d4; stroke-width: 2; }}
            .subscription {{ fill: #f3e5f5; stroke: #6b69d6; stroke-width: 2; }}
            .service-box {{ fill: #fff3e0; stroke: #d83b01; stroke-width: 1; cursor: pointer; }}
            .service-box:hover {{ fill: #ffebdd; }}
            .network-box {{ fill: #e8f5e8; stroke: #107c10; stroke-width: 1; cursor: pointer; }}
            .network-box:hover {{ fill: #f3fdf3; }}
            .security-box {{ fill: #ffebee; stroke: #d13438; stroke-width: 1; cursor: pointer; }}
            .security-box:hover {{ fill: #fdf3f4; }}
        </style>
    </defs>
    
    <!-- Background -->
    <rect width="100%" height="100%" fill="#f8f9fa"/>
    
    <!-- Title -->
    <text x="400" y="30" class="title" text-anchor="middle">Azure Landing Zone - {template_name}</text>
    
    <!-- Azure Tenant Container -->
    <rect x="50" y="60" width="700" height="520" fill="none" stroke="#0078d4" stroke-width="3" stroke-dasharray="5,5"/>
    <text x="60" y="80" class="group-title">Azure Tenant</text>
    
    <!-- Management Groups -->
    <rect x="70" y="100" width="200" height="150" class="mgmt-group" rx="5"/>
    <text x="80" y="120" class="group-title">Management Groups</text>
    
    <!-- Management Group Items -->
    <rect x="80" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="120" y="149" class="service" text-anchor="middle">Root MG</text>
    
    <rect x="170" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="210" y="149" class="service" text-anchor="middle">Platform</text>
    
    <rect x="80" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="120" y="189" class="service" text-anchor="middle">Landing Zones</text>
    
    <rect x="170" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="210" y="189" class="service" text-anchor="middle">Sandbox</text>
    
    <!-- Subscriptions -->
    <rect x="290" y="100" width="200" height="150" class="subscription" rx="5"/>
    <text x="300" y="120" class="group-title">Subscriptions</text>
    
    <!-- Subscription Items -->
    <rect x="300" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="340" y="149" class="service" text-anchor="middle">Connectivity</text>
    
    <rect x="390" y="130" width="80" height="30" class="service-box" rx="3"/>
    <text x="430" y="149" class="service" text-anchor="middle">Identity</text>
    
    <rect x="300" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="340" y="189" class="service" text-anchor="middle">Production</text>
    
    <rect x="390" y="170" width="80" height="30" class="service-box" rx="3"/>
    <text x="430" y="189" class="service" text-anchor="middle">Development</text>
    
    <!-- Network Architecture -->
    <rect x="520" y="100" width="200" height="150" class="network-box" rx="5"/>
    <text x="530" y="120" class="group-title">Network Architecture</text>
    
    <rect x="530" y="130" width="80" height="30" class="network-box" rx="3"/>
    <text x="570" y="149" class="service" text-anchor="middle">Hub VNet</text>
    
    <rect x="620" y="130" width="80" height="30" class="network-box" rx="3"/>
    <text x="660" y="149" class="service" text-anchor="middle">Spoke VNet</text>'''
    
    # Add selected services
    y_offset = 280
    if inputs.compute_services:
        svg_content += f'''
    <!-- Compute Services -->
    <rect x="70" y="{y_offset}" width="300" height="80" class="service-box" rx="5"/>
    <text x="80" y="{y_offset + 20}" class="group-title">Compute Services</text>'''
        
        x_pos = 80
        for i, service in enumerate(inputs.compute_services[:4]):  # Max 4 services
            service_name = service.replace('_', ' ').title()
            svg_content += f'''
    <rect x="{x_pos}" y="{y_offset + 30}" width="60" height="25" class="service-box" rx="3"/>
    <text x="{x_pos + 30}" y="{y_offset + 47}" class="service" text-anchor="middle" font-size="10">{service_name[:8]}</text>'''
            x_pos += 70
    
    if inputs.network_services:
        svg_content += f'''
    <!-- Network Services -->
    <rect x="390" y="{y_offset}" width="300" height="80" class="network-box" rx="5"/>
    <text x="400" y="{y_offset + 20}" class="group-title">Network Services</text>'''
        
        x_pos = 400
        for i, service in enumerate(inputs.network_services[:4]):  # Max 4 services
            service_name = service.replace('_', ' ').title()
            svg_content += f'''
    <rect x="{x_pos}" y="{y_offset + 30}" width="60" height="25" class="network-box" rx="3"/>
    <text x="{x_pos + 30}" y="{y_offset + 47}" class="service" text-anchor="middle" font-size="10">{service_name[:8]}</text>'''
            x_pos += 70
    
    # Security Services
    y_offset += 100
    svg_content += f'''
    <!-- Security & Identity -->
    <rect x="70" y="{y_offset}" width="620" height="80" class="security-box" rx="5"/>
    <text x="80" y="{y_offset + 20}" class="group-title">Security & Identity Services</text>
    
    <rect x="80" y="{y_offset + 30}" width="100" height="25" class="security-box" rx="3"/>
    <text x="130" y="{y_offset + 47}" class="service" text-anchor="middle">Azure AD</text>
    
    <rect x="190" y="{y_offset + 30}" width="100" height="25" class="security-box" rx="3"/>
    <text x="240" y="{y_offset + 47}" class="service" text-anchor="middle">Key Vault</text>
    
    <rect x="300" y="{y_offset + 30}" width="100" height="25" class="security-box" rx="3"/>
    <text x="350" y="{y_offset + 47}" class="service" text-anchor="middle">Security Center</text>'''
    
    # Add connections
    svg_content += '''
    <!-- Connections -->
    <line x1="170" y1="145" x2="290" y2="145" stroke="#666" stroke-width="2" marker-end="url(#arrowhead)"/>
    <line x1="390" y1="145" x2="520" y2="145" stroke="#666" stroke-width="2" marker-end="url(#arrowhead)"/>
    
    <!-- Arrow marker -->
    <defs>
        <marker id="arrowhead" markerWidth="10" markerHeight="7" refX="10" refY="3.5" orient="auto">
            <polygon points="0 0, 10 3.5, 0 7" fill="#666"/>
        </marker>
    </defs>
    
    <!-- Footer -->
    <text x="400" y="570" class="service" text-anchor="middle" fill="#8a8886">Generated by Azure Landing Zone Agent - Interactive Mode</text>
</svg>'''
    
    return svg_content

def _add_service_clusters(inputs: CustomerInputs, prod_vnet, workloads_mg):
    """Helper method to add service clusters to avoid code duplication and return service references for connectivity"""
    service_collections = {
        'compute_services': [],
        'storage_services': [],
        'database_services': [],
        'analytics_services': [],
        'monitoring_services': []
    }
    
    try:
        # Compute and Application Services - These go in the spoke with enhanced visualization
        if inputs.compute_services or inputs.workload:
            with Cluster("Compute & Applications (Spoke)", graph_attr={
                "bgcolor": "#fff8dc", 
                "style": "dashed", 
                "color": "#d83b01",
                "penwidth": "2",
                "label": "Spoke Workloads"
            }):
                compute_services = []
                
                # Add selected compute services
                if inputs.compute_services:
                    for service in inputs.compute_services:
                        if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                            diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                            service_name = AZURE_SERVICES_MAPPING[service]["name"]
                            service_instance = diagram_class(service_name)
                            compute_services.append(service_instance)
                
                # Fallback to workload if no specific compute services
                if not compute_services and inputs.workload:
                    if inputs.workload in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[inputs.workload]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[inputs.workload]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[inputs.workload]["name"]
                        compute_services.append(diagram_class(service_name))
                
                # Default to App Services if nothing specified
                if not compute_services:
                    compute_services.append(AppServices("Azure App Services"))
                
                # Store for connectivity
                service_collections['compute_services'] = compute_services
                
                # Connect compute services to production VNet with enhanced labels
                for cs in compute_services:
                    prod_vnet >> Edge(style="dashed", color="#d83b01", label="Spoke\nWorkload") >> cs
                    workloads_mg >> Edge(style="dotted", color="#666666") >> cs
        
        # Storage Services - Also in spoke
        if inputs.storage_services:
            with Cluster("Storage & Data (Spoke)", graph_attr={
                "bgcolor": "#f5f5dc", 
                "style": "dashed", 
                "color": "#d83b01",
                "penwidth": "2"
            }):
                storage_services = []
                for service in inputs.storage_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        service_instance = diagram_class(service_name)
                        storage_services.append(service_instance)
                
                if not storage_services:
                    storage_services.append(StorageAccounts("Storage Accounts"))
                
                # Store for connectivity
                service_collections['storage_services'] = storage_services
                
                # Connect storage to production VNet with enhanced labels
                for ss in storage_services:
                    prod_vnet >> Edge(style="dashed", color="#d83b01", label="Spoke\nData") >> ss
                    workloads_mg >> Edge(style="dotted", color="#666666") >> ss
        
        # Database Services - Also in spoke
        if inputs.database_services:
            with Cluster("Databases (Spoke)", graph_attr={
                "bgcolor": "#e6f3ff", 
                "style": "dashed", 
                "color": "#d83b01",
                "penwidth": "2"
            }):
                database_services = []
                for service in inputs.database_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        service_instance = diagram_class(service_name)
                        database_services.append(service_instance)
                
                # Store for connectivity
                service_collections['database_services'] = database_services
                
                # Connect databases to production VNet with enhanced labels
                for ds in database_services:
                    prod_vnet >> Edge(style="dashed", color="#d83b01", label="Spoke\nDatabase") >> ds
                    workloads_mg >> Edge(style="dotted", color="#666666") >> ds
        
        # Analytics Services
        if inputs.analytics_services:
            with Cluster("Analytics & AI", graph_attr={"bgcolor": "#f0e6ff", "style": "rounded"}):
                analytics_services = []
                for service in inputs.analytics_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        service_instance = diagram_class(service_name)
                        analytics_services.append(service_instance)
                
                # Store for connectivity
                service_collections['analytics_services'] = analytics_services
                
                # Connect analytics to production VNet and workloads with data flow labels
                for as_service in analytics_services:
                    prod_vnet >> Edge(label="Analytics Data", style="solid", color="purple") >> as_service
                    workloads_mg >> as_service
        
        # Integration Services
        if inputs.integration_services:
            with Cluster("Integration", graph_attr={"bgcolor": "#fff0e6", "style": "rounded"}):
                integration_services = []
                for service in inputs.integration_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        integration_services.append(diagram_class(service_name))
                
                # Connect integration services to production VNet and workloads
                for is_service in integration_services:
                    prod_vnet >> Edge(label="Integration Flow", style="solid", color="teal") >> is_service
                    workloads_mg >> is_service
        
        # DevOps Services  
        if inputs.devops_services:
            with Cluster("DevOps & Automation", graph_attr={"bgcolor": "#f5f5f5", "style": "rounded"}):
                devops_services = []
                for service in inputs.devops_services:
                    if service in AZURE_SERVICES_MAPPING and AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                        diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                        service_name = AZURE_SERVICES_MAPPING[service]["name"]
                        devops_services.append(diagram_class(service_name))
                
                # Connect DevOps services to management
                for ds in devops_services:
                    workloads_mg >> Edge(label="CI/CD Pipeline", style="dotted", color="gray") >> ds
        
        # Monitoring & Management Services
        if inputs.monitoring_services:
            with Cluster("Monitoring & Observability", graph_attr={"bgcolor": "#e8f4f8", "style": "rounded"}):
                monitoring_services_list = []
                for service in inputs.monitoring_services:
                    if service in AZURE_SERVICES_MAPPING:
                        if AZURE_SERVICES_MAPPING[service]["diagram_class"]:
                            diagram_class = AZURE_SERVICES_MAPPING[service]["diagram_class"]
                            service_name = AZURE_SERVICES_MAPPING[service]["name"]
                            service_instance = diagram_class(service_name)
                            monitoring_services_list.append(service_instance)
                        else:
                            logger.info(f"Monitoring service '{service}' included but no visual diagram component available")
                
                # Store for connectivity
                service_collections['monitoring_services'] = monitoring_services_list
                
                # Connect monitoring services to VNets and management with observability labels
                for ms in monitoring_services_list:
                    prod_vnet >> Edge(label="Telemetry", style="dotted", color="darkgreen") >> ms
                    workloads_mg >> ms
    
    except Exception as e:
        logger.warning(f"Error adding service clusters: {str(e)}")
        # Don't fail the entire diagram generation for service cluster issues
    
    return service_collections


def generate_architecture_template(inputs: CustomerInputs) -> Dict[str, Any]:
    """Generate architecture template based on inputs"""
    
    # Determine organization size and template
    if inputs.org_structure and "enterprise" in inputs.org_structure.lower():
        template = AZURE_TEMPLATES["enterprise"]
    elif inputs.org_structure and any(x in inputs.org_structure.lower() for x in ["small", "medium", "sme"]):
        template = AZURE_TEMPLATES["small_medium"]
    else:
        template = AZURE_TEMPLATES["startup"]
    
    # Build architecture components
    components = {
        "template": template,
        "identity": inputs.identity or "Azure AD",
        "network_model": inputs.network_model or "hub-spoke",
        "workload": inputs.workload or "app-services",
        "security": inputs.security_posture or "zero-trust",
        "monitoring": inputs.monitoring or "azure-monitor",
        "governance": inputs.governance or "azure-policy"
    }
    
    return components

def generate_professional_mermaid(inputs: CustomerInputs) -> str:
    """Generate professional Mermaid diagram for Azure Landing Zone with Hub-and-Spoke architecture"""
    
    template = generate_architecture_template(inputs)
    network_model = inputs.network_model or "hub-spoke"
    
    lines = [
        "graph TD",
        "    %% Professional Hub-and-Spoke Azure Landing Zone Architecture",
        "    subgraph \"Azure Landing Zone Architecture\"",
        "",
        "        %% Cross-Premises Connectivity",
        "        subgraph \"Cross-Premises\" [\"🏢 Cross-Premises Connectivity\"]",
        "            ONPREM[\"🏢 On-Premises<br/>Data Center\"]",
        "            INTERNET[\"🌐 Internet<br/>Users\"]"
    ]
    
    # Add ExpressRoute or VPN if specified
    if inputs.network_services and any(svc in ["expressroute", "vpn_gateway"] for svc in inputs.network_services):
        if "expressroute" in inputs.network_services:
            lines.append("            ER[\"⚡ ExpressRoute<br/>Private Connection\"]")
        if "vpn_gateway" in inputs.network_services:
            lines.append("            VPN[\"🔒 VPN Gateway<br/>Site-to-Site\"]")
    
    lines.extend([
        "        end",
        "",
        "        %% Hub Virtual Network (Central Hub)",
        "        subgraph \"Hub\" [\"🏛️ Hub Virtual Network\"]",
        "            direction TB",
        "",
        "            %% Network Security & Monitoring in Hub",
        "            subgraph \"HubSecurity\" [\"🛡️ Network Security\"]",
        "                FIREWALL[\"🛡️ Azure Firewall<br/>Central Security\"]"
    ])
    
    # Add network security services in hub
    if inputs.security_services:
        for service in inputs.security_services:
            if service in AZURE_SERVICES_MAPPING and service in ["security_center", "sentinel", "defender"]:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"                {service_id}[\"{service_info['icon']} {service_info['name']}<br/>Security Monitoring\"]")
    
    lines.extend([
        "            end",
        "",
        "            %% Shared Services in Hub",
        "            subgraph \"SharedServices\" [\"⚙️ Shared Services\"]",
        "                DNS[\"🌐 Private DNS<br/>Name Resolution\"]",
        "                BASTION[\"🔐 Azure Bastion<br/>Secure Access\"]"
    ])
    
    # Add monitoring services in hub
    if inputs.monitoring_services:
        for service in inputs.monitoring_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"                {service_id}[\"{service_info['icon']} {service_info['name']}<br/>Centralized Monitoring\"]")
    
    lines.extend([
        "            end",
        "",
        "            %% Gateway Services",
        "            subgraph \"Gateways\" [\"🚪 Gateway Services\"]"
    ])
    
    # Add gateways based on network services
    if inputs.network_services:
        for service in inputs.network_services:
            if service in ["application_gateway", "load_balancer", "vpn_gateway"]:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = service.upper().replace("_", "")
                lines.append(f"                {service_id}[\"{service_info['icon']} {service_info['name']}<br/>Traffic Management\"]")
    
    lines.extend([
        "            end",
        "        end",
        "",
        "        %% Production Spoke",
        "        subgraph \"ProdSpoke\" [\"🏭 Production Spoke\"]",
        "            direction TB",
        "            PRODVNET[\"🌐 Production VNet<br/>Workload Network\"]",
        "",
        "            %% Production Workloads",
        "            subgraph \"ProdWorkloads\" [\"💼 Production Workloads\"]"
    ])
    
    # Add compute services in production spoke
    if inputs.compute_services:
        for service in inputs.compute_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = f"PROD_{service.upper().replace('_', '')}"
                lines.append(f"                {service_id}[\"{service_info['icon']} {service_info['name']}<br/>Production\"]")
    
    # Add database services in production spoke
    if inputs.database_services:
        lines.append("            end")
        lines.append("")
        lines.append("            %% Production Data Services")
        lines.append("            subgraph \"ProdData\" [\"🗄️ Data Services\"]")
        for service in inputs.database_services:
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                service_id = f"PROD_{service.upper().replace('_', '')}"
                lines.append(f"                {service_id}[\"{service_info['icon']} {service_info['name']}<br/>Production Data\"]")
    
    lines.extend([
        "            end",
        "        end",
        "",
        "        %% Development Spoke",
        "        subgraph \"DevSpoke\" [\"👩‍💻 Development Spoke\"]",
        "            direction TB", 
        "            DEVVNET[\"🌐 Development VNet<br/>Dev/Test Network\"]",
        "",
        "            %% Development Workloads", 
        "            subgraph \"DevWorkloads\" [\"🧪 Development Workloads\"]"
    ])
    
    # Add simplified dev services
    if inputs.compute_services:
        # Just show primary compute service for dev to avoid clutter
        primary_compute = inputs.compute_services[0] if inputs.compute_services else "app_services"
        if primary_compute in AZURE_SERVICES_MAPPING:
            service_info = AZURE_SERVICES_MAPPING[primary_compute]
            lines.append(f"                DEV_COMPUTE[\"{service_info['icon']} Dev/Test<br/>{service_info['name']}\"]")
    
    lines.extend([
        "            end",
        "        end",
        "",
        "        %% Identity & Management Layer",
        "        subgraph \"Identity\" [\"🔐 Identity & Management\"]",
        "            AAD[\"👤 Azure Active Directory<br/>Identity Provider\"]",
        "            KEYVAULT[\"🔐 Azure Key Vault<br/>Secrets Management\"]"
    ])
    
    # Add governance and management services
    if template["template"]["name"] == "Enterprise Scale Landing Zone":
        lines.extend([
            "            POLICY[\"📋 Azure Policy<br/>Governance\"]",
            "            MGMTGROUPS[\"🏢 Management Groups<br/>Hierarchy\"]"
        ])
    
    lines.extend([
        "        end",
        "",
        "        %% Define Hub-and-Spoke Connections",
        "        %% Cross-premises to Hub",
        "        ONPREM -.->|\"Private Connection\"| FIREWALL",
        "        INTERNET -->|\"Public Access\"| FIREWALL"
    ])
    
    # Add specific gateway connections
    if inputs.network_services:
        if "expressroute" in inputs.network_services:
            lines.append("        ER -.->|\"Private Peering\"| FIREWALL")
        if "vpn_gateway" in inputs.network_services:
            lines.append("        VPN -.->|\"Site-to-Site\"| FIREWALL")
    
    lines.extend([
        "",
        "        %% Hub to Spokes (Hub-and-Spoke Topology)",
        "        FIREWALL -->|\"Secure Routing\"| PRODVNET",
        "        FIREWALL -->|\"Secure Routing\"| DEVVNET",
        "",
        "        %% Shared Services Connections",
        "        DNS -.->|\"Name Resolution\"| PRODVNET",
        "        DNS -.->|\"Name Resolution\"| DEVVNET",
        "        BASTION -.->|\"Secure Access\"| PRODVNET",
        "        BASTION -.->|\"Secure Access\"| DEVVNET",
        "",
        "        %% Identity Integration",
        "        AAD -.->|\"Authentication\"| PRODVNET",
        "        AAD -.->|\"Authentication\"| DEVVNET",
        "        KEYVAULT -.->|\"Secrets\"| PRODVNET",
        "",
        "        %% Security Monitoring",
    ])
    
    # Add security monitoring connections if services exist
    if inputs.security_services:
        for service in inputs.security_services:
            if service in ["security_center", "sentinel", "defender"]:
                service_id = service.upper().replace("_", "")
                lines.extend([
                    f"        {service_id} -.->|\"Monitor\"| PRODVNET",
                    f"        {service_id} -.->|\"Monitor\"| DEVVNET"
                ])
    
    # Add monitoring connections if services exist
    if inputs.monitoring_services:
        for service in inputs.monitoring_services:
            service_id = service.upper().replace("_", "")
            lines.extend([
                f"        {service_id} -.->|\"Telemetry\"| PRODVNET",
                f"        {service_id} -.->|\"Telemetry\"| DEVVNET"
            ])
    
    # Enhanced service-to-service connectivity
    lines.append("        %% Enhanced Service-to-Service Connectivity")
    
    # Database connectivity patterns
    if inputs.database_services:
        for db_service in inputs.database_services:
            if db_service == "sql_database":
                lines.extend([
                    "        %% SQL Database Connections",
                    "        FIREWALL -->|\"Database Security\"| PROD_SQLDATABASE",
                    "        KEYVAULT -.->|\"Connection Strings\"| PROD_SQLDATABASE",
                    "        AAD -.->|\"Database Authentication\"| PROD_SQLDATABASE"
                ])
                # Connect to compute services if they exist
                if inputs.compute_services:
                    for compute in inputs.compute_services:
                        if compute == "virtual_machines":
                            lines.append("        PROD_VIRTUALMACHINES -->|\"Application Data\"| PROD_SQLDATABASE")
                        elif compute == "app_services":
                            lines.append("        PROD_APPSERVICES -->|\"Application Data\"| PROD_SQLDATABASE")
            elif db_service == "cosmos_db":
                lines.extend([
                    "        %% Cosmos DB Connections", 
                    "        FIREWALL -->|\"NoSQL Security\"| PROD_COSMOSDB",
                    "        AAD -.->|\"Cosmos Authentication\"| PROD_COSMOSDB"
                ])
    
    # Storage service connections
    if inputs.storage_services:
        lines.extend([
            "        %% Storage Connectivity",
            "        KEYVAULT -.->|\"Storage Keys\"| PROD_STORAGEACCOUNTS",
            "        AAD -.->|\"Storage Access Control\"| PROD_STORAGEACCOUNTS"
        ])
        # Connect to compute services
        if inputs.compute_services:
            for compute in inputs.compute_services:
                if compute == "virtual_machines":
                    lines.append("        PROD_VIRTUALMACHINES -->|\"Data Storage\"| PROD_STORAGEACCOUNTS")
                elif compute == "app_services":
                    lines.append("        PROD_APPSERVICES -->|\"App Data\"| PROD_STORAGEACCOUNTS")
    
    # Analytics service connections (data flow patterns)
    if inputs.analytics_services:
        lines.extend([
            "        %% Analytics & Data Flow",
            "        ANALYTICS[\"🧠 Analytics Services<br/>Data Processing\"]"
        ])
        
        # Connect storage to analytics for data flow
        if inputs.storage_services:
            lines.append("        PROD_STORAGEACCOUNTS -->|\"Data Pipeline\"| ANALYTICS")
        
        # Connect databases to analytics
        if inputs.database_services:
            if "sql_database" in inputs.database_services:
                lines.append("        PROD_SQLDATABASE -->|\"Data Export\"| ANALYTICS")
            if "cosmos_db" in inputs.database_services:
                lines.append("        PROD_COSMOSDB -->|\"Document Analytics\"| ANALYTICS")
    
    # DevOps service connections
    if inputs.devops_services:
        lines.extend([
            "        %% DevOps & CI/CD",
            "        DEVOPS[\"⚙️ DevOps Services<br/>CI/CD Pipeline\"]",
            "        AAD -.->|\"DevOps Authentication\"| DEVOPS",
            "        KEYVAULT -.->|\"Deployment Secrets\"| DEVOPS"
        ])
        
        # Connect DevOps to compute services
        if inputs.compute_services:
            for compute in inputs.compute_services:
                if compute == "virtual_machines":
                    lines.append("        DEVOPS -->|\"VM Deployment\"| PROD_VIRTUALMACHINES")
                elif compute == "app_services":
                    lines.append("        DEVOPS -->|\"App Deployment\"| PROD_APPSERVICES")
                elif compute == "aks":
                    lines.append("        DEVOPS -->|\"Container Deployment\"| PROD_AKS")
    
    # Integration service connections
    if inputs.integration_services:
        lines.extend([
            "        %% Integration Services",
            "        INTEGRATION[\"🔗 Integration Services<br/>API Management\"]",
            "        AAD -.->|\"API Authentication\"| INTEGRATION"
        ])
        
        # Connect integration to databases and storage
        if inputs.database_services:
            lines.append("        INTEGRATION -->|\"Data Integration\"| PROD_SQLDATABASE")
        if inputs.storage_services:
            lines.append("        INTEGRATION -->|\"Storage Integration\"| PROD_STORAGEACCOUNTS")
    
    
    lines.extend([
        "",
        "    end",
        "",
        "    %% Professional Styling",
        "    classDef hubStyle fill:#e3f2fd,stroke:#1565c0,stroke-width:3px,font-weight:bold;",
        "    classDef spokeStyle fill:#f3e5f5,stroke:#7b1fa2,stroke-width:2px;",
        "    classDef securityStyle fill:#ffebee,stroke:#c62828,stroke-width:2px;",
        "    classDef networkStyle fill:#e8f5e8,stroke:#2e7d32,stroke-width:2px;",
        "    classDef identityStyle fill:#fff3e0,stroke:#ef6c00,stroke-width:2px;",
        "    classDef workloadStyle fill:#f1f8e9,stroke:#558b2f,stroke-width:2px;",
        "    classDef crossPremStyle fill:#fafafa,stroke:#616161,stroke-width:2px;",
        "",
        "    %% Apply Styles",
        "    class FIREWALL,DNS,BASTION securityStyle;",
        "    class PRODVNET,DEVVNET networkStyle;",
        "    class AAD,KEYVAULT,POLICY,MGMTGROUPS identityStyle;",
        "    class ONPREM,INTERNET,ER,VPN crossPremStyle;"
    ])
    
    # Apply workload styles - include new services with correct naming
    workload_services = ["PROD_VIRTUALMACHINES", "DEV_COMPUTE", "PROD_APPSERVICES", "PROD_AKS", "PROD_SQLDATABASE", "PROD_COSMOSDB", "PROD_STORAGEACCOUNTS", "ANALYTICS", "DEVOPS", "INTEGRATION"]
    
    # Only include services that actually exist in the diagram
    existing_workload_services = []
    for service_id in workload_services:
        if any(service_id in line for line in lines):
            existing_workload_services.append(service_id)
    
    if existing_workload_services:
        lines.append(f"    class {','.join(existing_workload_services)} workloadStyle;")
    if inputs.compute_services or inputs.database_services:
        workload_ids = []
        if inputs.compute_services:
            for service in inputs.compute_services:
                workload_ids.append(f"PROD_{service.upper().replace('_', '')}")
        if inputs.database_services:
            for service in inputs.database_services:
                workload_ids.append(f"PROD_{service.upper().replace('_', '')}")
        if "DEV_COMPUTE" in "\n".join(lines):
            workload_ids.append("DEV_COMPUTE")
        
        if workload_ids:
            lines.append(f"    class {','.join(workload_ids)} workloadStyle;")
    
    # Apply security and monitoring styles
    if inputs.security_services:
        security_ids = []
        for service in inputs.security_services:
            if service in ["security_center", "sentinel", "defender"]:
                security_ids.append(service.upper().replace("_", ""))
        if security_ids:
            lines.append(f"    class {','.join(security_ids)} securityStyle;")
    
    if inputs.monitoring_services:
        monitoring_ids = []
        for service in inputs.monitoring_services:
            monitoring_ids.append(service.upper().replace("_", ""))
        if monitoring_ids:
            lines.append(f"    class {','.join(monitoring_ids)} networkStyle;")
    
    if inputs.network_services:
        gateway_ids = []
        for service in inputs.network_services:
            if service in ["application_gateway", "load_balancer", "vpn_gateway"]:
                gateway_ids.append(service.upper().replace("_", ""))
        if gateway_ids:
            lines.append(f"    class {','.join(gateway_ids)} networkStyle;")
    
    return "\n".join(lines)

def generate_professional_mermaid_with_orchestration(inputs: CustomerInputs, orchestration_result: Dict[str, Any] = None) -> str:
    """Generate professional Mermaid diagram with LangGraph orchestration for clear hub-spoke separation"""
    
    # Use orchestration result if available, otherwise fall back to traditional method
    if orchestration_result and orchestration_result.get("diagram_components"):
        return _generate_orchestrated_mermaid(inputs, orchestration_result)
    else:
        # Fallback to traditional generation
        return generate_professional_mermaid(inputs)

def _generate_orchestrated_mermaid(inputs: CustomerInputs, orchestration_result: Dict[str, Any]) -> str:
    """Generate Mermaid diagram using orchestration results with clear hub-spoke separation"""
    
    components = orchestration_result["diagram_components"]
    hub_data = components.get("hub", {})
    spokes_data = components.get("spokes", {})
    connectivity = components.get("connectivity", {})
    
    lines = [
        "graph TB",
        "    %% Azure Landing Zone - Hub and Spoke Architecture (LangGraph Orchestrated)",
        "    %% Generated by Azure Landing Zone Agent with LangGraph Multi-Agent Workflow",
        "",
        "    %% Styling for different components",
        "    classDef hubStyle fill:#0078d4,stroke:#005a9e,stroke-width:3px,color:#ffffff",
        "    classDef spokeStyle fill:#00bcf2,stroke:#0078d4,stroke-width:2px,color:#ffffff", 
        "    classDef sharedStyle fill:#40e0d0,stroke:#008b8b,stroke-width:2px,color:#000000",
        "    classDef workloadStyle fill:#98fb98,stroke:#006400,stroke-width:2px,color:#000000",
        "",
        "    %% Hub Infrastructure (Shared Services)",
        "    subgraph \"Hub\" [\"🏢 Hub Infrastructure - Shared Services\"]",
        "        direction TB"
    ]
    
    # Add hub services
    hub_services = hub_data.get("services", [])
    hub_network = hub_data.get("network", {})
    
    # Core hub components
    lines.extend([
        "        HUBVNET[\"🌐 Hub VNet<br/>Shared Infrastructure\"]",
        "        FIREWALL[\"🔥 Azure Firewall<br/>Central Security\"]",
        "        BASTION[\"🔐 Azure Bastion<br/>Secure Access\"]",
        "        DNS[\"📧 Private DNS<br/>Name Resolution\"]"
    ])
    
    # Add additional hub services based on orchestration
    hub_service_mapping = {
        'azure_ad': "AD[\"👥 Azure AD<br/>Identity Provider\"]",
        'security_center': "SC[\"🛡️ Security Center<br/>Threat Protection\"]",
        'log_analytics': "LA[\"📊 Log Analytics<br/>Monitoring Hub\"]",
        'vpn_gateway': "VPN[\"🌐 VPN Gateway<br/>Hybrid Connectivity\"]",
        'expressroute': "ER[\"⚡ ExpressRoute<br/>Private Connection\"]"
    }
    
    for service in hub_services:
        if service in hub_service_mapping:
            lines.append(f"        {hub_service_mapping[service]}")
    
    lines.extend([
        "    end",
        "",
        "    %% Production Spoke",
        "    subgraph \"ProdSpoke\" [\"🏭 Production Spoke\"]",
        "        direction TB",
        "        PRODVNET[\"🌐 Production VNet<br/>Workload Network\"]",
        ""
    ])
    
    # Add spoke services based on orchestration
    spoke_services = spokes_data.get("services", [])
    workload_components = spokes_data.get("workloads", {})
    
    # Production workload components
    if workload_components.get("production_spoke"):
        prod_services = workload_components["production_spoke"].get("services", {})
        
        lines.append("        subgraph \"ProdWorkloads\" [\"💼 Production Workloads\"]")
        
        # Web tier
        if prod_services.get("web_tier"):
            for service in prod_services["web_tier"]:
                if service == "app_services":
                    lines.append("            PRODAPP[\"🌐 App Services<br/>Production Web Apps\"]")
        
        # Application tier  
        if prod_services.get("application_tier"):
            for service in prod_services["application_tier"]:
                if service == "virtual_machines":
                    lines.append("            PRODVM[\"💻 Virtual Machines<br/>Production Compute\"]")
                elif service == "aks":
                    lines.append("            PRODAKS[\"☸️ AKS Cluster<br/>Production Containers\"]")
        
        # Data tier
        if prod_services.get("data_tier"):
            for service in prod_services["data_tier"]:
                if service == "sql_database":
                    lines.append("            PRODDB[\"🗄️ SQL Database<br/>Production Data\"]")
                elif service == "cosmos_db":
                    lines.append("            PRODCOSMOS[\"🌍 Cosmos DB<br/>Global Database\"]")
        
        lines.extend([
            "        end",
            "    end",
            ""
        ])
    
    # Development spoke
    lines.extend([
        "    %% Development Spoke", 
        "    subgraph \"DevSpoke\" [\"🔧 Development Spoke\"]",
        "        direction TB",
        "        DEVVNET[\"🌐 Development VNet<br/>Dev/Test Network\"]",
        "",
        "        subgraph \"DevWorkloads\" [\"🧪 Development Workloads\"]"
    ])
    
    # Add dev services (simplified versions of prod)
    if workload_components.get("development_spoke"):
        dev_services = workload_components["development_spoke"].get("services", {})
        
        # Simplified dev components
        if dev_services.get("web_tier"):
            lines.append("            DEVAPP[\"🌐 App Services<br/>Development Apps\"]")
        if dev_services.get("application_tier"):
            lines.append("            DEVVM[\"💻 Virtual Machines<br/>Dev Compute\"]")
        if dev_services.get("data_tier"):
            lines.append("            DEVDB[\"🗄️ SQL Database<br/>Development Data\"]")
    
    lines.extend([
        "        end",
        "    end",
        "",
        "    %% External Connectivity",
        "    ONPREM[\"🏢 On-Premises<br/>Corporate Network\"]",
        "    INTERNET[\"🌍 Internet<br/>Public Access\"]",
        "",
        "    %% Hub-and-Spoke Connections (LangGraph Orchestrated)"
    ])
    
    # Add connectivity based on orchestration results
    traffic_flow = connectivity.get("traffic_flow", {})
    
    # Core hub-spoke connections
    lines.extend([
        "    %% Hub to Spokes (VNet Peering)",
        "    HUBVNET -.->|\"VNet Peering<br/>Private\"| PRODVNET",
        "    HUBVNET -.->|\"VNet Peering<br/>Private\"| DEVVNET",
        "",
        "    %% Shared Services to Spokes",
        "    FIREWALL -->|\"Secure Routing\"| PRODVNET",
        "    FIREWALL -->|\"Secure Routing\"| DEVVNET",
        "    DNS -.->|\"Name Resolution\"| PRODVNET", 
        "    DNS -.->|\"Name Resolution\"| DEVVNET",
        "    BASTION -.->|\"Secure Admin Access\"| PRODVNET",
        "    BASTION -.->|\"Secure Admin Access\"| DEVVNET"
    ])
    
    # External connectivity based on orchestration
    if "vpn_gateway" in hub_services:
        lines.append("    ONPREM -.->|\"Site-to-Site VPN\"| VPN")
        lines.append("    VPN -->|\"Secure Tunnel\"| HUBVNET")
    elif "expressroute" in hub_services:
        lines.append("    ONPREM -.->|\"Private Peering\"| ER")
        lines.append("    ER -->|\"Dedicated Connection\"| HUBVNET")
    else:
        lines.append("    ONPREM -.->|\"Hybrid Connection\"| FIREWALL")
    
    lines.extend([
        "    INTERNET -->|\"Public Traffic\"| FIREWALL",
        "",
        "    %% Enhanced Service-to-Service Connections"
    ])
    
    # Add enhanced service-to-service connectivity based on services present
    spoke_services = spokes_data.get("services", [])
    
    # Identity and security service connections
    if "azure_ad" in hub_services:
        lines.extend([
            "    AD -.->|\"Authentication\"| PRODVNET",
            "    AD -.->|\"Authentication\"| DEVVNET"
        ])
    
    # Key Vault connections (if present in security services)
    if inputs.security_services and "key_vault" in inputs.security_services:
        lines.extend([
            "    %% Key Vault Integration",
            "    KEYVAULT[\"🔐 Azure Key Vault<br/>Secrets Management\"]",
            "    KEYVAULT -.->|\"Secrets & Certificates\"| PRODVNET",
            "    KEYVAULT -.->|\"Secrets & Certificates\"| DEVVNET",
            "    KEYVAULT -.->|\"Firewall Certificates\"| FIREWALL"
        ])
    
    # Database and storage connectivity patterns
    if inputs.database_services:
        for db_service in inputs.database_services:
            if db_service == "sql_database":
                lines.extend([
                    "    %% Database Connectivity",
                    "    FIREWALL -->|\"Database Security\"| PRODDB",
                    "    PRODAPP -->|\"Application Data\"| PRODDB",
                    "    PRODVM -->|\"Database Access\"| PRODDB"
                ])
            elif db_service == "cosmos_db":
                lines.extend([
                    "    %% Cosmos DB Connectivity", 
                    "    FIREWALL -->|\"NoSQL Security\"| PRODCOSMOS",
                    "    PRODAPP -->|\"Document Data\"| PRODCOSMOS"
                ])
    
    # Storage service connections
    if inputs.storage_services:
        lines.extend([
            "    %% Storage Connectivity",
            "    PRODSTORAGE -->|\"Data Persistence\"| PRODAPP",
            "    PRODSTORAGE -->|\"VM Storage\"| PRODVM",
            "    FIREWALL -->|\"Storage Security\"| PRODSTORAGE"
        ])
    
    # Monitoring connections to all services
    if inputs.monitoring_services and "monitor" in inputs.monitoring_services:
        lines.extend([
            "    %% Monitoring & Observability",
            "    MONITOR[\"📊 Azure Monitor<br/>Observability\"]",
            "    MONITOR -.->|\"Telemetry\"| PRODVNET",
            "    MONITOR -.->|\"Telemetry\"| DEVVNET",
            "    MONITOR -.->|\"Hub Monitoring\"| HUBVNET"
        ])
        
        # Monitor specific services if they exist
        service_monitors = {
            "PRODAPP": "App Insights",
            "PRODVM": "VM Metrics", 
            "PRODDB": "DB Performance",
            "FIREWALL": "Security Logs"
        }
        
        for service_id, monitor_type in service_monitors.items():
            # Check if service exists in the diagram by looking for it in lines
            if any(service_id in line for line in lines):
                lines.append(f"    MONITOR -.->|\"{monitor_type}\"| {service_id}")
    
    # Analytics and AI service connections (data flow patterns)
    if inputs.analytics_services:
        lines.extend([
            "    %% Analytics & Data Flow",
            "    ANALYTICS[\"🧠 Analytics Services<br/>Data Processing\"]"
        ])
        
        # Connect storage to analytics for data flow
        if inputs.storage_services:
            lines.append("    PRODSTORAGE -->|\"Data Pipeline\"| ANALYTICS")
        
        # Connect databases to analytics
        if inputs.database_services:
            lines.append("    PRODDB -->|\"Data Export\"| ANALYTICS")
    
    # DevOps service connections
    if inputs.devops_services:
        lines.extend([
            "    %% DevOps & CI/CD",
            "    DEVOPS[\"⚙️ DevOps Services<br/>CI/CD Pipeline\"]",
            "    DEVOPS -->|\"Deployment\"| PRODAPP",
            "    DEVOPS -->|\"Deployment\"| DEVAPP",
            "    DEVOPS -.->|\"Infrastructure Updates\"| PRODVM"
        ])
    
    lines.extend([
        "",
        "    %% Apply Styling",
        "    class HUBVNET,FIREWALL,BASTION,DNS hubStyle",
        "    class PRODVNET,DEVVNET spokeStyle"
    ])
    
    # Add styling for hub services including new services
    hub_service_ids = ["AD", "SC", "LA", "VPN", "ER", "KEYVAULT", "MONITOR"]
    existing_hub_ids = [id for id in hub_service_ids if any(id in line for line in lines)]
    if existing_hub_ids:
        lines.append(f"    class {','.join(existing_hub_ids)} sharedStyle")
    
    # Add styling for workload and analytics components
    workload_ids = []
    analytics_ids = ["ANALYTICS", "DEVOPS"]
    
    for line in lines:
        # Extract workload IDs
        if any(wl in line for wl in ["PRODAPP", "PRODVM", "PRODAKS", "PRODDB", "PRODCOSMOS", "PRODSTORAGE", "DEVAPP", "DEVVM", "DEVDB"]):
            if "[" in line:
                id_part = line.split("[")[0].strip()
                if id_part and id_part not in workload_ids:
                    workload_ids.append(id_part)
    
    if workload_ids:
        lines.append(f"    class {','.join(workload_ids)} workloadStyle")
    
    # Add styling for analytics services
    existing_analytics_ids = [id for id in analytics_ids if any(id in line for line in lines)]
    if existing_analytics_ids:
        lines.append(f"    class {','.join(existing_analytics_ids)} sharedStyle")
    
    lines.extend([
        "",
        "    %% Metadata",
        f"    %%{{ Generated: {datetime.now().isoformat()} }}%%",
        "    %%{{ Architecture: Hub-and-Spoke with LangGraph Orchestration }}%%",
        f"    %%{{ Hub Services: {len(hub_services)} }}%%",
        f"    %%{{ Spoke Services: {len(spoke_services)} }}%%"
    ])
    
    return "\n".join(lines)


def generate_enhanced_drawio_xml(inputs: CustomerInputs) -> str:
    """Generate enhanced Draw.io XML with comprehensive Azure stencils based on user selections"""
    
    def esc(s): 
        return html.escape(s) if s else ""
    
    template = generate_architecture_template(inputs)
    diagram_id = str(uuid.uuid4())
    
    # Base layout coordinates
    y_start = 100
    current_y = y_start
    section_height = 350
    service_width = 100
    service_height = 80
    
    # Build dynamic XML content
    xml_parts = [
        f"""<mxfile host="app.diagrams.net" modified="2024-01-01T00:00:00.000Z" agent="Azure Landing Zone Agent" version="1.0.0">
  <diagram name="Azure Landing Zone Architecture" id="azure-lz-{diagram_id}">
    <mxGraphModel dx="1422" dy="794" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="2400" pageHeight="1600" math="0" shadow="0">
      <root>
        <mxCell id="0" />
        <mxCell id="1" parent="0" />
        
        <!-- Azure Tenant Container -->
        <mxCell id="tenant" value="Azure Tenant - {esc(inputs.org_structure or 'Enterprise')}" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#dae8fc;strokeColor=#6c8ebf;fontSize=16;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="50" y="50" width="2300" height="1500" as="geometry" />
        </mxCell>
        
        <!-- Management Groups -->
        <mxCell id="mgmt-groups" value="Management Groups" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e1d5e7;strokeColor=#9673a6;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="500" height="250" as="geometry" />
        </mxCell>"""
    ]
    
    # Management Group structure based on template
    mg_x = 150
    mg_y = current_y + 50
    xml_parts.append(f"""
        <mxCell id="root-mg" value="Root MG" style="shape=mxgraph.azure.management;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{mg_x}" y="{mg_y}" width="80" height="60" as="geometry" />
        </mxCell>""")
    
    mg_x += 120
    for i, mg in enumerate(template['template']['management_groups'][1:3]):  # Platform and Workloads
        mg_id = mg.lower().replace(' ', '-')
        xml_parts.append(f"""
        <mxCell id="{mg_id}-mg" value="{mg}" style="shape=mxgraph.azure.management;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{mg_x}" y="{mg_y}" width="80" height="60" as="geometry" />
        </mxCell>""")
        mg_x += 120
    
    # Subscriptions
    current_y += 300
    xml_parts.append(f"""
        <!-- Subscriptions -->
        <mxCell id="subscriptions" value="Subscriptions" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#d5e8d4;strokeColor=#82b366;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="700" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
    
    sub_x = 750
    sub_y = current_y + 50
    for sub in template['template']['subscriptions'][:4]:  # First 4 subscriptions
        xml_parts.append(f"""
        <mxCell id="{sub.lower().replace(' ', '-')}-sub" value="{sub}" style="shape=mxgraph.azure.subscription;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sub_x}" y="{sub_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
        sub_x += 130
        if sub_x > 1200:  # Wrap to next row
            sub_x = 750
            sub_y += 100
    
    # Network Architecture Section
    current_y += 300
    xml_parts.append(f"""
        <!-- Network Architecture -->
        <mxCell id="network" value="Network Architecture - {esc(inputs.network_model or 'Hub-Spoke')}" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff2cc;strokeColor=#d6b656;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="800" height="{section_height}" as="geometry" />
        </mxCell>""")
    
    # Hub VNet (always present)
    hub_x = 200
    hub_y = current_y + 80
    xml_parts.append(f"""
        <mxCell id="hub-vnet" value="Hub VNet\\nShared Services" style="shape=mxgraph.azure.virtual_network;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{hub_x}" y="{hub_y}" width="120" height="80" as="geometry" />
        </mxCell>""")
    
    # Spoke VNets
    spoke_x = 400
    spoke_y = hub_y - 50
    xml_parts.append(f"""
        <mxCell id="spoke1-vnet" value="Production VNet" style="shape=mxgraph.azure.virtual_network;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{spoke_x}" y="{spoke_y}" width="120" height="80" as="geometry" />
        </mxCell>
        <mxCell id="spoke2-vnet" value="Development VNet" style="shape=mxgraph.azure.virtual_network;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{spoke_x}" y="{spoke_y + 120}" width="120" height="80" as="geometry" />
        </mxCell>""")
    
    # Add selected network services
    if inputs.network_services:
        net_x = 600
        net_y = hub_y
        for i, service in enumerate(inputs.network_services[:4]):  # Max 4 network services
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="net-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{net_x}" y="{net_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                net_y += 100
    
    # Compute Services Section
    if inputs.compute_services or inputs.workload:
        current_y += section_height + 50
        xml_parts.append(f"""
        <!-- Compute Services -->
        <mxCell id="compute" value="Compute Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff0e6;strokeColor=#d79b00;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1000" y="{current_y}" width="600" height="{section_height}" as="geometry" />
        </mxCell>""")
        
        comp_x = 1050
        comp_y = current_y + 50
        
        # Add selected compute services
        services_to_add = inputs.compute_services or []
        if inputs.workload and inputs.workload not in services_to_add:
            services_to_add.append(inputs.workload)
            
        for i, service in enumerate(services_to_add[:6]):  # Max 6 compute services
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="compute-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{comp_x}" y="{comp_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                comp_x += 120
                if comp_x > 1450:  # Wrap to next row
                    comp_x = 1050
                    comp_y += 100
    
    # Storage Services Section
    if inputs.storage_services:
        current_y += section_height + 50
        xml_parts.append(f"""
        <!-- Storage Services -->
        <mxCell id="storage" value="Storage Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f0f0f0;strokeColor=#666666;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        
        stor_x = 150
        stor_y = current_y + 50
        for i, service in enumerate(inputs.storage_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'storage_accounts')
                xml_parts.append(f"""
        <mxCell id="storage-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{stor_x}" y="{stor_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                stor_x += 120
                if stor_x > 550:
                    stor_x = 150
                    stor_y += 100
    
    # Database Services Section
    if inputs.database_services:
        db_y = current_y if not inputs.storage_services else current_y
        if inputs.storage_services:
            xml_parts.append(f"""
        <!-- Database Services -->
        <mxCell id="database" value="Database Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e6f3ff;strokeColor=#0066cc;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="800" y="{db_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        else:
            current_y += section_height + 50
            xml_parts.append(f"""
        <!-- Database Services -->
        <mxCell id="database" value="Database Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#e6f3ff;strokeColor=#0066cc;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{current_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
            db_y = current_y
        
        db_x = 850 if inputs.storage_services else 150
        db_y += 50
        for i, service in enumerate(inputs.database_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'sql_database')
                xml_parts.append(f"""
        <mxCell id="database-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{db_x}" y="{db_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                db_x += 120
                if db_x > (1250 if inputs.storage_services else 550):
                    db_x = 850 if inputs.storage_services else 150
                    db_y += 100
    
    # Security Services Section (always present)
    current_y += 300
    xml_parts.append(f"""
        <!-- Security & Identity Services -->
        <mxCell id="security" value="Security &amp; Identity Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#ffebee;strokeColor=#c62828;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1700" y="{y_start}" width="600" height="600" as="geometry" />
        </mxCell>""")
    
    # Core security services (always present)
    sec_x = 1750
    sec_y = y_start + 50
    core_security = [
        ('azure-ad', 'Azure AD', 'azure_active_directory'),
        ('key-vault', 'Key Vault', 'key_vault'),
        ('security-center', 'Security Center', 'security_center')
    ]
    
    for sec_id, sec_name, sec_shape in core_security:
        xml_parts.append(f"""
        <mxCell id="{sec_id}" value="{sec_name}" style="shape=mxgraph.azure.{sec_shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sec_x}" y="{sec_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
        sec_x += 120
        if sec_x > 2100:
            sec_x = 1750
            sec_y += 100
    
    # Add additional selected security services
    if inputs.security_services:
        for i, service in enumerate(inputs.security_services):
            if service in AZURE_SERVICES_MAPPING and service not in ['active_directory', 'key_vault', 'security_center']:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="security-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{sec_x}" y="{sec_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                sec_x += 120
                if sec_x > 2100:
                    sec_x = 1750
                    sec_y += 100
    
    # Analytics Services Section
    if inputs.analytics_services:
        analytics_y = y_start + 650
        xml_parts.append(f"""
        <!-- Analytics & AI Services -->
        <mxCell id="analytics" value="Analytics &amp; AI Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f3e5f5;strokeColor=#9c27b0;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="1700" y="{analytics_y}" width="600" height="300" as="geometry" />
        </mxCell>""")
        
        ana_x = 1750
        ana_y = analytics_y + 50
        for i, service in enumerate(inputs.analytics_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="analytics-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{ana_x}" y="{ana_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                ana_x += 120
                if ana_x > 2100:
                    ana_x = 1750
                    ana_y += 100
    
    # Integration Services Section
    if inputs.integration_services:
        int_y = current_y
        xml_parts.append(f"""
        <!-- Integration Services -->
        <mxCell id="integration" value="Integration Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#fff8e1;strokeColor=#ff8f00;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="100" y="{int_y}" width="600" height="250" as="geometry" />
        </mxCell>""")
        
        int_x = 150
        int_y += 50
        for i, service in enumerate(inputs.integration_services[:4]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="integration-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{int_x}" y="{int_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                int_x += 120
                if int_x > 550:
                    int_x = 150
                    int_y += 100
    
    # DevOps Services Section
    if inputs.devops_services:
        devops_y = current_y if not inputs.integration_services else current_y
        devops_x_offset = 800 if inputs.integration_services else 100
        xml_parts.append(f"""
        <!-- DevOps Services -->
        <mxCell id="devops" value="DevOps Services" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#f5f5f5;strokeColor=#666666;fontSize=14;fontStyle=1;verticalAlign=top;spacingTop=10;" vertex="1" parent="1">
          <mxGeometry x="{devops_x_offset}" y="{devops_y}" width="400" height="250" as="geometry" />
        </mxCell>""")
        
        dev_x = devops_x_offset + 50
        dev_y = devops_y + 50
        for i, service in enumerate(inputs.devops_services[:3]):
            if service in AZURE_SERVICES_MAPPING:
                service_info = AZURE_SERVICES_MAPPING[service]
                shape = service_info.get('drawio_shape', 'generic_service')
                xml_parts.append(f"""
        <mxCell id="devops-service-{i}" value="{esc(service_info['name'])}" style="shape=mxgraph.azure.{shape};fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="1">
          <mxGeometry x="{dev_x}" y="{dev_y}" width="{service_width}" height="{service_height}" as="geometry" />
        </mxCell>""")
                dev_x += 120
                if dev_x > (devops_x_offset + 250):
                    dev_x = devops_x_offset + 50
                    dev_y += 100
    
    # Add basic connections
    xml_parts.append("""
        <!-- Key Connections -->
        <mxCell id="conn1" edge="1" source="root-mg" target="platform-mg" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn2" edge="1" source="root-mg" target="landing-zones-mg" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn3" edge="1" source="hub-vnet" target="spoke1-vnet" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        <mxCell id="conn4" edge="1" source="hub-vnet" target="spoke2-vnet" parent="1">
          <mxGeometry relative="1" as="geometry" />
        </mxCell>
        
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>""")
    
    return "".join(xml_parts)


def generate_professional_documentation(inputs: CustomerInputs) -> Dict[str, str]:
    """Generate professional TSD, HLD, and LLD documentation with AI enhancement"""
    
    template = generate_architecture_template(inputs)
    timestamp = datetime.now().strftime("%Y-%m-%d")
    
    # Generate AI insights if additional inputs are provided
    url_analysis = ""
    doc_analysis = ""
    ai_recommendations = ""
    
    try:
        if inputs.url_input:
            url_analysis = analyze_url_content(inputs.url_input)
            
        if inputs.uploaded_files_info:
            doc_analysis = "Document analysis results incorporated from uploaded files."
            
        # Generate AI-enhanced recommendations
        ai_recommendations = generate_ai_enhanced_recommendations(inputs, url_analysis, doc_analysis)
    except Exception as e:
        logger.warning(f"AI enhancement failed: {e}")
        ai_recommendations = "AI enhancement not available - using standard recommendations."
    
    # Technical Specification Document (TSD)
    tsd = f"""# Technical Specification Document (TSD)
## Azure Landing Zone Architecture - Enterprise Edition

**Document Version:** 2.0 (AI-Enhanced)
**Date:** {timestamp}
**Business Objective:** {inputs.business_objective or 'Not specified'}

### Executive Summary
This document outlines the technical specifications for implementing an Azure Landing Zone architecture based on comprehensive customer requirements analysis, including AI-powered insights and recommendations.

### Business Requirements Analysis
- **Primary Objective:** {inputs.business_objective or 'Cost optimization and operational efficiency'}
- **Industry:** {inputs.industry or 'General'}
- **Regulatory Requirements:** {inputs.regulatory or 'Standard compliance'}
- **Organization Structure:** {inputs.org_structure or 'Enterprise'}
- **Governance Model:** {inputs.governance or 'Centralized with delegated permissions'}

### Architecture Template Selection
**Selected Template:** {template['template']['name']}
**Justification:** Based on organizational size, complexity, and regulatory requirements.

### Core Architecture Components
- **Identity & Access Management:** {inputs.identity or 'Azure Active Directory with hybrid integration'}
- **Network Architecture:** {inputs.network_model or 'Hub-Spoke with Azure Virtual WAN'}
- **Security Framework:** {inputs.security_posture or 'Zero Trust with defense in depth'}
- **Connectivity Strategy:** {inputs.connectivity or 'Hybrid cloud with ExpressRoute'}
- **Primary Workloads:** {inputs.workload or 'Multi-tier applications with microservices'}
- **Monitoring & Observability:** {inputs.monitoring or 'Azure Monitor with Log Analytics'}

### Enhanced Requirements Analysis
{f"**Additional Context:** {inputs.free_text_input}" if inputs.free_text_input else "**Additional Context:** Standard requirements captured through structured inputs."}

{f"**URL Analysis Insights:** {url_analysis[:500]}..." if url_analysis else ""}

{f"**Document Analysis:** Document analysis completed on uploaded files." if inputs.uploaded_files_info else ""}

### AI-Powered Architecture Recommendations
{ai_recommendations[:2000] if ai_recommendations else "Standard architecture recommendations applied."}

### Compliance & Governance Framework
- **Governance Model:** {inputs.governance or 'Centralized with delegated permissions'}
- **Policy Framework:** Azure Policy for compliance enforcement
- **Security Framework:** {inputs.security_posture or 'Zero Trust'} security model
"""

    # High Level Design (HLD)
    hld = f"""# High Level Design (HLD)
## Azure Landing Zone Implementation

**Document Version:** 1.0
**Date:** {timestamp}

### Architecture Overview
The proposed Azure Landing Zone follows the {template['template']['name']} pattern.

### Management Group Structure
"""
    
    for mg in template['template']['management_groups']:
        hld += f"- **{mg}:** Management group for {mg.lower()} resources\n"
    
    hld += f"""
### Subscription Strategy
"""
    
    for sub in template['template']['subscriptions']:
        hld += f"- **{sub}:** Dedicated subscription for {sub.lower()} workloads\n"
    
    hld += f"""
### Network Architecture
**Topology:** {inputs.network_model or 'Hub-Spoke'}
- **Hub VNet:** Central connectivity and shared services
- **Spoke VNets:** Workload-specific virtual networks
- **Connectivity:** {inputs.connectivity or 'ExpressRoute and VPN Gateway'}

### Security Architecture
**Security Model:** {inputs.security_posture or 'Zero Trust'}
- **Identity:** {inputs.identity or 'Azure Active Directory'} with conditional access
- **Network Security:** Network Security Groups and Azure Firewall
- **Data Protection:** {inputs.key_vault or 'Azure Key Vault'} for secrets management
- **Threat Protection:** {inputs.threat_protection or 'Azure Security Center and Sentinel'}

### Workload Architecture
**Primary Workload:** {inputs.workload or 'Application Services'}
- **Compute:** {AZURE_SERVICES_MAPPING.get(inputs.workload or 'appservices', {'name': 'Azure App Services'})['name']}
- **Architecture Style:** {inputs.architecture_style or 'Microservices'}
- **Scalability:** {inputs.scalability or 'Auto-scaling enabled'}
"""

    # Low Level Design (LLD)
    lld = f"""# Low Level Design (LLD)
## Azure Landing Zone Technical Implementation

**Document Version:** 1.0
**Date:** {timestamp}

### Resource Configuration

#### Management Groups
"""
    
    for i, mg in enumerate(template['template']['management_groups']):
        lld += f"""
**{mg} Management Group:**
- Management Group ID: mg-{mg.lower().replace(' ', '-')}
- Parent: {template['template']['management_groups'][i-1] if i > 0 else 'Tenant Root'}
- Applied Policies: Azure Policy assignments for {mg.lower()}
"""

    lld += f"""
#### Subscriptions
"""
    
    for sub in template['template']['subscriptions']:
        lld += f"""
**{sub} Subscription:**
- Subscription Name: sub-{sub.lower().replace(' ', '-')}
- Resource Groups: Multiple RGs based on workload segregation
- RBAC: Custom roles and assignments
- Budget: Cost management and alerting configured
"""

    lld += f"""
#### Network Configuration

**Hub Virtual Network:**
- VNet Name: vnet-hub-{inputs.network_model or 'spoke'}-001
- Address Space: 10.0.0.0/16
- Subnets:
  - GatewaySubnet: 10.0.0.0/24 (VPN/ExpressRoute Gateway)
  - AzureFirewallSubnet: 10.0.1.0/24 (Azure Firewall)
  - SharedServicesSubnet: 10.0.2.0/24 (Domain Controllers, etc.)

**Spoke Virtual Networks:**
- Production Spoke: vnet-prod-{inputs.workload or 'app'}-001 (10.1.0.0/16)
- Development Spoke: vnet-dev-{inputs.workload or 'app'}-001 (10.2.0.0/16)

#### Security Configuration

**Azure Active Directory:**
- Tenant: {inputs.org_structure or 'enterprise'}.onmicrosoft.com
- Custom Domains: Configured as required
- Conditional Access: {inputs.security_posture or 'Zero Trust'} policies
- PIM: Privileged Identity Management for admin roles

**Network Security:**
- Azure Firewall: Central security appliance
- NSGs: Network Security Groups on all subnets
- UDRs: User Defined Routes for traffic steering

**Key Management:**
- Key Vault: {inputs.key_vault or 'Azure Key Vault'} for certificates and secrets
- Managed Identities: For secure service-to-service authentication

#### Workload Configuration

**Primary Workload: {inputs.workload or 'Application Services'}**
- Service: {AZURE_SERVICES_MAPPING.get(inputs.workload or 'appservices', {'name': 'Azure App Services'})['name']}
- SKU: Production-grade tier
- Scaling: {inputs.scalability or 'Auto-scaling based on CPU/memory'}
- Monitoring: {inputs.monitoring or 'Azure Monitor'} with custom dashboards

#### Operations Configuration

**Monitoring and Alerting:**
- Log Analytics Workspace: Central logging for all resources
- Azure Monitor: Metrics and alerting
- Application Insights: Application performance monitoring

**Backup and Recovery:**
- Azure Backup: {inputs.backup or 'Daily backups with 30-day retention'}
- Site Recovery: Disaster recovery as needed

**Cost Management:**
- Budget Alerts: Monthly budget monitoring
- Cost Optimization: {inputs.cost_priority or 'Regular cost reviews and optimization'}

#### Infrastructure as Code

**IaC Tool:** {inputs.iac or 'Bicep/ARM Templates'}
- Template Structure: Modular templates for each component
- Deployment: CI/CD pipeline using Azure DevOps
- Version Control: Git repository with proper branching strategy
"""

    return {
        "tsd": tsd,
        "hld": hld, 
        "lld": lld
    }


# ---------- API Endpoints ----------

@app.get("/")
def root():
    return {
        "message": "Azure Landing Zone Agent API",
        "version": "1.0.0",
        "features": [
            "Architecture validation with compliance scoring",
            "Intelligent diagram generation with validation feedback", 
            "Comprehensive Azure Landing Zone templates",
            "Hub-spoke topology optimization",
            "Security best practices validation"
        ],
        "endpoints": [
            "/docs - API Documentation",
            "/validate-architecture - Validate architecture against Azure Landing Zone rules",
            "/validate-and-generate-diagram - Validate and generate diagrams with feedback",
            "/generate-diagram - Generate architecture diagram (Mermaid + Draw.io)",
            "/generate-azure-diagram - Generate Azure architecture diagram with official Azure icons (Python Diagrams)",
            "/generate-hub-spoke-vm-firewall - Enhanced hub-spoke diagram with VM in spoke and Firewall in hub (with dashed lines)",
            "/generate-drawio - Generate Draw.io XML",
            "/generate-comprehensive-azure-architecture - Full architecture generation",
            "/generate-intelligent-diagram - AI-powered diagram generation",
            "/health - Health check"
        ],
        "enhanced_features": [
            "Hub-spoke architecture with visual distinction using dashed boxes",
            "VM placement in spoke VNets with enhanced borders",
            "Network services (Firewall) placement in hub with solid connections",
            "Dashed line connections showing hub-spoke topology",
            "Azure stencils using mingrammer Python library",
            "Enhanced enterprise resource visualization"
        ]
    }

@app.get("/health")
def health_check():
    """Enhanced health check that verifies system dependencies"""
    logger.info("Running health check...")
    status = "healthy"
    issues = []
    
    # Check Graphviz availability
    try:
        result = subprocess.run(['dot', '-V'], capture_output=True, text=True, timeout=10)
        if result.returncode != 0:
            issues.append(f"Graphviz 'dot' command failed with return code {result.returncode}")
            status = "degraded"
        else:
            logger.info(f"Graphviz check passed: {result.stderr.strip()}")
    except subprocess.TimeoutExpired:
        issues.append("Graphviz 'dot' command timed out")
        status = "degraded"
    except FileNotFoundError:
        issues.append("Graphviz not installed or not accessible")
        status = "degraded"
    except Exception as e:
        issues.append(f"Graphviz check failed: {str(e)}")
        status = "degraded"
    
    # Check diagrams library
    try:
        from diagrams import Diagram
        logger.info("Diagrams library import successful")
    except ImportError as e:
        issues.append(f"Diagrams library import failed: {str(e)}")
        status = "unhealthy"
    
    # Check output directory accessibility
    try:
        output_dir = get_safe_output_directory()
        logger.info(f"Output directory accessible: {output_dir}")
    except Exception as e:
        issues.append(f"Cannot access output directory: {str(e)}")
        status = "degraded"
    
    # Check available disk space
    try:
        import shutil
        output_dir = get_safe_output_directory()
        total, used, free = shutil.disk_usage(output_dir)
        free_mb = free // (1024*1024)
        if free_mb < 100:  # Less than 100MB free
            issues.append(f"Low disk space in output directory: {free_mb}MB free")
            status = "degraded"
        logger.info(f"Disk space check passed: {free_mb}MB free")
    except Exception as e:
        issues.append(f"Cannot check disk space: {str(e)}")
        status = "degraded"
    
    # Test a simple diagram generation
    try:
        from main import CustomerInputs
        test_inputs = CustomerInputs(business_objective="Health check test")
        # Just validate inputs, don't generate full diagram
        validate_customer_inputs(test_inputs)
        logger.info("Input validation test passed")
    except Exception as e:
        issues.append(f"Input validation test failed: {str(e)}")
        status = "degraded"
    
    logger.info(f"Health check completed with status: {status}")
    
    return {
        "status": status,
        "timestamp": datetime.now().isoformat(),
        "issues": issues,
        "dependencies": {
            "graphviz_available": "Graphviz" not in str(issues),
            "diagrams_available": "Diagrams library" not in str(issues),
            "output_directory_accessible": "output directory" not in str(issues),
            "sufficient_disk_space": "disk space" not in str(issues)
        }
    }

@app.post("/generate-diagram")
def generate_diagram(inputs: CustomerInputs):
    """Generate comprehensive Azure Landing Zone diagrams and documentation"""
    try:
        # Generate professional diagrams
        mermaid_diagram = generate_professional_mermaid(inputs)
        drawio_xml = generate_enhanced_drawio_xml(inputs)
        
        # Generate professional documentation
        docs = generate_professional_documentation(inputs)
        
        return {
            "success": True,
            "mermaid": mermaid_diagram,
            "drawio": drawio_xml,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent"
            }
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating diagram: {str(e)}")

@app.post("/generate-hub-spoke-diagram")
def generate_hub_spoke_diagram(inputs: CustomerInputs):
    """Generate Azure Landing Zone diagrams with LangGraph hub-spoke orchestration"""
    try:
        # Process through LangGraph orchestrator
        orchestration_result = None
        if langgraph_orchestrator:
            try:
                # Convert inputs to dict for orchestrator
                inputs_dict = inputs.dict()
                orchestration_result = langgraph_orchestrator.process_landing_zone_request(inputs_dict)
                logger.info("LangGraph orchestration completed successfully")
            except Exception as orch_error:
                logger.warning(f"LangGraph orchestration failed, falling back to traditional method: {orch_error}")
                orchestration_result = None
        
        # Generate orchestrated diagrams
        if orchestration_result and orchestration_result.get("success"):
            mermaid_diagram = generate_professional_mermaid_with_orchestration(inputs, orchestration_result)
            drawio_xml = generate_enhanced_drawio_xml_with_orchestration(inputs, orchestration_result)
        else:
            # Fallback to traditional generation
            mermaid_diagram = generate_professional_mermaid(inputs)
            drawio_xml = generate_enhanced_drawio_xml(inputs)
        
        # Generate professional documentation
        docs = generate_professional_documentation(inputs)
        
        # Prepare response
        response = {
            "success": True,
            "mermaid": mermaid_diagram,
            "drawio": drawio_xml,
            "tsd": docs["tsd"],
            "hld": docs["hld"], 
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent with LangGraph Orchestration",
                "orchestration_enabled": orchestration_result is not None,
                "hub_spoke_separation": True
            }
        }
        
        # Add orchestration details if available
        if orchestration_result:
            response["orchestration_details"] = {
                "hub_services_count": orchestration_result.get("final_result", {}).get("hub_services_count", 0),
                "spoke_services_count": orchestration_result.get("final_result", {}).get("spoke_services_count", 0),
                "architecture_pattern": orchestration_result.get("final_result", {}).get("architecture_pattern", "hub_and_spoke"),
                "execution_log": orchestration_result.get("execution_log", [])
            }
        
        return response
    
    except Exception as e:
        logger.error(f"Error in hub-spoke diagram generation: {e}")
        raise HTTPException(status_code=500, detail=f"Error generating hub-spoke diagram: {str(e)}")

def generate_enhanced_drawio_xml_with_orchestration(inputs: CustomerInputs, orchestration_result: Dict[str, Any] = None) -> str:
    """Generate enhanced Draw.io XML with LangGraph orchestration for clear hub-spoke separation"""
    
    # Use orchestration result if available, otherwise fall back to traditional method
    if orchestration_result and orchestration_result.get("diagram_components"):
        return _generate_orchestrated_drawio(inputs, orchestration_result)
    else:
        # Fallback to traditional generation
        return generate_enhanced_drawio_xml(inputs)

def _generate_orchestrated_drawio(inputs: CustomerInputs, orchestration_result: Dict[str, Any]) -> str:
    """Generate Draw.io XML using orchestration results with clear hub-spoke separation"""
    
    components = orchestration_result["diagram_components"]
    hub_data = components.get("hub", {})
    spokes_data = components.get("spokes", {})
    
    xml_parts = [
        '<?xml version="1.0" encoding="UTF-8"?>',
        '<mxfile host="app.diagrams.net" modified="',
        datetime.now().isoformat(),
        '" agent="Azure Landing Zone Agent with LangGraph" etag="',
        uuid.uuid4().hex[:12],
        '" version="24.2.5" type="device">',
        '  <diagram id="azure-hub-spoke-langgraph" name="Azure Hub-Spoke (LangGraph)">',
        '    <mxGraphModel dx="1422" dy="754" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="1920" pageHeight="1200" math="0" shadow="0">',
        '      <root>',
        '        <mxCell id="0" />',
        '        <mxCell id="1" parent="0" />',
        ''
    ]
    
    # Hub infrastructure section
    hub_x, hub_y = 200, 100
    xml_parts.append(f'''
        <!-- Hub Infrastructure (LangGraph Orchestrated) -->
        <mxCell id="hub-container" value="Hub Infrastructure - Shared Services (LangGraph)" style="swimlane;whiteSpace=wrap;html=1;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;startSize=30;" vertex="1" parent="1">
          <mxGeometry x="{hub_x}" y="{hub_y}" width="300" height="200" as="geometry" />
        </mxCell>''')
    
    # Hub VNet
    xml_parts.append(f'''
        <mxCell id="hub-vnet" value="Hub VNet\\nShared Services" style="shape=mxgraph.azure.virtual_network;fillColor=#0078d4;strokeColor=#005a9e;fontColor=#ffffff;" vertex="1" parent="hub-container">
          <mxGeometry x="20" y="40" width="120" height="80" as="geometry" />
        </mxCell>''')
    
    # Hub services based on orchestration
    hub_services = hub_data.get("services", [])
    service_y = 40
    for i, service in enumerate(hub_services[:3]):  # Limit to 3 for layout
        service_x = 160 + (i * 40)
        service_name = service.replace('_', ' ').title()
        xml_parts.append(f'''
        <mxCell id="hub-{service}" value="{service_name}" style="shape=mxgraph.azure.azure_service;fillColor=#40e0d0;strokeColor=#008b8b;" vertex="1" parent="hub-container">
          <mxGeometry x="{service_x}" y="{service_y}" width="30" height="30" as="geometry" />
        </mxCell>''')
        service_y += 40
    
    # Production Spoke
    prod_x, prod_y = 600, 50
    xml_parts.append(f'''
        <!-- Production Spoke (LangGraph Orchestrated) -->
        <mxCell id="prod-spoke-container" value="Production Spoke - Workloads" style="swimlane;whiteSpace=wrap;html=1;fillColor=#00bcf2;strokeColor=#0078d4;fontColor=#ffffff;startSize=30;" vertex="1" parent="1">
          <mxGeometry x="{prod_x}" y="{prod_y}" width="300" height="200" as="geometry" />
        </mxCell>''')
    
    xml_parts.append(f'''
        <mxCell id="prod-vnet" value="Production VNet" style="shape=mxgraph.azure.virtual_network;fillColor=#00bcf2;strokeColor=#0078d4;fontColor=#ffffff;" vertex="1" parent="prod-spoke-container">
          <mxGeometry x="20" y="40" width="120" height="80" as="geometry" />
        </mxCell>''')
    
    # Production workloads based on orchestration
    spoke_services = spokes_data.get("services", [])
    workload_y = 40
    for i, service in enumerate(spoke_services[:3]):  # Limit for layout
        workload_x = 160 + (i * 40)
        service_name = service.replace('_', ' ').title()
        xml_parts.append(f'''
        <mxCell id="prod-{service}" value="{service_name}" style="shape=mxgraph.azure.azure_service;fillColor=#98fb98;strokeColor=#006400;" vertex="1" parent="prod-spoke-container">
          <mxGeometry x="{workload_x}" y="{workload_y}" width="30" height="30" as="geometry" />
        </mxCell>''')
        workload_y += 40
    
    # Development Spoke
    dev_x, dev_y = 600, 300
    xml_parts.append(f'''
        <!-- Development Spoke -->
        <mxCell id="dev-spoke-container" value="Development Spoke - Dev/Test" style="swimlane;whiteSpace=wrap;html=1;fillColor=#00bcf2;strokeColor=#0078d4;fontColor=#ffffff;startSize=30;" vertex="1" parent="1">
          <mxGeometry x="{dev_x}" y="{dev_y}" width="300" height="150" as="geometry" />
        </mxCell>''')
    
    xml_parts.append(f'''
        <mxCell id="dev-vnet" value="Development VNet" style="shape=mxgraph.azure.virtual_network;fillColor=#00bcf2;strokeColor=#0078d4;fontColor=#ffffff;" vertex="1" parent="dev-spoke-container">
          <mxGeometry x="20" y="40" width="120" height="80" as="geometry" />
        </mxCell>''')
    
    # Connections
    xml_parts.extend([
        '''
        <!-- Hub-Spoke Connections (LangGraph Orchestrated) -->
        <mxCell id="hub-to-prod" edge="1" source="hub-vnet" target="prod-vnet" parent="1">
          <mxGeometry relative="1" as="geometry">
            <mxPoint x="400" y="180" as="sourcePoint" />
            <mxPoint x="600" y="150" as="targetPoint" />
          </mxGeometry>
        </mxCell>
        <mxCell id="hub-to-dev" edge="1" source="hub-vnet" target="dev-vnet" parent="1">
          <mxGeometry relative="1" as="geometry">
            <mxPoint x="400" y="200" as="sourcePoint" />
            <mxPoint x="600" y="375" as="targetPoint" />
          </mxGeometry>
        </mxCell>''',
        '''
        <!-- Metadata -->
        <mxCell id="metadata" value="Generated with LangGraph Multi-Agent Orchestration\\nHub-Spoke Architecture Pattern" style="text;html=1;strokeColor=none;fillColor=none;align=left;verticalAlign=middle;whiteSpace=wrap;rounded=0;fontSize=10;" vertex="1" parent="1">
          <mxGeometry x="50" y="550" width="300" height="40" as="geometry" />
        </mxCell>
        
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>'''
    ])
    
    return "".join(xml_parts)

@app.post("/generate-azure-diagram")
def generate_azure_diagram_endpoint(inputs: CustomerInputs):
    """Generate Azure architecture diagram using Python Diagrams library with proper Azure icons"""
    try:
        # Generate Azure architecture diagram with proper icons
        diagram_path = generate_azure_architecture_diagram(inputs)
        
        # Read the generated PNG file
        with open(diagram_path, "rb") as f:
            diagram_data = f.read()
        
        # Generate professional documentation
        docs = generate_professional_documentation(inputs)
        
        # Encode the diagram as base64 for JSON response
        import base64
        diagram_base64 = base64.b64encode(diagram_data).decode('utf-8')
        
        # Get enterprise resources information for user feedback
        enterprise_resources = _get_enterprise_resources()
        was_enterprise_included = _should_include_enterprise_resources(inputs, enterprise_resources)
        
        response = {
            "success": True,
            "diagram_path": diagram_path,
            "diagram_base64": diagram_base64,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Python Diagrams",
                "diagram_format": "PNG with Azure official icons"
            }
        }
        
        # Add enterprise resources information and prompt
        if was_enterprise_included:
            response["enterprise_resources"] = {
                "auto_included": True,
                "resources": enterprise_resources,
                "prompt": _get_user_prompt_for_enterprise_resources(),
                "current_mode": inputs.enterprise_resources_mode or "auto_when_missing",
                "show_connections": inputs.show_enterprise_connections if inputs.show_enterprise_connections is not None else True
            }
        
        return response
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating Azure diagram: {str(e)}")

@app.get("/generate-azure-diagram/download/{filename}")
def download_azure_diagram(filename: str):
    """Download generated Azure architecture diagram PNG file"""
    try:
        file_path = f"/tmp/{filename}"
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="Diagram file not found")
        
        with open(file_path, "rb") as f:
            diagram_data = f.read()
        
        return Response(
            content=diagram_data,
            media_type="image/png", 
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error downloading diagram: {str(e)}")

@app.post("/generate-hub-spoke-vm-firewall")
def generate_hub_spoke_vm_firewall_diagram(inputs: CustomerInputs):
    """Generate enhanced hub-spoke diagram specifically showing VM in spoke connected to Firewall in hub with dashed lines"""
    try:
        logger.info("Starting enhanced hub-spoke VM-Firewall diagram generation")
        
        # Force the configuration to show VM in spoke and Firewall in hub for demonstration
        enhanced_inputs = CustomerInputs(
            business_objective=inputs.business_objective or "Hub-Spoke Architecture with VM and Firewall",
            compute_services=["virtual_machines"],  # VM goes to spoke
            network_services=["firewall"],  # Firewall goes to hub
            scalability=inputs.scalability or "medium",
            security_posture=inputs.security_posture or "standard",
            show_enterprise_connections=True,  # Enhanced visualization
            enterprise_resources_mode="enabled"  # Show all enterprise features
        )
        
        diagram_path = generate_azure_architecture_diagram(enhanced_inputs)
        
        if not os.path.exists(diagram_path):
            raise HTTPException(status_code=500, detail="Failed to generate hub-spoke diagram file")
        
        # Read the generated PNG file
        with open(diagram_path, "rb") as f:
            diagram_data = f.read()
        
        # Encode the diagram as base64 for JSON response
        import base64
        diagram_base64 = base64.b64encode(diagram_data).decode('utf-8')
        
        filename = os.path.basename(diagram_path)
        
        return {
            "success": True,
            "message": "Enhanced hub-spoke diagram with VM and Firewall generated successfully",
            "description": "VM placed in spoke VNet with dashed box, Firewall in hub, connected with dashed lines",
            "diagram_path": diagram_path,
            "diagram_base64": diagram_base64,
            "filename": filename,
            "download_url": f"/generate-azure-diagram/download/{filename}",
            "file_size": os.path.getsize(diagram_path),
            "features": [
                "VM placed in Production Spoke with dashed border",
                "Azure Firewall placed in Hub with solid border", 
                "Dashed line connections showing hub-spoke topology",
                "Enhanced visual distinction between hub and spoke resources",
                "Proper Azure stencils using mingrammer python library"
            ],
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Enhanced Hub-Spoke",
                "diagram_type": "Hub-Spoke VM-Firewall with Dashed Lines"
            }
        }
    except Exception as e:
        logger.error(f"Failed to generate hub-spoke VM-Firewall diagram: {str(e)}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to generate hub-spoke diagram: {str(e)}")

@app.post("/generate-drawio", response_class=Response)
def generate_drawio_endpoint(inputs: CustomerInputs):
    """Generate Draw.io XML file for download"""
    try:
        xml = generate_enhanced_drawio_xml(inputs)
        return Response(
            content=xml, 
            media_type="application/xml",
            headers={"Content-Disposition": "attachment; filename=azure-landing-zone.drawio"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating Draw.io XML: {str(e)}")


@app.post("/generate-comprehensive-azure-architecture")
def generate_comprehensive_azure_architecture(inputs: CustomerInputs):
    """Generate comprehensive Azure architecture with both Draw.io XML and PNG diagram"""
    logger.info("Starting comprehensive Azure architecture generation")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate Draw.io XML with comprehensive Azure stencils
        logger.info("Generating Draw.io XML...")
        drawio_xml = generate_enhanced_drawio_xml(inputs)
        logger.info(f"Draw.io XML generated successfully (size: {len(drawio_xml)} characters)")
        
        # Generate Azure PNG diagram with proper Azure icons
        logger.info("Generating Azure PNG diagram...")
        diagram_path = generate_azure_architecture_diagram(inputs)
        logger.info(f"Azure PNG diagram generated successfully: {diagram_path}")
        
        # Read the PNG file
        try:
            with open(diagram_path, "rb") as f:
                diagram_data = f.read()
            diagram_base64 = base64.b64encode(diagram_data).decode('utf-8')
            logger.info(f"PNG file read and encoded successfully (size: {len(diagram_data)} bytes)")
        except Exception as e:
            logger.error(f"Failed to read PNG file {diagram_path}: {e}")
            raise Exception(f"Failed to read generated PNG file: {str(e)}")
        
        # Generate professional documentation
        logger.info("Generating professional documentation...")
        docs = generate_professional_documentation(inputs)
        logger.info("Professional documentation generated successfully")
        
        # Count Azure stencils used
        import re
        shapes = re.findall(r'shape=mxgraph\.azure\.[^;\"\s]*', drawio_xml)
        
        result = {
            "success": True,
            "drawio_xml": drawio_xml,
            "png_diagram_path": diagram_path,
            "png_diagram_base64": diagram_base64,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "azure_stencils": {
                "total_used": len(shapes),
                "unique_used": len(set(shapes)),
                "stencils_list": sorted(list(set(shapes)))
            },
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Comprehensive Generator",
                "drawio_size": len(drawio_xml),
                "png_size": len(diagram_data)
            }
        }
        
        logger.info("Comprehensive Azure architecture generated successfully")
        return result
    
    except ValueError as ve:
        # Input validation errors
        error_msg = f"Invalid input: {str(ve)}"
        logger.error(error_msg)
        raise HTTPException(status_code=400, detail=error_msg)
    
    except Exception as e:
        # Log the full error for debugging
        error_msg = f"Error generating comprehensive architecture: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        
        # Return a user-friendly error
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to generate architecture. Error: {str(e)}"
        )

@app.post("/generate-interactive-azure-architecture")
def generate_interactive_azure_architecture(inputs: CustomerInputs):
    """Generate interactive Azure architecture with SVG diagram for web display"""
    logger.info("Starting interactive Azure architecture generation")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate Mermaid diagram
        logger.info("Generating Mermaid diagram...")
        mermaid_diagram = generate_professional_mermaid(inputs)
        logger.info("Mermaid diagram generated successfully")
        
        # Generate Azure SVG diagram with proper Azure icons
        logger.info("Generating Azure SVG diagram...")
        svg_content = ""
        svg_diagram_path = ""
        
        try:
            svg_diagram_path = generate_azure_architecture_diagram(inputs, format="svg")
            # Read the SVG file
            with open(svg_diagram_path, "r", encoding="utf-8") as f:
                svg_content = f.read()
            logger.info(f"Azure SVG diagram generated successfully: {svg_diagram_path}")
        except Exception as svg_error:
            logger.warning(f"SVG generation failed, using fallback: {str(svg_error)}")
            # Fallback: Create a simple SVG representation of the architecture
            svg_content = generate_simple_svg_diagram(inputs)
            logger.info("Using simple SVG fallback diagram")
        
        if svg_content:
            logger.info(f"SVG content ready (size: {len(svg_content)} characters)")
        else:
            logger.warning("No SVG content available, falling back to Mermaid only")
        
        # Generate Draw.io XML for compatibility
        logger.info("Generating Draw.io XML...")
        drawio_xml = generate_enhanced_drawio_xml(inputs)
        logger.info(f"Draw.io XML generated successfully (size: {len(drawio_xml)} characters)")
        
        # Generate professional documentation
        logger.info("Generating professional documentation...")
        try:
            # Use a timeout for the documentation generation
            import signal
            
            def timeout_handler(signum, frame):
                raise Exception("Documentation generation timed out")
            
            signal.signal(signal.SIGALRM, timeout_handler)
            signal.alarm(10)  # 10 second timeout
            
            docs = generate_professional_documentation(inputs)
            
            signal.alarm(0)  # Clear the alarm
            logger.info("Professional documentation generated successfully")
        except Exception as e:
            signal.alarm(0)  # Clear the alarm
            logger.warning(f"Documentation generation failed, using fallback: {str(e)}")
            docs = {
                "tsd": f"# Technical Specification Document\n\n## Azure Landing Zone Architecture\n\n**Organization:** {inputs.org_structure or 'Enterprise'}\n**Business Objective:** {inputs.business_objective or 'Not specified'}\n\n### Selected Services\n- Compute: {', '.join(inputs.compute_services or [])}\n- Network: {', '.join(inputs.network_services or [])}\n- Security: {', '.join(inputs.security_services or [])}\n\n*Full documentation requires AI service availability.*",
                "hld": f"# High Level Design\n\n## Azure Architecture Overview\n\nThis document outlines the high-level design for an Azure Landing Zone.\n\n### Key Components\n- Management Groups\n- Subscriptions\n- Resource Groups\n- Network Architecture\n\n*Detailed design requires AI service availability.*",
                "lld": f"# Low Level Design\n\n## Implementation Details\n\nThis document provides implementation guidance for the Azure Landing Zone.\n\n### Implementation Steps\n1. Set up Management Groups\n2. Configure Subscriptions\n3. Deploy Network Infrastructure\n4. Implement Security Controls\n\n*Detailed implementation guide requires AI service availability.*"
            }
        
        # Count Azure stencils used in Draw.io XML
        import re
        shapes = re.findall(r'shape=mxgraph\.azure\.[^;\"\s]*', drawio_xml)
        
        result = {
            "success": True,
            "mermaid": mermaid_diagram,
            "svg_diagram": svg_content,
            "svg_diagram_path": svg_diagram_path,
            "drawio_xml": drawio_xml,
            "tsd": docs["tsd"],
            "hld": docs["hld"],
            "lld": docs["lld"],
            "architecture_template": generate_architecture_template(inputs),
            "azure_stencils": {
                "total_used": len(shapes),
                "unique_used": len(set(shapes)),
                "stencils_list": sorted(list(set(shapes)))
            },
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Interactive Generator",
                "diagram_format": "SVG with Azure official icons",
                "svg_size": len(svg_content),
                "drawio_size": len(drawio_xml)
            }
        }
        
        logger.info("Interactive Azure architecture generated successfully")
        return result
    
    except ValueError as ve:
        # Input validation errors
        error_msg = f"Invalid input: {str(ve)}"
        logger.error(error_msg)
        raise HTTPException(status_code=400, detail=error_msg)
    
    except Exception as e:
        # Log the full error for debugging
        error_msg = f"Error generating interactive architecture: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        
        # Return a user-friendly error
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to generate interactive architecture. Error: {str(e)}"
        )

@app.post("/generate-png-diagram")
def generate_png_diagram(inputs: CustomerInputs):
    """Generate PNG diagram for download"""
    logger.info("Starting PNG diagram generation for download")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate PNG diagram
        logger.info("Generating PNG diagram...")
        png_path = generate_azure_architecture_diagram(inputs, format="png")
        logger.info(f"PNG diagram generated successfully: {png_path}")
        
        # Read and encode the PNG file
        with open(png_path, "rb") as f:
            png_content = f.read()
        
        png_base64 = base64.b64encode(png_content).decode("utf-8")
        
        return {
            "success": True,
            "png_base64": png_base64,
            "png_path": png_path,
            "file_size": len(png_content),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "format": "PNG"
            }
        }
        
    except Exception as e:
        logger.error(f"Error generating PNG diagram: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate PNG diagram: {str(e)}")

@app.post("/generate-svg-diagram")
def generate_svg_diagram(inputs: CustomerInputs):
    """Generate SVG diagram for download"""
    logger.info("Starting SVG diagram generation for download")
    
    try:
        # Validate inputs early
        validate_customer_inputs(inputs)
        logger.info("Input validation completed successfully")
        
        # Generate SVG diagram
        logger.info("Generating SVG diagram...")
        svg_path = generate_azure_architecture_diagram(inputs, format="svg")
        logger.info(f"SVG diagram generated successfully: {svg_path}")
        
        # Read the SVG file
        with open(svg_path, "r", encoding="utf-8") as f:
            svg_content = f.read()
        
        return {
            "success": True,
            "svg_content": svg_content,
            "svg_path": svg_path,
            "file_size": len(svg_content),
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "format": "SVG"
            }
        }
        
    except Exception as e:
        logger.error(f"Error generating SVG diagram: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate SVG diagram: {str(e)}")

@app.post("/upload-file")
async def upload_file(file: UploadFile = File(...)):
    """Upload and process files (PDF, Excel, PowerPoint) for AI analysis"""
    try:
        # Validate file type
        allowed_extensions = ['.pdf', '.xlsx', '.xls', '.pptx', '.ppt']
        file_extension = os.path.splitext(file.filename.lower())[1]
        
        if file_extension not in allowed_extensions:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type. Allowed types: {', '.join(allowed_extensions)}"
            )
        
        # Validate file size (max 10MB)
        max_file_size = 10 * 1024 * 1024  # 10MB
        file_content = await file.read()
        
        if len(file_content) > max_file_size:
            raise HTTPException(
                status_code=400,
                detail="File too large. Maximum size is 10MB."
            )
        
        # Process the file with AI
        file_type = file_extension[1:]  # Remove the dot
        analysis_result = process_uploaded_document(file_content, file.filename, file_type)
        
        # Return file info and analysis
        return {
            "success": True,
            "filename": file.filename,
            "file_type": file_type,
            "file_size": len(file_content),
            "analysis": analysis_result,
            "upload_timestamp": datetime.now().isoformat()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading file {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.post("/analyze-url")
def analyze_url(request: Dict[str, str]):
    """Analyze URL content for Azure architecture insights"""
    try:
        url = request.get("url")
        if not url:
            raise HTTPException(status_code=400, detail="URL is required")
        
        # Validate URL format
        if not url.startswith(('http://', 'https://')):
            raise HTTPException(status_code=400, detail="URL must start with http:// or https://")
        
        # Analyze the URL with AI
        analysis_result = analyze_url_content(url)
        
        return {
            "success": True,
            "url": url,
            "analysis": analysis_result,
            "analysis_timestamp": datetime.now().isoformat()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing URL {url}: {e}")
        raise HTTPException(status_code=500, detail=f"Error analyzing URL: {str(e)}")

@app.get("/templates")
def get_templates():
    """Get available Azure Landing Zone templates"""
    return {
        "templates": AZURE_TEMPLATES,
        "azure_services": AZURE_SERVICES_MAPPING
    }

@app.post("/generate-intelligent-diagram")
def generate_intelligent_diagram(request: Dict[str, str]):
    """Generate intelligent Azure architecture diagram from natural language requirements"""
    logger.info("Starting intelligent diagram generation")
    
    try:
        # Get requirements text from request
        requirements_text = request.get("requirements")
        if not requirements_text:
            raise HTTPException(status_code=400, detail="Requirements text is required")
        
        # Validate requirements text length
        if len(requirements_text) < 10:
            raise HTTPException(status_code=400, detail="Requirements text too short (minimum 10 characters)")
        
        if len(requirements_text) > 5000:
            raise HTTPException(status_code=400, detail="Requirements text too long (maximum 5000 characters)")
        
        logger.info(f"Processing requirements: {requirements_text[:100]}...")
        
        # Check if intelligent generator is available
        if not intelligent_generator:
            # Check if the issue is missing API key
            openai_api_key = os.getenv('OPENAI_API_KEY')
            if not openai_api_key:
                error_detail = "OpenAI API key is required for intelligent diagram generation. Please set the OPENAI_API_KEY environment variable with a valid OpenAI API key."
            else:
                error_detail = "Intelligent diagram generator failed to initialize. Please check your OpenAI API key and try again."
            raise HTTPException(status_code=503, detail=error_detail)
        
        # Generate diagram from natural language
        result = intelligent_generator.generate_from_natural_language(requirements_text)
        
        # Execute the generated Python code to create the actual diagram
        diagram_path = None
        diagram_base64 = None
        execution_error = None
        
        try:
            # Execute the Python code in a safe environment
            logger.info("Executing generated diagram code...")
            
            # Create a safe execution environment with necessary built-ins
            # Include essential built-ins that are needed for diagram generation
            import builtins
            safe_builtins = {
                '__import__': builtins.__import__,
                '__build_class__': builtins.__build_class__,
                'len': len,
                'str': str,
                'dict': dict,
                'list': list,
                'tuple': tuple,
                'set': set,
                'bool': bool,
                'int': int,
                'float': float,
                'range': range,
                'enumerate': enumerate,
                'zip': zip,
                'abs': abs,
                'min': min,
                'max': max,
                'sum': sum,
                'sorted': sorted,
                'reversed': reversed,
                'getattr': getattr,
                'setattr': setattr,
                'hasattr': hasattr,
                'isinstance': isinstance,
                'issubclass': issubclass,
                'type': type,
                'repr': repr,
                'format': format,
                # Note: Deliberately excluding potentially dangerous functions like:
                # - open, exec, eval, compile, __import__ with custom hooks
                # - file system operations, network operations, etc.
            }
            
            exec_globals = {
                '__builtins__': safe_builtins,
                'Diagram': Diagram,
                'Cluster': Cluster,
                'Edge': Edge,
                'VM': VM,
                'AKS': AKS,
                'AppServices': AppServices,
                'VirtualNetworks': VirtualNetworks,
                'ApplicationGateway': ApplicationGateway,
                'LoadBalancers': LoadBalancers,
                'Firewall': Firewall,
                'StorageAccounts': StorageAccounts,
                'SQLDatabases': SQLDatabases,
                'CosmosDb': CosmosDb,
                'KeyVaults': KeyVaults,
                'ActiveDirectory': ActiveDirectory,
                'SynapseAnalytics': SynapseAnalytics,
                'DataFactories': DataFactories,
                'Databricks': Databricks,
                'LogicApps': LogicApps,
                'ServiceBus': ServiceBus,
                'APIManagement': APIManagement,
                'Devops': Devops,
                'SecurityCenter': SecurityCenter,
                'Sentinel': Sentinel
            }
            
            # Set output directory for diagram
            output_dir = get_safe_output_directory()
            exec_globals['output_dir'] = output_dir
            
            # Modify the code to save to our output directory
            modified_code = result.python_code.replace(
                'filename="azure_architecture_intelligent"',
                f'filename="{output_dir}/azure_architecture_intelligent"'
            )
            
            # Execute the code
            exec(modified_code, exec_globals)
            
            # Find the generated diagram file
            possible_paths = [
                f"{output_dir}/azure_architecture_intelligent.png",
                f"{output_dir}/azure_architecture_intelligent.svg"
            ]
            
            for path in possible_paths:
                if os.path.exists(path):
                    diagram_path = path
                    break
            
            if diagram_path:
                # Read and encode the diagram
                with open(diagram_path, "rb") as f:
                    diagram_data = f.read()
                diagram_base64 = base64.b64encode(diagram_data).decode('utf-8')
                logger.info(f"Diagram generated successfully: {diagram_path}")
            else:
                execution_error = "Generated diagram file not found"
                logger.warning("Generated diagram file not found")
                
        except Exception as e:
            execution_error = f"Error executing diagram code: {str(e)}"
            logger.error(f"Error executing diagram code: {e}")
        
        # Return comprehensive result
        response = {
            "success": True,
            "requirements_processed": requirements_text,
            "generated_code": result.python_code,
            "description": result.description,
            "review_comments": result.review_comments,
            "enterprise_compliance_score": result.enterprise_compliance_score,
            "diagram_path": diagram_path,
            "diagram_base64": diagram_base64,
            "execution_error": execution_error,
            "intelligent_features": {
                "natural_language_parsing": True,
                "enterprise_review": True,
                "code_generation": True,
                "compliance_scoring": True
            },
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Intelligent Azure Architecture Generator",
                "llm_mode": "openai"
            }
        }
        
        logger.info("Intelligent diagram generation completed successfully")
        return response
        
    except HTTPException:
        raise
    except Exception as e:
        error_msg = f"Error generating intelligent diagram: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        
        # Create detailed error response for different types of errors
        error_detail = {
            "error": error_msg,
            "error_type": type(e).__name__
        }
        
        # Add specific handling for common errors
        if "OpenAI API" in str(e) or "API call failed" in str(e):
            error_detail["error_category"] = "api_error"
            error_detail["suggestion"] = "Please check your OpenAI API key and internet connection."
        elif "execution_error" in str(e):
            error_detail["error_category"] = "execution_error"
            error_detail["suggestion"] = "There was an error executing the generated diagram code."
        else:
            error_detail["error_category"] = "general_error"
            error_detail["suggestion"] = "Please try again or contact support if the issue persists."
        
        raise HTTPException(status_code=500, detail=error_detail)

@app.post("/enhance-diagram")
def enhance_diagram(request: Dict[str, str]):
    """Enhance existing diagram with new requirements"""
    logger.info("Starting diagram enhancement")
    
    try:
        # Get parameters from request
        existing_code = request.get("existing_code")
        enhancement_requirements = request.get("enhancement_requirements")
        
        if not existing_code:
            raise HTTPException(status_code=400, detail="Existing code is required")
        if not enhancement_requirements:
            raise HTTPException(status_code=400, detail="Enhancement requirements are required")
        
        # Check if intelligent generator is available
        if not intelligent_generator:
            raise HTTPException(status_code=503, detail="Intelligent diagram generator not available")
        
        # Enhance the diagram
        result = intelligent_generator.enhance_existing_diagram(existing_code, enhancement_requirements)
        
        return {
            "success": True,
            "original_code": existing_code,
            "enhancement_requirements": enhancement_requirements,
            "enhanced_code": result.python_code,
            "description": result.description,
            "review_comments": result.review_comments,
            "enterprise_compliance_score": result.enterprise_compliance_score,
            "metadata": {
                "enhanced_at": datetime.now().isoformat(),
                "version": "1.0.0"
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        error_msg = f"Error enhancing diagram: {str(e)}"
        logger.error(error_msg)
        raise HTTPException(status_code=500, detail=error_msg)

@app.post("/validate-architecture")
def validate_architecture_endpoint(inputs: CustomerInputs):
    """
    Validate Azure Landing Zone architecture based on customer inputs.
    
    This endpoint converts customer inputs to architecture format and validates
    all resources against Azure Landing Zone best practices and compliance rules.
    """
    try:
        logger.info("Starting architecture validation from customer inputs")
        
        # Convert customer inputs to architecture format
        architecture = convert_customer_inputs_to_architecture(inputs)
        
        # Validate the architecture
        validation_result = validate_architecture(architecture)
        
        # Generate diagram structure with validation results
        diagram_structure = generate_diagram_structure(architecture, validation_result)
        
        # Prepare response
        response = {
            "success": True,
            "validation": {
                "passed": validation_result.passed,
                "compliance_score": validation_result.compliance_score,
                "total_resources": validation_result.total_resources,
                "total_issues": validation_result.issues_count,
                "summary": validation_result.summary,
                "issues": [
                    {
                        "resource_name": issue.resource_name,
                        "resource_type": issue.resource_type,
                        "issue_type": issue.issue_type,
                        "severity": issue.severity.value,
                        "message": issue.message,
                        "recommendation": issue.recommendation,
                        "rule_id": issue.rule_id,
                        "compliance_impact": issue.compliance_impact
                    }
                    for issue in validation_result.issues
                ]
            },
            "diagram_structure": {
                "nodes": len(diagram_structure.nodes),
                "connections": len(diagram_structure.connections), 
                "layout_type": diagram_structure.layout.get("type", "hub-spoke"),
                "metadata": diagram_structure.metadata
            },
            "architecture": architecture,
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Architecture Validator"
            }
        }
        
        logger.info(f"Architecture validation completed. Score: {validation_result.compliance_score}%, Issues: {validation_result.issues_count}")
        return response
        
    except Exception as e:
        error_msg = f"Error validating architecture: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=error_msg)

@app.post("/validate-and-generate-diagram")
def validate_and_generate_diagram(inputs: CustomerInputs):
    """
    Validate architecture and generate comprehensive diagram with validation feedback.
    
    This endpoint combines validation with diagram generation, providing both
    architectural validation results and visual diagram output.
    """
    try:
        logger.info("Starting validation and diagram generation")
        
        # Convert and validate architecture
        architecture = convert_customer_inputs_to_architecture(inputs)
        validation_result = validate_architecture(architecture)
        diagram_structure = generate_diagram_structure(architecture, validation_result)
        
        # Generate traditional outputs (mermaid, drawio)
        mermaid_diagram = generate_professional_mermaid(inputs)
        drawio_xml = generate_enhanced_drawio_xml(inputs)
        
        # Generate documentation
        docs = generate_professional_documentation(inputs)
        
        # Create comprehensive response
        response = {
            "success": True,
            "validation": {
                "passed": validation_result.passed,
                "compliance_score": validation_result.compliance_score,
                "total_issues": validation_result.issues_count,
                "critical_issues": validation_result.summary.get("critical_issues", 0),
                "high_issues": validation_result.summary.get("high_issues", 0),
                "top_recommendations": validation_result.summary.get("top_recommendations", [])[:5]
            },
            "diagrams": {
                "mermaid": mermaid_diagram,
                "drawio": drawio_xml
            },
            "documentation": {
                "tsd": docs["tsd"],
                "hld": docs["hld"], 
                "lld": docs["lld"]
            },
            "architecture_template": generate_architecture_template(inputs),
            "validation_details": [
                {
                    "resource": issue.resource_name,
                    "type": issue.resource_type,
                    "severity": issue.severity.value,
                    "message": issue.message,
                    "recommendation": issue.recommendation
                }
                for issue in validation_result.issues
                if issue.severity.value in ["critical", "high"]
            ][:10],  # Top 10 most important issues
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0.0",
                "agent": "Azure Landing Zone Agent - Validation + Generation",
                "validation_enabled": True
            }
        }
        
        logger.info("Validation and diagram generation completed successfully")
        return response
        
    except Exception as e:
        error_msg = f"Error in validation and diagram generation: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=error_msg)

@app.get("/enterprise-resources-prompt")
def get_enterprise_resources_prompt():
    """Get information about enterprise resource auto-inclusion and user prompt"""
    enterprise_resources = _get_enterprise_resources()
    
    return {
        "prompt": _get_user_prompt_for_enterprise_resources(),
        "enterprise_resources": enterprise_resources,
        "resource_descriptions": {
            "key_vault": "Azure Key Vault for secrets, keys, and certificate management",
            "active_directory": "Azure Active Directory for identity and access management", 
            "firewall": "Azure Firewall for network security and traffic filtering",
            "monitor": "Azure Monitor for observability, logging, and alerting"
        },
        "modes": {
            "always_include": "Always include enterprise resources in every diagram",
            "auto_when_missing": "Include enterprise resources only when not explicitly specified (default)",
            "never_auto_include": "Never automatically include enterprise resources"
        },
        "default_settings": {
            "enterprise_resources_mode": "auto_when_missing",
            "show_enterprise_connections": True
        }
    }

@app.get("/services")
def get_services():
    """Get available Azure services categorized for form selection"""
    services_by_category = {}
    
    for service_key, service_info in AZURE_SERVICES_MAPPING.items():
        category = service_info["category"]
        if category not in services_by_category:
            services_by_category[category] = []
        
        services_by_category[category].append({
            "key": service_key,
            "name": service_info["name"],
            "icon": service_info["icon"],
            "azure_icon": service_info.get("azure_icon", ""),
        })
    
    return {
        "categories": services_by_category,
        "category_mapping": {
            "compute": "Compute Services",
            "network": "Networking Services", 
            "storage": "Storage Services",
            "database": "Database Services",
            "security": "Security Services",
            "monitoring": "Monitoring Services",
            "ai": "AI & Machine Learning",
            "analytics": "Data & Analytics",
            "integration": "Integration Services",
            "devops": "DevOps & Governance",
            "backup": "Backup & Recovery"
        }
    }
